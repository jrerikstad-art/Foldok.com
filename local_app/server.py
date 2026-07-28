#!/usr/bin/env python3
"""Foldok local workbench — the REAL engine behind a browser UI.

The missing "start a project" experience: create a project from a folder,
index its files (with cost estimate), confirm the artifact model
(checkpoint A), pick a template, generate (checkpoints B+C), read the
draft. Every call goes through foldok_compile.py — no mock data.

Start:  python local_app/server.py        (or scripts/workbench.ps1)
Open:   http://127.0.0.1:8766

State:
  local_app/projects.json          — registered projects
  <folder>/.foldok_state.json     — artifact model, confirmation, gaps
  <folder>/.foldok_cache/         — per-file index cache (engine-owned)
  <folder>/draft.md                — generated draft (latest)
  <folder>/.foldok_drafts/        — one archived draft per template
  <folder>/Rapporter/              — named exports (Konstruksjonsrapport.md, …)
"""
import json, os, re, sys, threading, time, traceback, uuid
from datetime import datetime, timezone
from concurrent.futures import ThreadPoolExecutor, as_completed
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
from urllib.parse import unquote

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

# Background jobs (WORKORDER 0.55 — cancel / heartbeat / budget)
JOBS: dict = {}
LOCK = threading.Lock()
INDEX_CHUNK_SIZE = 100
HEARTBEAT_TIMEOUT_S = 60
DEFAULT_INDEX_BUDGET_EUR = 10.0


def load_env_file():
    """Load KEY=value lines from foldok-engine/.env (not committed)."""
    env_path = ROOT / ".env"
    if not env_path.exists():
        return
    for line in env_path.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, _, val = line.partition("=")
        key, val = key.strip(), val.strip().strip('"').strip("'")
        if key and key not in os.environ:
            os.environ[key] = val


load_env_file()
KEY_SET = len(os.environ.get("ANTHROPIC_API_KEY", "")) > 30
if not os.environ.get("ANTHROPIC_API_KEY"):
    # allow import + browsing cached data without a key; API calls will fail clearly
    os.environ["ANTHROPIC_API_KEY"] = "missing-key"


def pick_folder_dialog(title="Velg prosjektmappe"):
    """Native folder picker — prefer in-UI /api/browse; this is a fallback only."""
    if sys.platform == "win32":
        import subprocess, tempfile
        # Write result to a temp file so a hung UI can't corrupt stdout parsing.
        out = Path(tempfile.gettempdir()) / f"foldok-folder-{uuid.uuid4().hex[:8]}.txt"
        ps = (
            "Add-Type -AssemblyName System.Windows.Forms; "
            "[System.Windows.Forms.Application]::EnableVisualStyles(); "
            "$f = New-Object System.Windows.Forms.FolderBrowserDialog; "
            f"$f.Description = '{title.replace(chr(39), '')}'; "
            "$f.ShowNewFolderButton = $true; "
            "$f.RootFolder = [Environment+SpecialFolder]::MyComputer; "
            "if ($f.ShowDialog() -eq 'OK') { "
            f"  Set-Content -LiteralPath '{out.as_posix()}' -Value $f.SelectedPath -Encoding UTF8 "
            "}"
        )
        try:
            subprocess.run(
                ["powershell", "-NoProfile", "-STA", "-WindowStyle", "Normal", "-Command", ps],
                timeout=120, creationflags=getattr(subprocess, "CREATE_NEW_CONSOLE", 0),
            )
            if out.exists():
                chosen = out.read_text(encoding="utf-8").strip() or None
                out.unlink(missing_ok=True)
                return chosen
            return None
        except Exception:
            try: out.unlink(missing_ok=True)
            except Exception: pass
            return None
    try:
        import tkinter as tk
        from tkinter import filedialog
        root = tk.Tk()
        root.withdraw()
        root.attributes("-topmost", True)
        chosen = filedialog.askdirectory(title=title) or None
        root.destroy()
        return chosen
    except Exception:
        return None


def browse_dirs(path):
    """List subfolders for the in-UI picker (no modal, no server block)."""
    home = Path.home()
    docs = home / "Documents"
    onedrive = home / "OneDrive" / "Documents"
    if not path or not str(path).strip():
        roots = []
        for label, p in (("Documents", docs), ("Cloud Documents", onedrive), ("Home", home)):
            if p.is_dir():
                roots.append({"name": label, "path": str(p), "kind": "root"})
        # Drive letters on Windows
        if sys.platform == "win32":
            import string
            for letter in string.ascii_uppercase:
                drive = Path(f"{letter}:/")
                if drive.exists():
                    roots.append({"name": f"{letter}:", "path": str(drive), "kind": "drive"})
        return {"path": None, "parent": None, "entries": roots}

    root = Path(path)
    if not root.is_dir():
        return {"error": f"Mappen finnes ikke: {path}"}
    entries = []
    try:
        for child in sorted(root.iterdir(), key=lambda p: p.name.lower()):
            if not child.is_dir():
                continue
            if child.name.startswith(".") or child.name.lower() in {
                "capture", "foldok-engine", "node_modules", "__pycache__", "releases",
                ".git", ".venv", "venv", "dist", "build",
            }:
                continue
            entries.append({"name": child.name, "path": str(child), "kind": "dir"})
    except PermissionError:
        return {"error": f"Ingen tilgang: {path}"}
    parent = str(root.parent) if root.parent != root else None
    return {"path": str(root), "parent": parent, "name": root.name, "entries": entries[:200]}

import foldok_compile as fc
import foldok_paths as fpaths
import bom_engine as be

APP_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(APP_DIR))
import doc_state as ds
import hub_chat as hub
import editor_chat as edchat
import chat_attach as chattach
import local_learning as learning
import index_tools as idxtools
import document_type_registry as dtr
import index_prescan as prescan
import account_metering as acct
try:
    acct.install_compile_hooks()
except Exception:
    pass
import project_io as pio
import export_formats as expfmt
PROJECTS_FILE = APP_DIR / "projects.json"
COMPANY_TEMPLATES_DIR = APP_DIR / "company_templates"
ATTACH_STAGING = APP_DIR / ".attach_staging"
OUTPUT_NAMES = {"draft.md", "mapping.json", "compile_log.txt", "expected.json",
                "project_findings.xlsx"}
# Never treat these as field sources (product UI shots, engine tree, deps, …)
SKIP_DIR_NAMES = {
    *fpaths.SKIP_PRODUCT_DIRS,
    ".git", ".cursor", "agent-transcripts", "terminals", "assets",
    *fpaths.SKIP_CACHE_DIR_NAMES,
}
def open_knowledge_engine(primary_folder, enable_vectors=True):
    """Lazy import — pandas/openpyxl required; LanceDB optional."""
    from hybrid_knowledge_engine import HybridKnowledgeEngine
    return HybridKnowledgeEngine(primary_folder, enable_vectors=enable_vectors)


# ── persistence ──────────────────────────────────────────────────────
def load_projects():
    if not PROJECTS_FILE.exists():
        return []
    projects = json.loads(PROJECTS_FILE.read_text(encoding="utf-8"))
    for p in projects:  # migrate pre-0.13 single-folder records
        if "folders" not in p:
            p["folders"] = [p.pop("folder")]
    # Never surface the Foldok product repo as a customer project
    # (UI mock screenshots in capture/ become «Løfteverktøy for batteripakke»).
    cleaned, changed = [], False
    for p in projects:
        folders = p.get("folders") or []
        if folders and looks_like_product_repo(folders[0]):
            changed = True
            continue
        cleaned.append(p)
    if changed:
        save_projects(cleaned)
    return cleaned


def save_projects(projects):
    PROJECTS_FILE.write_text(json.dumps(projects, indent=2, ensure_ascii=False), encoding="utf-8")


def get_project(pid):
    for p in load_projects():
        if p["id"] == pid:
            return p
    return None


# ── ISOLATION RULE (BUGFIX 0.19) ─────────────────────────────────────
# No module-level project / index / state cache keyed by "current".
# Every request carries a project id. Folder, state path, and index are
# derived from that id inside the handler — same resolver the compile
# pipeline uses. Mismatch → loud error, never a silent wrong answer.


class IsolationError(Exception):
    """Raised when a resolved folder does not belong to the requested project."""


def resolve_project(pid):
    """Sole project resolver for handlers. Returns (project, folders, primary_folder)."""
    if not pid or not str(pid).strip():
        raise ValueError("missing project id")
    p = get_project(str(pid).strip())
    if not p:
        raise LookupError(f"unknown project: {pid}")
    folders = list(p.get("folders") or [])
    if not folders:
        raise ValueError(f"project {pid} has no folders")
    primary = folders[0]
    assert_folder_on_project(p, primary)
    return p, folders, primary


def assert_folder_on_project(p, folder):
    folders = p.get("folders") or []
    # Normalize for Windows path compare
    want = str(Path(folder))
    owned = {str(Path(f)) for f in folders}
    if want not in owned:
        raise IsolationError(
            f"ISOLATION: resolved folder {want!r} not on project {p['id']} "
            f"({p.get('name')}); owned={sorted(owned)}"
        )


def project_banner(p, folder):
    return f"PROSJEKT: {p.get('name') or p['id']} · MAPPE: {folder}"


def load_project_index(pid, lang="no", user_facts=None):
    """Index loader keyed only by project id — one resolver, one truth."""
    p, folders, primary = resolve_project(pid)
    index = load_index(folders, lang, user_facts, project_name=p.get("name"))
    return p, folders, primary, index


def log_chat_isolation(handler_name, pid, p, folders, primary, index):
    first = (index[0].get("file") if index else None) or "(empty)"
    print(
        f"[isolation:{handler_name}] id={pid!r} name={p.get('name')!r} "
        f"folder={primary!r} state={state_path(primary)} "
        f"index_n={len(index)} first_file={first!r}",
        flush=True,
    )


def fact_key_inventory(index, limit=40):
    """Top fact keys by count — keys only, never values (token-cheap)."""
    from collections import Counter
    c = Counter()
    for e in index or []:
        for f in e.get("facts") or []:
            k = f.get("key")
            if k:
                c[k] += 1
    return [{"key": k, "count": n} for k, n in c.most_common(limit)]


def build_project_chat_context(p, folders, primary, state, index=None, *,
                               lang="no", conversation_limit=80):
    """Non-negotiable context for EVERY in-project chat model call (WO 0.20 A1).

    Includes: project name + folder + file count + indexed count; full artifact
    (incl. confidence); document list + gap counts + active document; fact-key
    inventory (keys+counts only); open conversation history.
    """
    index = index if index is not None else load_index(
        folders, lang, (state or {}).get("user_facts"), project_name=p.get("name"))
    rows = file_rows(folders)
    file_count = sum(1 for r in rows if r.get("kind") != "skipped")
    indexed_count = sum(1 for e in (index or []) if e.get("kind") != "skipped")
    art = (state or {}).get("artifact")
    if not isinstance(art, dict):
        art = {}
    # Full artifact — never a prose summary
    art_json = json.dumps(art, ensure_ascii=False, indent=2)

    templates = templates_list()
    docs = list_documents(primary, state or {}, templates)
    active_tf = (state or {}).get("active_template") or (state or {}).get("template")
    doc_lines = []
    for d in docs:
        label = d.get("name_no") or d.get("name") or d.get("template") or "?"
        marker = " ← ACTIVE" if d.get("template") == active_tf else ""
        doc_lines.append(
            f"- {label} · template={d.get('template')} · "
            f"gaps={d.get('gaps', 0)} blocking={d.get('blocking', 0)}{marker}"
        )
    if not doc_lines:
        gs = ds.gaps_summary((state or {}).get("gaps") or [])
        if active_tf:
            doc_lines.append(
                f"- (active) {active_tf} · gaps={gs.get('total', 0)} "
                f"blocking={gs.get('blocking', 0)} ← ACTIVE"
            )
        else:
            doc_lines.append("(ingen dokumenter ennå)")

    inv = fact_key_inventory(index, limit=40)
    inv_txt = ", ".join(f"{x['key']}:{x['count']}" for x in inv) or "(ingen fakta i indeks)"
    brief = edchat.corpus_brief(index, file_count)

    # BUGFIX_0.19 §A extended — conversation scoped to this project_id only
    pid = p.get("id")
    conv = edchat.conversation_for_project(state or {}, pid)
    if conversation_limit and len(conv) > conversation_limit:
        conv = conv[-conversation_limit:]
    hist_lines = []
    for t in conv:
        role = (t.get("role") or "?").upper()
        text = (t.get("text") or "").strip()
        if not text:
            continue
        if len(text) > 800:
            text = text[:800] + "…"
        hist_lines.append(f"{role}: {text}")
    hist_txt = "\n".join(hist_lines) if hist_lines else "(tom samtale)"

    folder_list = " | ".join(str(f) for f in folders)
    text = (
        "=== PROJECT CHAT CONTEXT (engine-owned; always attached) ===\n"
        f"PROJECT: {p.get('name') or p.get('id')}\n"
        f"PROJECT_ID: {p.get('id')}\n"
        f"FOLDER: {primary}\n"
        f"FOLDERS: {folder_list}\n"
        f"FILE_COUNT: {file_count}\n"
        f"INDEXED_COUNT: {indexed_count}\n"
        f"ACTIVE_DOCUMENT: {active_tf or '(none)'}\n"
        f"CORPUS BRIEF: {brief}\n\n"
        f"ARTIFACT MODEL (full JSON, incl. confidence):\n{art_json}\n\n"
        f"DOCUMENTS (+ gap counts; ← ACTIVE marks current):\n" + "\n".join(doc_lines) + "\n\n"
        f"FACT KEY INVENTORY (keys only, top {min(40, len(inv) or 40)} by count — NO values):\n"
        f"{inv_txt}\n\n"
        f"CONVERSATION HISTORY:\n{hist_txt}\n"
        "=== END PROJECT CHAT CONTEXT ==="
    )
    return {
        "text": text,
        "file_count": file_count,
        "indexed_count": indexed_count,
        "corpus_brief": brief,
        "fact_keys": inv,
        "documents": docs,
        "artifact": art,
        "active_template": active_tf,
        "index": index,
    }


def chat_turn_extras(message, index, artifact, file_count=None):
    """§7 — zero-token grounding pack for this user message."""
    brief = edchat.corpus_brief(index, file_count if file_count is not None else 0)
    known = edchat.known_from_index(
        message, index, artifact, search_fn=fc.search_fact_candidates)
    open_ended = edchat.is_open_ended_create(message)
    # Rough generate estimate: ~€0.022–0.04 per section × typical 8–10
    estimate_eur = 0.22
    return {
        "corpus_brief": brief,
        "known_block": (
            f"ALREADY KNOWN FROM INDEX (zero-token search — state these, do not re-ask):\n{known}"
        ),
        "open_ended": open_ended,
        "estimate_eur": estimate_eur,
        "policy": edchat.CHAT_AGENT_POLICY,
    }


def build_artifact_assist_sources(pid, lang="no", user_facts=None):
    """Chat context for Checkpoint A — always bound to the requested project id."""
    p, folders, primary, index = load_project_index(pid, lang, user_facts)
    log_chat_isolation("artifact/assist", pid, p, folders, primary, index)
    banner = project_banner(p, primary)
    caps = "\n".join(f"[{e['file']}] {e.get('caption', '')}" for e in index)[:10000]
    state = load_state(primary)
    chat_ctx = build_project_chat_context(p, folders, primary, state, index, lang=lang)
    return {
        "project": p,
        "folders": folders,
        "primary": primary,
        "index": index,
        "banner": banner,
        "captions": caps,
        "state_path": str(state_path(primary)),
                "chat_context": chat_ctx["text"],
        "chat_context_meta": {
            k: chat_ctx[k] for k in ("file_count", "fact_keys", "corpus_brief") if k in chat_ctx
        },
    }


def state_path(folder):
    return fpaths.state_path(folder)


def drafted_templates_dir(folder):
    return fpaths.templates_dir(folder)


def load_template(template_file, folder=None):
    if not template_file:
        return None
    name = Path(template_file).name
    if folder:
        local = drafted_templates_dir(folder) / name
        if local.exists():
            return json.loads(local.read_text(encoding="utf-8"))
    owned = COMPANY_TEMPLATES_DIR / name
    if owned.exists():
        return json.loads(owned.read_text(encoding="utf-8"))
    tpath = ROOT / "templates" / name
    if tpath.exists():
        return json.loads(tpath.read_text(encoding="utf-8"))
    # WORKORDER 0.59/0.60 — never fail «Tomt dokument» if the json is missing from a stale tree
    if name == "sketch_document.json":
        return {
            "template_key": "sketch_document",
            "name": "Sketch Document",
            "name_no": "Tomt dokument (skisse)",
            "description": "Blank A4 canvas",
            "version": 1,
            "language_default": "no",
            "export_price_tier": "standard",
            "document_species": "sketch",
            "origin": "system",
            "sections": [{
                "section_key": "canvas", "title": "Canvas", "title_no": "Lerret",
                "position": 1, "required": False, "gap_severity": "info",
                "required_facts": [], "required_media": {}, "required_content": [],
            }],
        }
    return None


def _ensure_sketch_fields(state, template, prev_doc=None):
    """Keep sketch canvas metadata when (re)building a sketch document shell."""
    doc = state.get("doc") or {}
    if (template or {}).get("document_species") != "sketch" and doc.get("document_species") != "sketch":
        return
    prev = prev_doc or {}
    sketch = doc.get("sketch") or prev.get("sketch") or {"placeholders": [], "mode": True}
    sketch.setdefault("placeholders", [])
    sketch["mode"] = True
    doc["sketch"] = sketch
    doc["document_species"] = "sketch"
    state["doc"] = doc


def ensure_doc_for_template(folder, state, template_file):
    """Ensure state.doc holds sections for the requested template (hydrate from archive if needed).

    WORKORDER 0.58 §0 — never leave a foreign template's doc in place when switching.
    folder may be None for folder-less projects (WORKORDER 0.61).
    """
    if not template_file:
        return state
    doc = state.get("doc") or {}
    if doc.get("template_file") == template_file and doc.get("sections"):
        template = load_template(template_file, folder)
        _ensure_sketch_fields(state, template, doc)
        return state
    template = load_template(template_file, folder)
    if not template:
        return state
    import form_model as fm
    raw = None
    if folder:
        raw = read_draft_content(folder, template_file)
    if raw:
        state["doc"] = {
            "template_file": template_file,
            "sections": ds.split_draft_to_sections(raw, template),
            "generated_at": doc.get("generated_at") or ds.iso_now(),
            "document_species": template.get("document_species") or "narrative",
        }
    else:
        # Empty shells for THIS template only — do not keep previous doc
        # Preserve sketch if same template rebuild
        keep_sketch = doc.get("sketch") if doc.get("template_file") == template_file else None
        sections = {}
        for s in template.get("sections") or []:
            sk = s["section_key"]
            sections[sk] = (
                fm.init_section_shell(s) if fm.is_form_fill(template)
                else {"md": "", "files": []}
            )
        state["doc"] = {
            "template_file": template_file,
            "sections": sections,
            "generated_at": ds.iso_now(),
            "document_species": template.get("document_species") or "narrative",
        }
        if keep_sketch:
            state["doc"]["sketch"] = keep_sketch
    _ensure_sketch_fields(state, template, doc)
    return state


def save_drafted_template(folder, template):
    d = drafted_templates_dir(folder)
    d.mkdir(exist_ok=True)
    key = template.get("template_key") or "drafted"
    fname = f"{key}.json"
    (d / fname).write_text(json.dumps(template, indent=2, ensure_ascii=False), encoding="utf-8")
    return fname


def save_company_template(template):
    """L1 — owned template in local company library (inspectable)."""
    COMPANY_TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)
    key = template.get("template_key") or "imported"
    fname = f"{key}.json"
    template = {**template, "owner": "local", "import_status": "confirmed"}
    (COMPANY_TEMPLATES_DIR / fname).write_text(
        json.dumps(template, indent=2, ensure_ascii=False), encoding="utf-8")
    return fname


def regen_capabilities():
    try:
        import importlib.util
        spec = importlib.util.spec_from_file_location(
            "build_caps", ROOT / "scripts" / "build_caps.py")
        mod = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(mod)
        # Include company templates in the scan if build() accepts extra dirs
        if hasattr(mod, "build"):
            try:
                return mod.build(extra_dirs=[COMPANY_TEMPLATES_DIR])
            except TypeError:
                return mod.build()
    except Exception as e:
        return {"error": str(e)}
    return {}


def fact_aliases():
    return learning.merged_fact_aliases(fc.FACT_ALIASES)


def load_state(folder, template_file=None, project_id=None):
    p = state_path(folder)
    if p.exists():
        state = json.loads(p.read_text(encoding="utf-8"))
    else:
        state = ds.default_state()
    if project_id:
        state["project_id"] = project_id
    tf = template_file or state.get("active_template") or state.get("template")
    if tf:
        template = load_template(tf, folder)
        if template:
            if not (state.get("doc") and state["doc"].get("sections")):
                content = read_draft_content(folder, tf)
                if content:
                    state["_legacy_draft_md"] = content
            state = ds.migrate_state(state, template)
    return state


def isolated_conversation(state, project_id):
    """Return conversation turns belonging only to this project (BUGFIX_0.19 §A)."""
    return edchat.conversation_for_project(state or {}, project_id)


def save_state(folder, state):
    state_path(folder).write_text(json.dumps(state, indent=2, ensure_ascii=False), encoding="utf-8")


def primary_folder(p):
    """WORKORDER 0.61 — never IndexError on folder-less projects."""
    return pio.primary_folder(p)


def project_state_load(p, template_file=None):
    """Load state from folder or in-app memory (folder-less)."""
    folder = primary_folder(p)
    if folder is not None and folder.is_dir():
        return load_state(folder, template_file, project_id=p.get("id")), folder
    state = pio.load_memory_state(p.get("id") or "")
    if not state:
        state = ds.default_state()
    state["project_id"] = p.get("id")
    state["folderless"] = True
    tf = template_file or state.get("active_template") or state.get("template") or p.get("preferred_template")
    if tf:
        template = load_template(tf)
        if template:
            if not (state.get("doc") and (state["doc"] or {}).get("template_file") == tf):
                import template_lifecycle as tl
                if not state.get("doc"):
                    tl.create_document_shell(state, tf, template)
            ensure_doc_for_template(None, state, tf)
            _ensure_sketch_fields(state, template, state.get("doc"))
    return state, None


def project_state_save(p, state, folder=None):
    folder = folder if folder is not None else primary_folder(p)
    if folder is not None and Path(folder).is_dir():
        save_state(folder, state)
    else:
        state["folderless"] = True
        pio.save_memory_state(p.get("id") or "", state)


def capture_status_for_folder(folder):
    """Read-only Capture bridge status from project folder (.foldok/*)."""
    if folder is None or not Path(folder).is_dir():
        return None
    try:
        from local_app import capture_bridge as cbr

        return cbr.capture_status(folder)
    except Exception as exc:
        return {"error": str(exc)}


def create_folderless_project(name, template_file=None, *, output_format="pdf"):
    template = load_template(template_file) if template_file else None
    import template_lifecycle as tl
    return pio.create_folderless_project(
        name,
        template_file=template_file,
        template=template,
        output_format=output_format,
        load_projects=load_projects,
        save_projects=save_projects,
        create_document_shell=tl.create_document_shell,
        default_state=ds.default_state,
    )


# ── engine helpers ───────────────────────────────────────────────────
def looks_like_product_repo(folder):
    """True when the path is Foldok itself (or its engine), not a customer case."""
    return fpaths.is_product_tree(Path(folder))


def source_files(folders):
    """Recursive walk over every project folder. Returns (path, rel_name,
    cache_dir) triples; rel_name is the file's identity, prefixed with the
    folder name when the project spans several folders. Per-folder cache
    so a folder shared between projects indexes once, free everywhere."""
    out = []
    for folder in folders:
        root = Path(folder)
        cache_dir = fpaths.cache_dir(root)
        cache_dir.mkdir(exist_ok=True)
        prefix = f"{root.name}/" if len(folders) > 1 else ""
        for p in sorted(root.rglob("*")):
            if not p.is_file():
                continue
            rel_parts = p.relative_to(root).parts
            if any(part.startswith(".") for part in rel_parts):
                continue
            if any(part in SKIP_DIR_NAMES for part in rel_parts[:-1]):
                continue
            if len(rel_parts) >= 2 and rel_parts[0] == "Rapporter" and rel_parts[1] == "media":
                continue
            if p.name in OUTPUT_NAMES:
                continue
            out.append((p, prefix + p.relative_to(root).as_posix(), cache_dir))
    return out


def file_kind(path):
    ext = path.suffix.lower()
    if ext in fc.PHOTO_EXT:
        return "photo"
    if ext in fc.DOC_EXT:
        return "doc"
    if ext in getattr(fc, "CAD_EXT", set()):
        return "cad"
    return "skipped"


def file_rows(folders):
    import hashlib
    rows = []
    for p, rel, cache_dir in source_files(folders):
        kind = file_kind(p)
        sha = hashlib.sha256(p.read_bytes()).hexdigest()
        cache = cache_dir / f"{sha}.json"
        caption = None
        facts = 0
        if cache.exists():
            entry = fc.read_json_file(cache)
            caption = entry.get("caption")
            facts = len(entry.get("facts", []))
        rows.append({"name": rel, "size": p.stat().st_size, "kind": kind,
                     "indexed": cache.exists(), "caption": caption, "facts": facts})
    return rows


def load_index(folders, lang, user_facts=None, project_name=None, *, cache_only=False):
    """Load the project index; cached files are free.

    cache_only=True (page loads / gaps): never call Anthropic — skip uncached
    files so a dead API key or empty Anthropic balance cannot 500 the UI.
    """
    index = []
    for p, rel, cache_dir in source_files(folders or []):
        if cache_only:
            import hashlib
            try:
                sha = hashlib.sha256(p.read_bytes()).hexdigest()
            except OSError:
                continue
            cache = cache_dir / f"{sha}.json"
            if not cache.exists():
                continue
            try:
                entry = fc.read_json_file(cache)
                entry["cached"] = True
                entry.setdefault("file", rel)
                index.append(entry)
            except Exception:
                continue
            continue
        e = fc.index_file(p, lang, cache_dir, rel_name=rel)
        index.append(e)
    index = [e for e in index if e.get("kind") != "skipped"]
    name = (project_name or "").strip() or (Path(folders[0]).name if folders else "")
    if name and folders:
        cache_dir = fpaths.cache_dir(folders[0])
        try:
            synth = fc.index_project_name(name, cache_dir, lang)
            if synth:
                index = [synth] + [e for e in index if e.get("kind") != "project_name"]
        except Exception:
            pass
    return fc.inject_user_facts(index, user_facts or [])


def get_draft_md(folder, state, template, template_file):
    """Assembled draft from v2 state, or legacy file fallback."""
    ensure_doc_for_template(folder, state, template_file)
    if state.get("doc") and state["doc"].get("sections"):
        folders = [str(folder)] if folder else []
        full_index = load_index(
            folders, "no", state.get("user_facts"),
            project_name=(Path(folder).name if folder else ""),
            cache_only=True,
        ) if folders else []
        return ds.assemble_draft(state, template, state.get("artifact"), full_index=full_index)
    return read_draft_content(folder, template_file) if folder else None


def quarantine_document(folder, state, template_file, reason="WORKORDER_0.58"):
    """Move a stray document's drafts aside and drop it from the project ledger."""
    folder_path = Path(folder)
    stem = template_stem(template_file)
    qdir = fpaths.quarantine_dir(folder_path) / stem
    qdir.mkdir(parents=True, exist_ok=True)
    meta = {
        "template": template_file,
        "reason": reason,
        "quarantined_at": ds.iso_now(),
    }
    (qdir / "quarantine.json").write_text(
        json.dumps(meta, ensure_ascii=False, indent=2), encoding="utf-8")
    drafts = drafts_dir(folder_path)
    for ext in (".md", ".html"):
        src = drafts / f"{stem}{ext}"
        if src.exists():
            dest = qdir / src.name
            if dest.exists():
                dest.unlink()
            src.rename(dest)
    state["documents"] = [
        d for d in (state.get("documents") or [])
        if d.get("template") != template_file
    ]
    if (state.get("active_template") or "") == template_file or (
            (state.get("doc") or {}).get("template_file") == template_file):
        # Prefer another document if present
        nxt = None
        for d in state.get("documents") or []:
            if d.get("template"):
                nxt = d["template"]
                break
        if not nxt:
            for md in sorted(drafts.glob("*.md")) if drafts.is_dir() else []:
                nxt = f"{md.stem}.json"
                break
        if nxt:
            state["active_template"] = nxt
            state["template"] = nxt
            ensure_doc_for_template(folder, state, nxt)
        else:
            state["active_template"] = None
            state["template"] = None
            state["doc"] = None
            state["gaps"] = []
    return state


def sync_draft_files(folder, state, template, template_file, content):
    """Write draft.md, archive, and Rapporter export (+ HTML for form_fill)."""
    folder_path = Path(folder)
    drafts = drafts_dir(folder_path)
    drafts.mkdir(exist_ok=True)
    stem = template_stem(template_file)
    (drafts / f"{stem}.md").write_text(content, encoding="utf-8")
    (folder_path / "draft.md").write_text(content, encoding="utf-8")
    out_path, display = write_report_export(folder, template, content)
    # WORKORDER_0.29 §D — print-faithful HTML alongside markdown
    try:
        import form_model as fm
        if template and fm.is_form_fill(template):
            write_form_html_export(folder, template, state, display)
    except Exception:
        pass
    return out_path, display


def write_form_html_export(folder, template, state, display=None):
    """Write Rapporter/<name>.html from form_engine (print-faithful)."""
    import form_engine as fe
    folder_path = Path(folder)
    display = display or report_display_name(template)
    company = {}
    profile = (state or {}).get("company") or (state or {}).get("company_profile") or {}
    if isinstance(profile, dict) and profile.get("name"):
        company = {"name": profile["name"]}
    elif (state or {}).get("artifact", {}).get("owner"):
        company = {"name": state["artifact"]["owner"]}
    html_body = fe.export_form_html(
        template, state or {},
        artifact=(state or {}).get("artifact"),
        company=company,
        lang=(state or {}).get("lang") or "no",
    )
    out_dir = reports_dir(folder_path)
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / f"{safe_filename(display)}.html"
    out_path.write_text(html_body, encoding="utf-8")
    drafts = drafts_dir(folder_path)
    drafts.mkdir(exist_ok=True)
    stem = template_stem(template.get("file") or template.get("template_key") or display)
    # Prefer explicit template file stem when available via state active_template
    active = (state or {}).get("active_template") or ""
    if active.endswith(".json"):
        stem = template_stem(active)
    (drafts / f"{stem}.html").write_text(html_body, encoding="utf-8")
    return out_path


def estimate_cost(rows):
    """Rough per-file estimate from observed runs (€0.014/doc, €0.008/photo)."""
    docs = sum(1 for r in rows if not r["indexed"] and r["kind"] == "doc")
    photos = sum(1 for r in rows if not r["indexed"] and r["kind"] == "photo")
    return round(docs * 0.02 + photos * 0.01, 2)


def reports_dir(folder):
    return Path(folder) / "Rapporter"


def safe_filename(name):
    for ch in '\\/:*?"<>|':
        name = name.replace(ch, "-")
    return name.strip()[:120] or "rapport"


def report_display_name(template):
    """Short Norwegian label for filenames, e.g. 'Konstruksjonsrapport'."""
    name = template.get("name_no") or template.get("name") or template.get("template_key", "rapport")
    return name.split(" / ")[0].strip()


def write_report_export(folder, template, content):
    """Save a human-named copy under <project>/Rapporter/. Returns path + display name."""
    folder_path = Path(folder)
    display = report_display_name(template)
    out_dir = reports_dir(folder_path)
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / f"{safe_filename(display)}.md"
    out_path.write_text(content, encoding="utf-8")
    return out_path, display


def drafts_dir(folder):
    return fpaths.drafts_dir(folder)


def _doc_export_meta(folder, tpl, t_meta):
    """export_name + export_path for a template, if the Rapporter/ file exists."""
    if not t_meta or not folder:
        return (report_display_name(t_meta) if t_meta else None), None
    display = report_display_name(t_meta)
    out_path = reports_dir(folder) / f"{safe_filename(display)}.md"
    return display, str(out_path) if out_path.exists() else None


def _enrich_doc(folder, doc, by_file, state=None):
    tpl = doc.get("template")
    t = by_file.get(tpl, {})
    display, export_path = _doc_export_meta(folder, tpl, t)
    doc.setdefault("name_no", t.get("name_no") or t.get("name") or doc.get("key", ""))
    if display:
        doc["export_name"] = display
    if export_path:
        doc["export_path"] = export_path
    # WORKORDER 0.60 — payment / draft status chip
    try:
        blocking = int(doc.get("blocking") or 0)
        st_for_doc = state
        if state and (state.get("doc") or {}).get("template_file") not in (None, tpl):
            # Fingerprint only when this doc is the loaded shell
            st_for_doc = {"doc": {"sections": {}, "sketch": {}}}
            if (state.get("doc") or {}).get("template_file") == tpl:
                st_for_doc = state
        elif state and (state.get("active_template") == tpl or state.get("template") == tpl):
            st_for_doc = state
        doc["status"] = acct.document_status(
            doc, blocking_gaps=blocking, state=st_for_doc)
        tier, price = acct.export_price_for_template(t, hub.load_capabilities())
        doc["export_tier"] = tier
        doc["export_price_eur"] = price
        ent = acct.export_entitlement(doc, st_for_doc)
        doc["export_entitlement"] = ent
    except Exception:
        doc.setdefault("status", {"key": "draft", "label": "○ Utkast", "class": "st-draft"})
    return doc


def template_stem(template_file):
    return template_file.replace(".json", "")


def draft_path(folder, template_file):
    return drafts_dir(folder) / f"{template_stem(template_file)}.md"


def read_draft_content(folder, template_file):
    """Load draft text from archive, export folder, or legacy draft.md."""
    folder_path = Path(folder)
    src = draft_path(folder_path, template_file)
    if src.exists():
        return src.read_text(encoding="utf-8")
    tpath = ROOT / "templates" / template_file
    if tpath.exists():
        t = json.loads(tpath.read_text(encoding="utf-8"))
        exported = reports_dir(folder_path) / f"{safe_filename(report_display_name(t))}.md"
        if exported.exists():
            return exported.read_text(encoding="utf-8")
    state = load_state(folder)
    lone = folder_path / "draft.md"
    if lone.exists() and state.get("template") == template_file:
        return lone.read_text(encoding="utf-8")
    return None


def list_documents(folder, state, templates):
    """Built reports for this project — from state + on-disk .foldok_drafts/."""
    by_file = {t["file"]: t for t in templates}
    docs, seen = [], set()
    drafts = drafts_dir(folder)

    # migrate lone draft.md into archive so prior runs are not lost on next generate
    lone = Path(folder) / "draft.md"
    if lone.exists() and state.get("template"):
        drafts.mkdir(exist_ok=True)
        stem = template_stem(state["template"])
        archived = drafts / f"{stem}.md"
        if not archived.exists():
            archived.write_text(lone.read_text(encoding="utf-8"), encoding="utf-8")

    for d in state.get("documents", []):
        tpl = d.get("template")
        if not tpl:
            continue
        stem = d.get("key") or template_stem(tpl)
        if not (drafts / f"{stem}.md").exists():
            continue
        seen.add(tpl)
        docs.append(_enrich_doc(folder, {**d, "key": stem}, by_file, state))

    if drafts.is_dir():
        for md in sorted(drafts.glob("*.md"), key=lambda p: p.stat().st_mtime, reverse=True):
            tpl = f"{md.stem}.json"
            if tpl in seen:
                continue
            seen.add(tpl)
            t = by_file.get(tpl, {})
            docs.append(_enrich_doc(folder, {
                "template": tpl, "key": md.stem,
                "name_no": t.get("name_no") or t.get("name") or md.stem,
                "gaps": 0, "blocking": 0,
                "generated_at": datetime.fromtimestamp(md.stat().st_mtime, tz=timezone.utc).isoformat(),
            }, by_file, state))

    # Shell docs without draft yet (sketch / form create) — still list them
    for d in state.get("documents", []) or []:
        tpl = d.get("template")
        if not tpl or tpl in seen:
            continue
        stem = d.get("key") or template_stem(tpl)
        seen.add(tpl)
        docs.append(_enrich_doc(folder, {**d, "key": stem}, by_file, state))

    # legacy: only draft.md + template in state
    if not docs and state.get("template") and (Path(folder) / "draft.md").exists():
        tpl = state["template"]
        stem = template_stem(tpl)
        t = by_file.get(tpl, {})
        docs.append(_enrich_doc(folder, {
            "template": tpl, "key": stem,
            "name_no": t.get("name_no") or t.get("name") or stem,
            "gaps": len(state.get("gaps", [])),
            "blocking": sum(1 for g in state.get("gaps", []) if g.get("severity") == "blocking"),
            "generated_at": datetime.fromtimestamp((Path(folder) / "draft.md").stat().st_mtime,
                                                     tz=timezone.utc).isoformat(),
        }, by_file, state))

    docs.sort(key=lambda d: d.get("generated_at") or "", reverse=True)
    # backfill Rapporter/ for drafts that predate auto-export
    for doc in docs:
        if doc.get("export_path"):
            continue
        tpl = doc.get("template")
        src = draft_path(folder, tpl)
        if not tpl or not src.exists():
            continue
        tpath = ROOT / "templates" / tpl
        if not tpath.exists():
            continue
        t = json.loads(tpath.read_text(encoding="utf-8"))
        out_path, display = write_report_export(folder, t, src.read_text(encoding="utf-8"))
        doc["export_name"] = display
        doc["export_path"] = str(out_path)
    return docs


def _materialize_sketch_sections(state, sk):
    """Write sketch placeholders into doc.sections for assemble/export."""
    doc = state.setdefault("doc", {"sections": {}})
    phs = ((doc.get("sketch") or {}).get("placeholders") or [])
    built = sk.placeholders_to_sections_md(phs)
    sections = doc.setdefault("sections", {})
    # Keep canvas shell; merge sketch sections
    for sk_key, sec in built.items():
        sections[sk_key] = {
            **(sections.get(sk_key) or {}),
            "md": sec.get("md") or "",
            "files": sec.get("files") or [],
            "title_override": sec.get("title_override"),
            "sketch_id": sec.get("sketch_id"),
        }
    # Assembled preview on canvas
    parts = []
    for ph in sk.sort_placeholders(phs):
        title = ph.get("label") or ph.get("prompt") or ph.get("type")
        md = ph.get("md") or f"*({ph.get('prompt') or 'tom'})*"
        parts.append(f"## {title}\n\n{md}\n")
    sections.setdefault("canvas", {"md": "", "files": []})
    sections["canvas"]["md"] = "\n".join(parts)


def fm_is_form(template):
    try:
        import form_model as fm
        return fm.is_form_fill(template)
    except Exception:
        return False


def templates_list(project=None, *, tags=None):
    """List system + company templates, filtered for project domain.

    Domain-locked vehicle fixtures (sample_multipoint, …) are hidden unless
    the project carries a vehicle tag (see form_model.filter_templates_for_project).
    """
    import form_model as fm
    out = []
    seen = set()
    dirs = [ROOT / "templates", COMPANY_TEMPLATES_DIR]
    for d in dirs:
        if not d.is_dir():
            continue
        for p in sorted(d.glob("*.json")):
            if p.name in seen:
                continue
            seen.add(p.name)
            t = json.loads(p.read_text(encoding="utf-8"))
            out.append({"file": p.name, "key": t.get("template_key"),
                        "name_no": t.get("name_no") or t.get("name"), "name": t.get("name"),
                        "description": (t.get("description") or "")[:400],
                        "applies_to": t.get("applies_to") or [],
                        "document_species": t.get("document_species") or "narrative",
                        "origin": t.get("origin") or (
                            "owned" if d == COMPANY_TEMPLATES_DIR else "system"
                        ),
                        "badge": t.get("badge"),
                        "system_default": bool(t.get("system_default")),
                        "owned": d == COMPANY_TEMPLATES_DIR or t.get("origin") == "sketched",
                        "section_count": len(t.get("sections") or []),
                        "sections": [{"section_key": s.get("section_key"),
                                      "title_no": s.get("title_no"), "title": s.get("title"),
                                      "block_type": s.get("block_type"),
                                      "columns": s.get("columns") or 1,
                                      "fields": s.get("fields") or []}
                                     for s in t.get("sections", [])]})
    return fm.filter_templates_for_project(out, project, tags=tags)


def template_file_for_key(template_key):
    if not template_key:
        return None
    for t in templates_list():
        if t.get("key") == template_key:
            return t["file"]
    return None


def active_suggestions(state, folders, template_file=None):
    """Zero-token completeness suggestions (dismissals persisted in state)."""
    if not state.get("artifact"):
        return []
    index = load_index(folders, "no", state.get("user_facts"), cache_only=True)
    template = load_template(template_file) if template_file else None
    tkey = (template or {}).get("template_key")
    dismissed = {d.get("name") for d in state.get("dismissed_suggestions", [])}
    raw = be.detect_suggestions(index, state.get("artifact"), current_template_key=tkey)
    excluded = ds.excluded_source_files(state)
    raw.extend(fc.detect_revision_supersede(index, excluded_files=excluded, folders=folders))
    out = []
    for s in raw:
        if s.get("name") in dismissed:
            continue
        if s.get("type") == "supersede_revision":
            s = {**s, "name": s.get("name") or f"supersede|{s.get('drawing_no')}|{s.get('file')}"}
        out.append(s)
    return out


def load_active_index(state, folders, lang="no", *, cache_only=True):
    """Index minus sources toggled off (document exclude + project source_selection).

    Default cache_only=True so UI page loads never trigger live Anthropic calls.
    """
    index = load_index(folders, lang, state.get("user_facts"), cache_only=cache_only)
    excluded = set(ds.excluded_source_files(state))
    sel = state.get("source_selection")
    if sel:
        for e in index:
            rel = e.get("file") or ""
            if rel and not prescan.source_is_enabled(rel, sel):
                excluded.add(rel)
    if not excluded:
        return index
    return [e for e in index if e.get("file") not in excluded]


def _store_table_section(state, section_key, data, lang="no"):
    """Apply sovereign cell overrides, render md, store both md + structure."""
    data = fc.apply_cell_overrides(data, state.get("cell_overrides"), section_key)
    text = fc.render_table_md(data, lang)
    sec = state["doc"]["sections"].setdefault(section_key, {"md": "", "files": []})
    sec["md"] = text
    sec["table"] = data
    sec["updated"] = ds.iso_now()
    return sec


def refresh_doc_control_section(state, folders, template_file, lang="no"):
    """Re-render doc_control from index — zero model calls."""
    template = load_template(template_file)
    if not template or not state.get("doc"):
        return False
    if not any(s["section_key"] == "doc_control" for s in template.get("sections", [])):
        return False
    index = load_active_index(state, folders, lang)
    data = fc.compile_doc_control_data(index, state.get("artifact"), lang)
    _store_table_section(state, "doc_control", data, lang)
    return True


def refresh_spec_overview_section(state, folders, template_file, lang="no"):
    """Re-render spec_overview from index — zero model calls."""
    template = load_template(template_file)
    if not template or not state.get("doc"):
        return False
    if not any(s["section_key"] == "spec_overview" for s in template.get("sections", [])):
        return False
    index = load_active_index(state, folders, lang)
    data = fc.compile_spec_overview_data(index, state.get("artifact"), lang)
    _store_table_section(state, "spec_overview", data, lang)
    return True


def refresh_code_tables(state, folders, template_file, lang="no"):
    """Re-run every zero-token compiler this template has."""
    changed = refresh_doc_control_section(state, folders, template_file, lang)
    changed = refresh_spec_overview_section(state, folders, template_file, lang) or changed
    template = load_template(template_file)
    if template and state.get("doc"):
        sec = (state["doc"].get("sections") or {}).get("drawings_register")
        if sec is not None and any(s["section_key"] == "drawings_register"
                                   for s in template.get("sections", [])):
            index = load_active_index(state, folders, lang)
            sec["md"] = fc.compile_drawings_register(index, lang)
            sec["updated"] = ds.iso_now()
            changed = True
    return changed


def excluded_fact_ids(state, folders, lang="no"):
    """Fact ids belonging to sources toggled off for this document."""
    import source_citations as sc
    excluded = ds.excluded_source_files(state)
    if not excluded:
        return set()
    full = load_index(folders, lang, state.get("user_facts"))
    out = set()
    for f in excluded:
        out |= sc.fact_ids_for_file(full, f)
    return out


def refresh_bom_section(state, folders, template_file, lang="no"):
    """Re-render bom section from cached index — zero model calls."""
    template = load_template(template_file)
    if not template or not state.get("doc"):
        return False
    bom_sec = next((s for s in template.get("sections", []) if s["section_key"] == "bom"), None)
    if not bom_sec:
        return False
    cond = bom_sec.get("condition")
    if cond:
        holds, recognized = fc._condition_holds(cond, state.get("artifact") or {})
        if recognized and not holds:
            return False
    index = load_active_index(state, folders, lang)
    mapping = {"section": bom_sec, "files": [], "fact_ids": []}
    text = fc.generate_section("bom", mapping, index, state.get("artifact"), lang)
    text, cited, _violations = fc.postprocess(
        "bom", text, index, state.get("artifact"),
        excluded_fact_ids=excluded_fact_ids(state, folders, lang))
    # WO 0.22 C3 — merge photo-derived component rows from document state
    try:
        import bom_engine as be
        text = be.merge_bom_markdown(text, state.get("bom_components") or [], lang)
    except Exception:
        pass
    sec = state["doc"]["sections"].setdefault("bom", {"md": "", "files": []})
    sec["md"] = text
    sec["cited_fact_ids"] = cited
    # Attach image files referenced by components
    img_files = [c.get("file") for c in (state.get("bom_components") or []) if c.get("file")]
    if img_files:
        existing = list(sec.get("files") or [])
        for f in img_files:
            if f not in existing:
                existing.append(f)
        sec["files"] = existing[:12]
    sec["updated"] = ds.iso_now()
    return True


def store_connection_diagram(state, spec, svg, *, lang="no"):
    """Persist confirmed connection_spec + SVG into document section."""
    import connection_diagram as cdiag
    if not state.get("doc"):
        state["doc"] = {"sections": {}}
    sections = state["doc"].setdefault("sections", {})
    md = cdiag.embed_svg_markdown(svg, lang=lang)
    sec = sections.setdefault("connection_diagram", {"md": "", "files": []})
    sec["md"] = md
    sec["block_type"] = cdiag.BLOCK_TYPE
    sec["connection_spec"] = spec
    sec["svg"] = svg
    sec["updated"] = ds.iso_now()
    # Collect images for KILDER / hover trace
    imgs = [c.get("image") for c in (spec.get("components") or []) if c.get("image")]
    if imgs:
        existing = list(sec.get("files") or [])
        for f in imgs:
            if f not in existing:
                existing.append(f)
        sec["files"] = existing[:12]
    state["connection_spec"] = spec
    ds.add_version(state, "user", "connection_diagram", "Confirmed connection block diagram")
    return sec


def ensure_template_sections(state, template, artifact=None):
    """Add section shells for template sections that apply to this artifact."""
    if not template or not state.get("doc"):
        return False
    artifact = artifact or state.get("artifact") or {}
    sections = state["doc"].setdefault("sections", {})
    changed = False
    for s in template.get("sections", []):
        sk = s["section_key"]
        cond = s.get("condition")
        if cond:
            holds, recognized = fc._condition_holds(cond, artifact)
            if recognized and not holds:
                continue
        if sk not in sections:
            sections[sk] = {"md": "", "files": []}
            changed = True
    return changed


# ── background jobs ──────────────────────────────────────────────────
def start_job(target, *args, **job_extra):
    job_id = uuid.uuid4().hex[:8]
    with LOCK:
        JOBS[job_id] = {
            "id": job_id, "status": "running", "step": "", "done": 0, "total": 0,
            "detail": "", "cost_eur": 0.0, "error": None, "started_at": time.time(),
            "eta_s": None,
            "cancel_requested": False,
            "pause_requested": False,
            "paused": False,
            "last_heartbeat": time.time(),
            "spent_eur": 0.0,
            "budget_eur": None,
            "scope": None,
            "current_folder": "",
            "kind": "generic",
        }
        if job_extra:
            JOBS[job_id].update(job_extra)
    t = threading.Thread(target=_job_wrapper, args=(job_id, target) + args, daemon=True)
    t.start()
    return job_id


def _job_wrapper(job_id, target, *args):
    ledger_start = len(fc.LEDGER)
    try:
        target(job_id, *args)
        with LOCK:
            j = JOBS.get(job_id) or {}
            if j.get("cancel_requested"):
                JOBS[job_id]["status"] = "cancelled"
            elif j.get("paused") or j.get("pause_requested"):
                JOBS[job_id]["status"] = "paused"
            elif j.get("status") == "running":
                JOBS[job_id]["status"] = "done"
    except Exception as e:
        traceback.print_exc()
        with LOCK:
            JOBS[job_id]["status"] = "error"
            JOBS[job_id]["error"] = str(e)
    finally:
        cost = sum(l["eur"] for l in fc.LEDGER[ledger_start:])
        with LOCK:
            JOBS[job_id]["cost_eur"] = round(cost, 3)
            JOBS[job_id]["spent_eur"] = round(cost, 3)


def job_update(job_id, **kw):
    with LOCK:
        if job_id in JOBS:
            JOBS[job_id].update(kw)


def job_snapshot(job_id):
    with LOCK:
        j = JOBS.get(job_id)
        return dict(j) if j else None


def request_cancel(job_id) -> bool:
    with LOCK:
        j = JOBS.get(job_id)
        if not j:
            return False
        j["cancel_requested"] = True
        return True


def request_pause(job_id, reason: str = "") -> bool:
    with LOCK:
        j = JOBS.get(job_id)
        if not j:
            return False
        j["pause_requested"] = True
        if reason:
            j["pause_reason"] = reason
        return True


def heartbeat(job_id) -> bool:
    with LOCK:
        j = JOBS.get(job_id)
        if not j:
            return False
        j["last_heartbeat"] = time.time()
        return True


def _job_should_stop(job_id) -> str | None:
    """Return 'cancel' | 'pause' | 'budget' | None."""
    with LOCK:
        j = JOBS.get(job_id) or {}
        if j.get("cancel_requested"):
            return "cancel"
        if j.get("pause_requested"):
            return "pause"
        last = j.get("last_heartbeat") or j.get("started_at") or time.time()
        if j.get("kind") == "index" and (time.time() - last) > HEARTBEAT_TIMEOUT_S:
            j["pause_requested"] = True
            j["pause_reason"] = "Pauset — ingen klient tilkoblet"
            return "pause"
        budget = j.get("budget_eur")
        spent = j.get("spent_eur") or 0.0
        if budget is not None and spent >= float(budget):
            j["pause_requested"] = True
            j["pause_reason"] = f"Budsjettgrensen på €{budget:g} er nådd"
            return "budget"
    return None


def _refresh_job_spend(job_id, ledger_start: int):
    spent = sum(l["eur"] for l in fc.LEDGER[ledger_start:])
    job_update(job_id, spent_eur=round(spent, 3), cost_eur=round(spent, 3))


def run_index(job_id, folders, lang, scope=None):
    """Chunked, cancellable, budget-aware indexing (WORKORDER 0.55 §C)."""
    scope = scope or {}
    mode = scope.get("mode") or "all"
    subfolders = scope.get("subfolders")
    newest_n = scope.get("newest_n")
    budget = scope.get("budget_eur")
    if budget is None:
        # project setting
        try:
            st0 = load_state(folders[0]) if folders else {}
            budget = st0.get("index_budget_eur", DEFAULT_INDEX_BUDGET_EUR)
        except Exception:
            budget = DEFAULT_INDEX_BUDGET_EUR

    job_update(
        job_id, kind="index", step="Forbereder", scope=scope,
        budget_eur=float(budget) if budget is not None else None,
        last_heartbeat=time.time(),
    )

    all_files = source_files(folders)
    # Drop already-cached (free) before chunking spend
    pending = []
    for p, rel, cd in all_files:
        try:
            sha = __import__("hashlib").sha256(p.read_bytes()).hexdigest()
            if (cd / f"{sha}.json").exists():
                continue
        except OSError:
            continue
        pending.append((p, rel, cd))

    pending = prescan.filter_pending(
        pending, mode=mode, subfolders=subfolders, newest_n=newest_n,
        skip_oversize=True,
        disabled_files=(scope.get("disabled_files")
                        or (scope.get("source_selection") or {}).get("disabled_files")),
        disabled_folders=(scope.get("disabled_folders")
                          or (scope.get("source_selection") or {}).get("disabled_folders")),
    )

    total = len(pending)
    job_update(job_id, step="Indekserer", total=total, done=0)
    if total == 0:
        job_update(job_id, detail="Ingenting å indeksere (alt i cache / scope tom)")
        _commit_index_manifest(job_id, folders)
        return

    ledger_start = len(fc.LEDGER)
    done = 0
    t0 = time.time()
    stop_reason = None

    for chunk_start in range(0, total, INDEX_CHUNK_SIZE):
        stop_reason = _job_should_stop(job_id)
        if stop_reason:
            break
        chunk = pending[chunk_start: chunk_start + INDEX_CHUNK_SIZE]
        # Process chunk with limited workers; check cancel between completions
        with ThreadPoolExecutor(max_workers=5) as pool:
            futs = {}
            for p, rel, cd in chunk:
                stop_reason = _job_should_stop(job_id)
                if stop_reason:
                    break
                futs[pool.submit(fc.index_file, p, lang, cd, rel)] = rel
            for fut in as_completed(futs):
                try:
                    fut.result()
                except Exception:
                    traceback.print_exc()
                done += 1
                rel = futs[fut]
                folder_hint = rel.split("/", 1)[0] if "/" in rel else ""
                _refresh_job_spend(job_id, ledger_start)
                elapsed = max(0.001, time.time() - t0)
                remaining = total - done
                eta = int((elapsed / done) * remaining) if done else None
                job_update(
                    job_id, done=done,
                    eta_s=eta if remaining else 0,
                    detail=f"{rel}" + (f" · ~{eta}s igjen" if eta else ""),
                    current_folder=folder_hint,
                )
                stop_reason = _job_should_stop(job_id)
                if stop_reason:
                    # Cancel pending futures in this chunk (running ones finish ≤1)
                    for f in futs:
                        f.cancel()
                    break
        if stop_reason:
            break

    _refresh_job_spend(job_id, ledger_start)
    spent = job_snapshot(job_id).get("spent_eur") or 0.0

    if stop_reason == "cancel":
        job_update(
            job_id, status="cancelled", paused=False,
            detail=(
                f"Stoppet — {done} av {total} filer indeksert (€{spent:.2f}). "
                f"Fortsett når du vil."
            ),
            resume_remaining=total - done,
        )
    elif stop_reason in ("pause", "budget"):
        snap = job_snapshot(job_id) or {}
        reason = snap.get("pause_reason") or "Pauset"
        job_update(
            job_id, status="paused", paused=True, pause_requested=True,
            detail=f"{reason} — {done} av {total} filer (€{spent:.2f}).",
            resume_remaining=total - done,
        )
    else:
        # Persist throughput for next pre-scan estimate
        if done > 0 and folders:
            elapsed = max(0.001, time.time() - t0)
            thr = done / elapsed / 5.0
            try:
                st = load_state(folders[0])
                st["index_last_throughput"] = round(thr, 4)
                st["index_scope"] = scope
                save_state(folders[0], st)
            except Exception:
                pass
        job_update(job_id, detail=f"Ferdig — {done} filer (€{spent:.2f})")

    _commit_index_manifest(job_id, folders)


def _commit_index_manifest(job_id, folders):
    try:
        primary = folders[0] if folders else None
        if primary:
            result = idxtools.commit_manifest_after_index(
                primary, folders, fc, source_files)
            job_update(job_id, index_diff=result, index_version=result.get("index_version"))
    except Exception as e:
        job_update(job_id, index_manifest_error=str(e)[:200])


def run_artifact(job_id, folders, lang):
    job_update(job_id, step="Bygger artefaktmodell", total=1, detail="Checkpoint A (Sonnet)")
    index = load_index(folders, lang)
    # Index always includes (prosjektnavn); allow build from name alone
    artifact = fc.build_artifact_model(index, lang)
    state = load_state(folders[0])
    state["artifact"] = artifact
    state["confirmed"] = False
    save_state(folders[0], state)
    job_update(job_id, done=1)


def run_form_fill(job_id, folders, template_file, lang):
    """WORKORDER_0.29 — form_fill: prefill from index, zero model calls."""
    import form_model as fm
    import template_lifecycle as tl
    folder = folders[0]
    state = load_state(folder, template_file)
    template = load_template(template_file, folder)
    if not template or not fm.is_form_fill(template):
        raise RuntimeError(f"Ikke et form_fill-skjema: {template_file}")
    index = load_index(folders, lang, state.get("user_facts"),
                       project_name=Path(folder).name)
    job_update(job_id, step="Forhåndsutfyller skjema", total=1, detail="0 tokens")
    tl.create_document_shell(state, template_file, template)
    pref = fm.prefill_form(state, template, index)
    state["gaps"] = ds.gaps_for_document(state, template, index, state.get("artifact") or {}, fc)
    content = ds.assemble_draft(state, template, state.get("artifact"))
    export_path, export_name = sync_draft_files(folder, state, template, template_file, content)
    blocking = sum(1 for g in state["gaps"] if g.get("severity") == "blocking")
    stem = template_stem(template_file)
    doc = {"template": template_file, "key": stem,
           "name_no": template.get("name_no") or template.get("name") or stem,
           "document_species": "form_fill",
           "export_name": export_name, "export_path": str(export_path),
           "gaps": len(state["gaps"]), "blocking": blocking,
           "gap_list": state["gaps"],
           "generated_at": datetime.now(timezone.utc).isoformat()}
    docs = [d for d in state.get("documents", []) if d.get("template") != template_file]
    docs.append(doc)
    state["template"] = template_file
    state["active_template"] = template_file
    state["documents"] = docs
    state["violations"] = []
    ds.add_version(state, "system", "doc",
                   f"Skjema {export_name} (prefill={pref.get('prefilled', 0)})", section=None)
    job_update(job_id, done=1, detail=f"prefilled={pref.get('prefilled', 0)}")
    save_state(folder, state)


def run_generate(job_id, folders, template_file, lang):
    folder = folders[0]
    state = load_state(folder, template_file)
    artifact = state.get("artifact")
    template = load_template(template_file, folder)
    if not template:
        raise RuntimeError(f"Mal ikke funnet: {template_file}")
    import form_model as fm
    if fm.is_form_fill(template):
        return run_form_fill(job_id, folders, template_file, lang)
    if not artifact:
        raise RuntimeError("Artefaktmodellen mangler — beskriv prosjektet i chatten først")
    if not state.get("confirmed"):
        raise RuntimeError("Artefaktmodellen må bekreftes før generering (checkpoint A)")
    index = load_index(folders, lang, state.get("user_facts"),
                       project_name=Path(folder).name)

    job_update(job_id, step="Kartlegger seksjoner", detail="Checkpoint B")
    mappings, gaps = fc.map_sections(template, index, artifact)
    pre_gaps = fc.template_gaps(template, index, artifact)

    job_update(job_id, step="Genererer seksjoner", total=len(mappings))
    sections_data, all_violations = [], []
    t0 = time.time()
    for n, (sec_key, mapping) in enumerate(mappings.items(), 1):
        s = mapping["section"]
        title = s.get("title_no") if lang == "no" else s.get("title")
        title = title or s.get("title") or sec_key
        remaining = len(mappings) - (n - 1)
        eta = int((time.time() - t0) / max(n - 1, 1) * remaining) if n > 1 else None
        detail = f"{title}" + (f" · ~{eta}s igjen" if eta is not None else "")
        job_update(job_id, done=n - 1, detail=detail, eta_s=eta)
        text = fc.generate_section_with_structure(sec_key, mapping, index, artifact, lang, pre_gaps=pre_gaps)
        text, cited, violations = fc.postprocess(sec_key, text, index, artifact)
        if violations:
            text2 = fc.generate_section_with_structure(sec_key, mapping, index, artifact, lang)
            text2, cited2, violations2 = fc.postprocess(sec_key, text2, index, artifact)
            if len(violations2) < len(violations):
                text, cited, violations = text2, cited2, violations2
            if violations:
                text = fc.redact_uncited(text)
                all_violations.append(sec_key)
        sections_data.append((sec_key, text, cited, violations, mapping.get("files", [])))
    job_update(job_id, done=len(mappings), detail="", eta_s=0)

    state["doc"] = ds.build_doc_from_generation(template_file, sections_data)
    # Prefer model {{fig:}} / structure fallbacks; fill remaining required_media gaps
    enrich_section_media(state, folders, template)
    ensure_figures_in_doc(state, folders, template)
    state["gaps"] = ds.gaps_for_document(state, template, index, artifact, fc)
    content = ds.assemble_draft(state, template, artifact)
    content = materialize_export_figures(folder, content)
    export_path, export_name = sync_draft_files(folder, state, template, template_file, content)

    blocking = sum(1 for g in state["gaps"] if g.get("severity") == "blocking")
    stem = template_stem(template_file)
    doc = {"template": template_file, "key": stem,
           "name_no": template.get("name_no") or template.get("name") or stem,
           "export_name": export_name, "export_path": str(export_path),
           "gaps": len(state["gaps"]), "blocking": blocking,
           "gap_list": state["gaps"],
           "generated_at": datetime.now(timezone.utc).isoformat()}
    docs = [d for d in state.get("documents", []) if d.get("template") != template_file]
    docs.append(doc)

    state["template"] = template_file
    state["active_template"] = template_file
    state["documents"] = docs
    state["violations"] = all_violations
    ds.add_version(state, "system", "doc", f"Genererte {export_name}", section=None)
    save_state(folder, state)


def run_regenerate_section(job_id, folders, template_file, section_key, lang, instruction=None):
    folder = folders[0]
    state = load_state(folder, template_file)
    artifact = state.get("artifact")
    template = load_template(template_file)
    index = load_index(folders, lang, state.get("user_facts"))
    mappings, _ = fc.map_sections(template, index, artifact)
    if section_key not in mappings:
        raise RuntimeError(f"Ukjent seksjon: {section_key}")
    old_md = state.get("doc", {}).get("sections", {}).get(section_key, {}).get("md", "")
    job_update(job_id, step="Regenererer seksjon", total=1, detail=section_key)
    result = fc.regenerate_one_section(section_key, mappings[section_key], index, artifact, lang,
                                       instruction=instruction)
    job_update(job_id, done=1, result={"new_md": result["md"], "old_md": old_md,
                                      "section_key": section_key, "cited": result["cited"],
                                      "violations": result["violations"]})


def find_source(folders, rel_name):
    for p, rel, cache_dir in source_files(folders):
        if rel == rel_name:
            return p, cache_dir
    return None, None


def open_path_in_os(path, reveal=False):
    """Open a file/folder with the OS default app (or reveal in Explorer)."""
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(str(path))
    if sys.platform == "win32":
        import subprocess
        if reveal:
            subprocess.Popen(["explorer", "/select,", str(path)])
        else:
            os.startfile(str(path))  # noqa: S606 — local workbench only
        return
    if sys.platform == "darwin":
        import subprocess
        if reveal:
            subprocess.Popen(["open", "-R", str(path)])
        else:
            subprocess.Popen(["open", str(path)])
        return
    import subprocess
    subprocess.Popen(["xdg-open", str(path)])


def load_cache_entry(folders, rel_name):
    path_f, cache_dir = find_source(folders, rel_name)
    if not path_f or not cache_dir:
        return None, None, None
    import hashlib
    st = path_f.stat()
    # Fast path: reuse cache keyed by path+mtime+size (avoid re-hashing 20MB PDFs on every view)
    fast = cache_dir / f"path-{hashlib.sha256(f'{rel_name}:{st.st_size}:{st.st_mtime_ns}'.encode()).hexdigest()[:16]}.json"
    if fast.exists():
        entry = fc.read_json_file(fast)
        return path_f, cache_dir, entry
    sha = hashlib.sha256(path_f.read_bytes()).hexdigest()
    cache = cache_dir / f"{sha}.json"
    entry = fc.read_json_file(cache) if cache.exists() else None
    if entry and not fast.exists():
        try:
            fast.write_text(json.dumps(entry, ensure_ascii=False), encoding="utf-8")
        except Exception:
            pass
    return path_f, cache_dir, entry


DRAWING_ROLES = fc.DRAWING_ROLES
DRAWING_NAME = fc.DRAWING_NAME


def is_visual_file(path_or_name, entry=None):
    name = path_or_name.name if isinstance(path_or_name, Path) else str(path_or_name)
    return fc.is_visual_source(name, entry)


def pptx_cutout_count(path_f, cache_dir):
    """Extract embedded images from PPTX into previews/; return cutout count."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import hashlib
    from io import BytesIO
    sha = hashlib.sha256(path_f.read_bytes()).hexdigest()[:16]
    prev_dir = cache_dir / "previews"
    prev_dir.mkdir(exist_ok=True)
    prs = Presentation(str(path_f))
    n = 0
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            out = prev_dir / f"{sha}-p{n}.jpg"
            if not out.exists():
                blob = shape.image.blob
                try:
                    from PIL import Image
                    im = Image.open(BytesIO(blob))
                    if im.mode not in ("RGB", "L"):
                        im = im.convert("RGB")
                    im.thumbnail((1200, 1600))
                    im.save(out, format="JPEG", quality=85)
                except Exception:
                    out.write_bytes(blob)
            n += 1
    # If no embedded pictures, still one placeholder page via first slide note
    return max(n, 1 if len(prs.slides) else 0)


def media_page_count(path_f, cache_dir=None):
    ext = path_f.suffix.lower()
    if ext in fc.PHOTO_EXT:
        return 1
    if ext == ".pdf":
        return pdf_page_count(path_f)
    if ext == ".pptx" and cache_dir is not None:
        return pptx_cutout_count(path_f, cache_dir)
    return 1


def render_media_preview(path_f, cache_dir, page=0, max_w=1200):
    """Return JPEG bytes for a photo, PDF page, or PPTX cutout."""
    import hashlib
    from io import BytesIO
    sha = hashlib.sha256(path_f.read_bytes()).hexdigest()[:16]
    prev_dir = cache_dir / "previews"
    prev_dir.mkdir(exist_ok=True)
    out = prev_dir / f"{sha}-p{page}.jpg"
    if out.exists() and out.stat().st_size > 0:
        return out.read_bytes(), "image/jpeg"

    ext = path_f.suffix.lower()
    if ext in fc.PHOTO_EXT:
        from PIL import Image
        im = Image.open(path_f)
        if im.mode not in ("RGB", "L"):
            im = im.convert("RGB")
        im.thumbnail((max_w, max_w * 2))
        buf = BytesIO()
        im.save(buf, format="JPEG", quality=85)
        data = buf.getvalue()
        out.write_bytes(data)
        return data, "image/jpeg"

    if ext == ".pdf":
        import fitz
        doc = fitz.open(path_f)
        try:
            page = max(0, min(page, len(doc) - 1))
            pix = doc[page].get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
            data = pix.tobytes("jpeg")
        finally:
            doc.close()
        out.write_bytes(data)
        return data, "image/jpeg"

    if ext == ".pptx":
        pptx_cutout_count(path_f, cache_dir)  # materialize all
        if out.exists():
            return out.read_bytes(), "image/jpeg"
        # fallback: blank board with filename
        from PIL import Image, ImageDraw
        im = Image.new("RGB", (800, 450), (242, 240, 234))
        ImageDraw.Draw(im).text((40, 200), path_f.name[:60], fill=(90, 100, 114))
        buf = BytesIO()
        im.save(buf, format="JPEG", quality=85)
        data = buf.getvalue()
        out.write_bytes(data)
        return data, "image/jpeg"

    raise RuntimeError(f"Kan ikke forhåndsvise: {path_f.suffix}")


def pdf_page_count(path_f):
    if path_f.suffix.lower() != ".pdf":
        return 1
    import fitz
    doc = fitz.open(path_f)
    try:
        return len(doc)
    finally:
        doc.close()


def illustration_section_key(template):
    """One section carries the shared anchor pack — not every chapter."""
    if not template:
        return None
    keys = {s["section_key"] for s in template.get("sections", [])}
    for sk in ("summary", "scope", "overview", "parties_scope"):
        if sk in keys:
            return sk
    return None


def apply_cover_image(state, folders, template, rel, caption=None, *, record=True):
    """WORKORDER_0.20 B — pin an indexed image as document forside (free write)."""
    sk = edchat.cover_section_key(template)
    if not sk:
        return {"ok": False, "error": "no_cover_section"}
    cap = (caption or "").strip() or Path(rel).name
    # Prefer live index caption when available
    for pth, r, cache_dir in source_files(folders):
        if r == rel or Path(r).name == Path(rel).name:
            try:
                _pf, _cd, entry = load_cache_entry(folders, r)
                if entry and entry.get("caption"):
                    cap = entry["caption"]
                rel = r
            except Exception:
                pass
            break
    state["cover_image"] = {"file": rel, "section": sk, "caption": cap}
    doc = state.setdefault("doc", {})
    sections = doc.setdefault("sections", {})
    sec = sections.setdefault(sk, {"md": "", "files": []})
    body = fc.strip_illustration_block(sec.get("md") or "").lstrip()
    # Drop a previous **Forside** block if present
    body = re.sub(r"(?is)^\*\*Forside\*\*\s*\n+(?:\{\{figure:[^}]+\}\}\s*\n*)?", "", body).lstrip()
    mark = f"{{{{figure:{rel}:0|{cap}}}}}"
    sec["md"] = f"**Forside**\n\n{mark}\n\n{body}".rstrip() + "\n"
    files = [rel] + [f for f in (sec.get("files") or []) if f != rel]
    sec["files"] = files[:4]
    if record:
        ds.add_version(state, "user", "cover", f"Forside satt: {Path(rel).name}", section=sk)
    return {"ok": True, "file": rel, "section": sk, "caption": cap}


def inject_figures_for_section(md, file_names, folders, max_figures=4, max_pages_per_file=1, excluded=None):
    """Generate {{figure:...}} markers from section's mapped visual files."""
    excluded = excluded or set()
    entries, counts = [], {}
    # Prefer raster cutouts first (PNG/JPG), then compact PDFs
    ordered = sorted(file_names or [], key=lambda r: (
        0 if Path(r).suffix.lower() in fc.PHOTO_EXT else 1,
        0 if re.search(r"s[øo]knad|rev", r, re.I) else 1,
        len(r),
    ))
    for rel in ordered:
        if figure_key(rel, 0) in excluded:
            continue
        path_f, cache_dir, entry = load_cache_entry(folders, rel)
        if not path_f:
            continue
        if not (is_visual_file(path_f, entry) or DRAWING_NAME.search(rel)
                or path_f.suffix.lower() in fc.PHOTO_EXT | {".pdf", ".pptx"}):
            continue
        try:
            total = media_page_count(path_f, cache_dir)
        except Exception:
            total = 1
        counts[rel] = min(total, max_pages_per_file)
        entries.append((rel, entry or {"caption": Path(rel).name}))
        if len(entries) >= max_figures:
            break
    md2 = fc.inject_illustration_markers(md, entries, counts)
    # Cap total figure markers so a section stays readable
    marks = list(fc.FIGURE_MARK.finditer(md2))
    if len(marks) <= max_figures:
        return md2
    body = fc.strip_illustration_block(md or "")
    kept_lines = [m.group(0) for m in marks[:max_figures]]
    return body.rstrip() + "\n\n### Illustrasjoner\n\n" + "\n\n".join(kept_lines) + "\n"


def trim_excess_figures(state, template=None, max_per_section=4):
    """Cheap cleanup — one illustrated section, strip the rest."""
    doc = state.get("doc") or {}
    ill_sk = illustration_section_key(template) if template else None
    excluded = ds.excluded_figure_keys(state)
    changed = False
    for sk, sec in (doc.get("sections") or {}).items():
        if ill_sk and sk != ill_sk:
            md = sec.get("md") or ""
            if fc.FIGURE_MARK.search(md) or "### Illustrasjoner" in md:
                sec["md"] = fc.strip_illustration_block(md)
                changed = True
            if sec.get("files"):
                sec["files"] = []
                changed = True
            continue
        files = [f for f in (sec.get("files") or []) if figure_key(f, 0) not in excluded][:max_per_section]
        if files != (sec.get("files") or []):
            sec["files"] = files
            changed = True
        md = sec.get("md") or ""
        marks = [m for m in fc.FIGURE_MARK.finditer(md)
                 if figure_key(m.group(1), m.group(2)) not in excluded]
        if len(marks) > max_per_section:
            body = fc.strip_illustration_block(md)
            kept = [m.group(0) for m in marks[:max_per_section]]
            sec["md"] = body.rstrip() + ("\n\n### Illustrasjoner\n\n" + "\n\n".join(kept) + "\n" if kept else "\n")
            changed = True
        elif len(marks) < len(list(fc.FIGURE_MARK.finditer(md))):
            body = fc.strip_illustration_block(md)
            kept = [m.group(0) for m in marks]
            sec["md"] = body.rstrip() + ("\n\n### Illustrasjoner\n\n" + "\n\n".join(kept) + "\n" if kept else "\n")
            changed = True
    return changed


def figure_key(file, page=0):
    return ds.figure_key(file, page)


def ensure_figures_in_doc(state, folders, template):
    """Inject figures into sections that need them (WORKORDER 0.48 Bug 3).

    Prefer model {{fig:}} / {{figure:}} markers. For sections with
    required_media.min_photos > 0 and no markers, insert mapped photos.
    Do not strip figures from media-required sections.
    """
    enrich_section_media(state, folders, template)
    cover = state.get("cover_image") or {}
    cover_sk = cover.get("section")
    cover_file = cover.get("file")
    excluded = ds.excluded_figure_keys(state)
    excluded_src = ds.excluded_source_files(state)
    doc = state.get("doc") or {}
    sec_defs = {
        (s.get("section_key") or s.get("key")): s
        for s in (template.get("sections") or [])
    }
    for sk, sec in (doc.get("sections") or {}).items():
        if cover_file and sk == cover_sk:
            continue
        sdef = sec_defs.get(sk) or {}
        min_photos = int((sdef.get("required_media") or {}).get("min_photos") or 0)
        md = sec.get("md") or ""
        has_fig = bool(fc.FIGURE_MARK.search(md) or fc.FIG_SHORT_MARK.search(md))
        files = [f for f in (sec.get("files") or [])
                 if figure_key(f, 0) not in excluded and f not in excluded_src]
        if min_photos > 0 or has_fig:
            if files and (not has_fig or min_photos > fc.count_figures(md)):
                sec["md"] = inject_figures_for_section(
                    md, files, folders,
                    max_figures=max(min_photos, 4),
                    excluded=excluded,
                )
            else:
                sec["md"] = md
            continue
        # Leave model-requested figures in place
        if has_fig:
            sec["md"] = md
    if cover_file and cover_sk and template:
        apply_cover_image(state, folders, template, cover_file, cover.get("caption"),
                          record=False)
    return state


def materialize_export_figures(folder, content):
    """Render {{figure}} markers to Rapporter/media/*.jpg and rewrite as ![](media/...)."""
    folder = Path(folder)
    media_dir = reports_dir(folder) / "media"
    media_dir.mkdir(parents=True, exist_ok=True)

    def url_fn(rel, page):
        path_f, cache_dir = find_source([str(folder)], rel)
        if not path_f:
            # multi-folder projects — try all projects? only this folder for now
            for p, r, cd in source_files([str(folder)]):
                if r == rel:
                    path_f, cache_dir = p, cd
                    break
        if not path_f:
            return f"media/missing-{page}.jpg"
        try:
            data, _ = render_media_preview(path_f, cache_dir, page=page)
        except Exception:
            return f"media/missing-{page}.jpg"
        safe = re.sub(r"[^\w.\-]+", "_", Path(rel).stem)[:60]
        name = f"{safe}-p{page}.jpg"
        (media_dir / name).write_bytes(data)
        return f"media/{name}"

    return fc.expand_figures_to_markdown_images(content, url_fn)


def pick_anchor_drawings(visuals, folders, max_n=5, approved=None):
    """Choose a small set of drawings that ground the text — prefer PNG/JPG
    application packs (Til Søknad / REV) over huge multi-page PDFs.
    Higher REV numbers and explicitly approved drawings win."""
    kinds = [
        ("plantegning", re.compile(r"plantegning", re.I)),
        ("fasade", re.compile(r"fasade", re.I)),
        ("situasjon", re.compile(r"situasjon|situvasjon", re.I)),
        ("snitt", re.compile(r"snitt", re.I)),
        ("perspektiv", re.compile(r"perspektiv", re.I)),
    ]
    approved = approved or []
    approved_norms = [a.replace("\\", "/").lower() for a in approved]

    def is_approved(name):
        n = name.replace("\\", "/").lower()
        return any(a in n or n.endswith(a.split("/")[-1]) for a in approved_norms)

    scored = []
    for e in visuals:
        name = e["file"]
        path_f, _ = find_source(folders, name)
        if not path_f:
            continue
        ext = path_f.suffix.lower()
        score = 0
        if ext in fc.PHOTO_EXT:
            score += 80  # already a cutout — show this
        rev_m = re.search(r"rev\s*(\d+)", name, re.I)
        if rev_m:
            score += 45 + int(rev_m.group(1)) * 12  # REV2/REV3 beat unrevised packs
        if re.search(r"s[øo]knad", name, re.I):
            score += 20
        if is_approved(name):
            score += 120  # user-/kommune-approved governing drawing
        if DRAWING_NAME.search(name) or fc.is_visual_source(name, e):
            score += 15
        # Prefer cutouts of the same approved set over older TOPP/unnamed packs
        if re.search(r"\btopp\b", name, re.I) and not rev_m:
            score -= 35
        try:
            sz = path_f.stat().st_size
            if sz < 2_000_000:
                score += 25
            elif sz > 8_000_000:
                score -= 30  # avoid huge PDFs when a PNG exists
        except OSError:
            pass
        kind = None
        for kid, rx in kinds:
            if rx.search(name):
                kind = kid
                break
        scored.append((score, kind, e))
    scored.sort(key=lambda x: -x[0])
    by_kind = {}
    for score, kind, e in scored:
        if not kind:
            continue
        by_kind.setdefault(kind, []).append((score, e))
    picked = []
    for kid, _rx in kinds:
        cands = by_kind.get(kid) or []
        if not cands:
            continue
        photos = [(s, e) for s, e in cands
                  if Path(e["file"]).suffix.lower() in fc.PHOTO_EXT]
        # Prefer raster cutout of the winning revision; fall back to PDF/SVG
        _score, chosen = max(photos or cands, key=lambda x: x[0])
        picked.append(chosen)
        if len(picked) >= max_n:
            break
    return picked


def enrich_section_media(state, folders, template):
    """Attach visual files (photos + drawing/presentation PDFs) to sections by role.
    Does not require regenerate — fills empty or sparse section.files lists."""
    if not state.get("doc") or not template:
        return state
    index = []
    for p, rel, cache_dir in source_files(folders):
        import hashlib
        st = p.stat()
        fast = cache_dir / f"path-{hashlib.sha256(f'{rel}:{st.st_size}:{st.st_mtime_ns}'.encode()).hexdigest()[:16]}.json"
        if fast.exists():
            try:
                e = fc.read_json_file(fast)
                e["file"] = rel
                index.append(e)
                continue
            except Exception:
                pass
        if p.suffix.lower() not in fc.PHOTO_EXT and not DRAWING_NAME.search(rel):
            continue
        sha = hashlib.sha256(p.read_bytes()).hexdigest()
        cache = cache_dir / f"{sha}.json"
        if cache.exists():
            try:
                e = fc.read_json_file(cache)
                e["file"] = rel
                index.append(e)
            except Exception:
                pass
        elif p.suffix.lower() in fc.PHOTO_EXT or DRAWING_NAME.search(rel):
            # uncached drawing still usable as visual anchor
            index.append({"file": rel, "doc_role_hints": ["drawing"], "caption": Path(rel).name, "facts": []})
    visuals = []
    for e in index:
        path_f, _cd = find_source(folders, e["file"])
        if path_f and (is_visual_file(path_f, e) or DRAWING_NAME.search(e["file"])
                       or path_f.suffix.lower() in fc.PHOTO_EXT):
            visuals.append(e)
    approved = list(state.get("approved_drawings") or [])
    anchors = pick_anchor_drawings(visuals, folders, max_n=4, approved=approved)
    ill_sk = illustration_section_key(template)
    sections = state["doc"].setdefault("sections", {})
    excluded = ds.excluded_figure_keys(state)
    anchor_files = [e["file"] for e in anchors if figure_key(e["file"], 0) not in excluded][:4]

    for s in template.get("sections", []):
        sk = s["section_key"]
        if sk not in sections:
            sections[sk] = {"md": "", "files": []}
        if sk != ill_sk:
            sections[sk]["files"] = []
            continue
        roles = set(s.get("required_media", {}).get("preferred_roles") or [])
        sec = sections[sk]
        files = list(anchor_files)
        for e in visuals:
            name = e["file"]
            if name in files or figure_key(name, 0) in excluded:
                continue
            eroles = set(e.get("doc_role_hints") or [])
            if roles and (eroles & roles) and len(files) < 4:
                files.append(name)
        sec["files"] = files[:4]
    return state


def persist_doc(folder, state, template_file):
    """Reassemble draft from v2 sections and write all export paths."""
    folder = str(folder if not isinstance(folder, list) else folder[0])
    template = load_template(template_file)
    ensure_figures_in_doc(state, [folder], template)
    content = ds.assemble_draft(state, template, state.get("artifact"))
    content = materialize_export_figures(folder, content)
    export_path, export_name = sync_draft_files(folder, state, template, template_file, content)
    for d in state.get("documents", []):
        if d.get("template") == template_file:
            d["gaps"] = len(state.get("gaps", []))
            d["blocking"] = sum(1 for g in state.get("gaps", []) if g.get("severity") == "blocking")
            d["gap_list"] = state.get("gaps", [])
            d["export_path"] = str(export_path)
            d["export_name"] = export_name
    save_state(folder, state)
    return content


def _pid(body):
    return body.get("id") or body.get("project")


def add_bom_component(state, *, part_no, file, caption=None, confidence=0.0,
                      fact_id=None, status="ok", verified_by_user=False):
    """WO 0.22 C3 — one BOM row in document state (never templates)."""
    rows = state.setdefault("bom_components", [])
    # Upsert by file
    for r in rows:
        if r.get("file") == file:
            r.update({
                "part_no": part_no, "caption": caption, "confidence": confidence,
                "fact_id": fact_id, "status": status,
                "verified_by_user": verified_by_user or r.get("verified_by_user"),
            })
            return r
    row = {
        "part_no": part_no, "file": file, "caption": caption,
        "confidence": confidence, "fact_id": fact_id, "status": status,
        "verified_by_user": verified_by_user,
    }
    rows.append(row)
    ds.add_version(state, "user", "bom", f"BOM-komponent: {part_no or file}")
    return row


def run_component_scan(folders, index, state, *, lang="no"):
    """Scan unscanned photos for part IDs; update index cache flags + bom_components."""
    import agent_truth as atruth
    photos = atruth.photo_entries(index)
    pending = [e for e in photos if atruth.needs_component_scan(e)]
    ok = uncertain = no_id = 0
    # Re-use existing extraction facts when present; mark scanned
    for e in pending or photos:
        parts = atruth.part_facts_from_entry(e)
        e["component_scanned"] = True
        # Persist flag on cache if possible
        try:
            rel = e.get("file")
            path_f, cache_dir, entry = load_cache_entry(folders, rel)
            if entry is not None and cache_dir:
                entry["component_scanned"] = True
                sha = entry.get("sha")
                if sha:
                    (Path(cache_dir) / f"{sha}.json").write_text(
                        json.dumps(entry, ensure_ascii=False, indent=2), encoding="utf-8")
        except Exception:
            pass
        if not parts:
            no_id += 1
            continue
        best = max(parts, key=lambda f: float(f.get("confidence") or 0))
        conf = float(best.get("confidence") or 0)
        val = str(best.get("value") or "").strip()
        if not val:
            no_id += 1
            continue
        status = "ok" if conf >= atruth.CONF_OK else "uncertain"
        add_bom_component(
            state, part_no=val, file=e.get("file"), caption=e.get("caption"),
            confidence=conf, fact_id=best.get("id"), status=status)
        if status == "ok":
            ok += 1
        else:
            uncertain += 1
    return {"ok": ok, "uncertain": uncertain, "no_id": no_id,
            "scanned": len(pending) or len(photos)}


def assert_source_immutable_write(dest: Path, *, engine_root: Path):
    """WO 0.22 C — refuse template writes; callers must use additive_dest for sources."""
    import agent_truth as atruth
    atruth.assert_not_template_write(dest, engine_root)


def create_project_with_skeleton(name, template_key=None, *, base_dir=None,
                                 conversation=None):
    """Shared create path for API button and hub chat execute (WO 0.21 A1)."""
    settings = hub.load_settings()
    base = (base_dir or settings.get("projects_base_dir") or "").strip().strip('"')
    if not base:
        return {"error": "Velg hvor nye prosjektmapper skal ligge (én gang).",
                "code": "need_base_dir"}
    base_path = Path(base)
    if not base_path.is_dir():
        return {"error": f"Basismappe finnes ikke: {base}", "code": "bad_base_dir"}
    if base_dir:
        settings["projects_base_dir"] = str(base_path)
        hub.save_settings(settings)
    try:
        folder = hub.create_skeleton_folder(base_path, name)
    except Exception as e:
        return {"error": str(e)}
    caps = hub.load_capabilities()
    tkey = (template_key or "").strip() or None
    caps_entry = next((t for t in caps.get("templates") or [] if t.get("key") == tkey), None)
    hub.write_checklist(folder, caps_entry)
    projects = load_projects()
    proj = {"id": uuid.uuid4().hex[:8], "name": name, "folders": [str(folder)]}
    if tkey and caps_entry:
        proj["preferred_template"] = caps_entry.get("file")
    projects.append(proj)
    save_projects(projects)
    state = load_state(folder, project_id=proj["id"])
    for turn in conversation or []:
        if isinstance(turn, dict) and turn.get("text"):
            edchat.append_turn(state, turn.get("role") or "user", turn["text"],
                               html=turn.get("html"), project_id=proj["id"])
    if not state.get("conversation"):
        edchat.append_turn(state, "bot",
            f"Prosjektmappen «{name}» er klar. Fortsett her — samme samtale.",
            project_id=proj["id"])
    save_state(folder, state)
    checklist_path = folder / "SJEKKLISTE.txt"
    checklist_n = 0
    if caps_entry and caps_entry.get("checklist"):
        checklist_n = len(caps_entry["checklist"])
    elif checklist_path.exists():
        checklist_n = sum(1 for ln in checklist_path.read_text(encoding="utf-8").splitlines()
                          if ln.strip().startswith("□"))
    return {
        **proj,
        "checklist": str(checklist_path),
        "checklist_count": checklist_n,
        "checklist_items": list((caps_entry or {}).get("checklist") or [])[:12],
        "skeleton": list(hub.SKELETON_DIRS),
        "template_key": tkey,
        "caps_entry": caps_entry,
        "conversation": state.get("conversation") or [],
    }


def create_project_from_staged(token, *, name=None, template_key=None, base_dir=None,
                               conversation=None):
    """WORKORDER_0.25 — create project and copy hub-staged indexed file in."""
    import hub_session as hses
    import shutil
    session = hses.load_session()
    staged = next((s for s in (session.get("staged") or [])
                   if s.get("token") == token), None)
    if not staged:
        return {"error": "Fant ikke den indekserte filen — dra den inn på nytt."}
    pname = name or staged.get("name") or "Prosjekt"
    pname = Path(pname).stem
    # Guess template from caption/facts
    tkey = template_key
    if not tkey:
        blob = " ".join([
            staged.get("caption") or "",
            " ".join(staged.get("fact_keys") or []),
            staged.get("name") or "",
        ]).lower()
        if any(w in blob for w in ("forsikring", "insurance", "kontrakt", "contract",
                                   "vilkår", "policy")):
            tkey = "contract_review"
    created = create_project_with_skeleton(
        pname, tkey, base_dir=base_dir, conversation=conversation)
    if created.get("error"):
        return created
    folder = Path(created["folders"][0])
    src = Path(staged["path"])
    if src.exists():
        dest_dir = folder / "Notater"
        dest_dir.mkdir(parents=True, exist_ok=True)
        dest = dest_dir / staged["name"]
        shutil.copy2(src, dest)
        # Copy cache entry if present
        try:
            cache_src = ATTACH_STAGING / "hub_cache"
            if cache_src.is_dir():
                cache_dst = fpaths.cache_dir(folder)
                cache_dst.mkdir(exist_ok=True)
                for f in cache_src.glob("*.json"):
                    shutil.copy2(f, cache_dst / f.name)
        except Exception:
            pass
    state = load_state(folder)
    edchat.append_turn(
        state, "system",
        f"[project_created] {pname} from staged {staged.get('name')} | "
        f"Indeksert som: {staged.get('caption')}",
        meta={"kind": "project_created"})
    save_state(folder, state)
    hses.clear_pending(session)
    # Drop this staged file
    session["staged"] = [s for s in session.get("staged") or [] if s.get("token") != token]
    hses.append_event(session, "system", f"[project_created] {pname}",
                      meta={"kind": "project_created", "id": created["id"]})
    hses.save_session(session)
    created["staged_file"] = staged.get("name")
    created["from_hub_stage"] = True
    return created


def create_demo_project(kind="contract", *, base_dir=None, name=None,
                        template_key=None, conversation=None):
    """WORKORDER_0.23 C — marked synthetic demo; not exportable as paid."""
    import demo_project as demo
    kind = (kind or "contract").lower().strip()
    if kind in ("tech", "lifting"):
        kind = "technical"
    if kind not in ("contract", "technical"):
        kind = "contract"
    default_name = "DEMO_Løfteverktøy" if kind == "technical" else "DEMO_Kontraktsak"
    tkey = template_key or (
        "technical_doc_package" if kind == "technical" else "contract_review")
    created = create_project_with_skeleton(
        name or default_name, tkey, base_dir=base_dir, conversation=conversation)
    if created.get("error"):
        return created
    folder = Path(created["folders"][0])
    files = demo.create_demo_files(folder, kind)
    state = load_state(folder)
    state["demo"] = True
    state["demo_kind"] = kind
    state["export_paid_blocked"] = True
    caps = hub.load_capabilities()
    entry = next((t for t in caps.get("templates") or [] if t.get("key") == tkey), None)
    if entry and entry.get("file"):
        state["active_template"] = entry["file"]
        state["template"] = entry["file"]
    save_state(folder, state)
    created["demo"] = True
    created["demo_files"] = files
    created["kind"] = kind
    return created


class Handler(BaseHTTPRequestHandler):
    def log_message(self, fmt, *args):
        pass  # keep console clean; engine prints its own progress

    def _send(self, code, body, ctype="application/json; charset=utf-8"):
        data = body if isinstance(body, bytes) else json.dumps(body, ensure_ascii=False).encode("utf-8")
        self.send_response(code)
        self.send_header("Content-Type", ctype)
        self.send_header("Content-Length", str(len(data)))
        self.end_headers()
        self.wfile.write(data)

    def _json_body(self):
        length = int(self.headers.get("Content-Length", 0))
        return json.loads(self.rfile.read(length) or b"{}")

    def do_GET(self):
        path, _, query = self.path.partition("?")
        params = {unquote(kv.split("=", 1)[0]): unquote(kv.split("=", 1)[1])
                  for kv in query.split("&") if "=" in kv}

        if path in ("/", "/index.html"):
            html = (APP_DIR / "app.html").read_bytes()
            return self._send(200, html, "text/html; charset=utf-8")

        if path == "/VERSION":
            ver = ROOT / "VERSION"
            if ver.exists():
                return self._send(200, ver.read_text(encoding="utf-8").strip() + "\n",
                                  "text/plain; charset=utf-8")
            return self._send(404, "missing\n", "text/plain; charset=utf-8")

        # Brand favicon (Foldok mark […]) — files live in public/
        favicon_map = {
            "/favicon.ico": ("favicon.ico", "image/x-icon"),
            "/favicon.svg": ("favicon.svg", "image/svg+xml"),
            "/favicon-16.png": ("favicon-16.png", "image/png"),
            "/favicon-32.png": ("favicon-32.png", "image/png"),
            "/apple-touch-icon.png": ("apple-touch-icon.png", "image/png"),
        }
        if path in favicon_map:
            name, ctype = favicon_map[path]
            fav = ROOT / "public" / name
            if fav.exists():
                return self._send(200, fav.read_bytes(), ctype)
            return self._send(404, b"missing", "text/plain; charset=utf-8")

        if path == "/site-meta.json":
            meta_path = ROOT / "site-meta.json"
            if meta_path.exists():
                return self._send(200, meta_path.read_bytes(), "application/json; charset=utf-8")
            return self._send(404, {"error": "site-meta missing"})

        if path == "/design-dummy.html":
            dummy = ROOT / "web" / "design-dummy.html"
            if dummy.exists():
                return self._send(200, dummy.read_bytes(), "text/html; charset=utf-8")

        if path == "/diagram.html":
            page = ROOT / "web" / "diagram.html"
            if page.exists():
                return self._send(200, page.read_bytes(), "text/html; charset=utf-8")

        if path == "/boxes-demo.html":
            page = ROOT / "web" / "boxes-demo.html"
            if page.exists():
                return self._send(200, page.read_bytes(), "text/html; charset=utf-8")

        if path == "/foldok-box-editor.js":
            js = ROOT / "web" / "foldok-box-editor.js"
            if not js.exists():
                js = ROOT / "foldok_boxes" / "editor" / "foldok-box-editor.js"
            if js.exists():
                return self._send(200, js.read_bytes(), "application/javascript; charset=utf-8")

        if path == "/api/layout/session":
            sid = (params.get("id") or "").strip()
            try:
                import box_sessions as bses

                if sid and bses.get_session(sid):
                    return self._send(200, bses.session_payload(sid))
                return self._send(200, bses.create_session())
            except Exception as e:
                return self._send(500, {"error": str(e)})

        if path == "/api/capture/status":
            p = get_project(params.get("id", ""))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = primary_folder(p)
            if folder is None or not folder.is_dir():
                return self._send(400, {"error": "Capture bridge requires a project folder"})
            return self._send(200, capture_status_for_folder(folder))

        if path == "/electrical-dummy.html":
            self.send_response(302)
            self.send_header("Location", "/diagram.html")
            self.end_headers()
            return

        if path == "/api/diagram":
            fixture = (params.get("fixture") or "water_heater").strip()
            profile = (params.get("profile") or "").strip() or None
            width = float(params.get("width") or params.get("target_width_pt") or 420)
            try:
                import diagram_sessions as dses

                return self._send(200, dses.render_fixture_svg(fixture, profile=profile, target_width_pt=width))
            except Exception as e:
                return self._send(500, {"error": str(e), "fixture": fixture})

        if path == "/api/diagram/project":
            p = get_project(params.get("id", ""))
            if not p or not (p.get("folders") or []):
                return self._send(404, {"error": "unknown project"})
            try:
                from diagram_store import list_diagrams

                return self._send(200, {"diagrams": list_diagrams(p["folders"][0])})
            except Exception as e:
                return self._send(500, {"error": str(e)})

        if path == "/api/diagram/open":
            p = get_project(params.get("id", ""))
            if not p or not (p.get("folders") or []):
                return self._send(404, {"error": "unknown project"})
            graph_id = (params.get("graph_id") or params.get("graph") or "").strip()
            if not graph_id:
                return self._send(400, {"error": "graph_id required"})
            try:
                import diagram_sessions as dses

                payload = dses.open_project_diagram(
                    p["folders"][0],
                    graph_id,
                    profile=params.get("profile"),
                    target_width_pt=float(params.get("width") or 420),
                )
                return self._send(200, payload)
            except FileNotFoundError as e:
                return self._send(404, {"error": str(e)})
            except Exception as e:
                return self._send(500, {"error": str(e)})

        if path == "/api/diagram/session":
            sid = (params.get("id") or "").strip()
            try:
                import diagram_sessions as dses

                ed = dses.get_editor(sid) if sid else None
                if not ed:
                    return self._send(404, {"error": "unknown session"})
                return self._send(200, dses.session_payload(sid))
            except Exception as e:
                return self._send(500, {"error": str(e)})

        if path == "/sample_multipoint.html":
            preview = ROOT / "fixtures" / "sample_multipoint" / "sample_multipoint.html"
            if preview.exists():
                return self._send(200, preview.read_bytes(), "text/html; charset=utf-8")

        if path == "/api/media":
            p = get_project(params.get("id", ""))
            if not p:
                return self._send(404, {"error": "unknown project"})
            rel = params.get("file", "")
            page = int(params.get("page", "0") or 0)
            path_f, cache_dir = find_source(p["folders"], rel)
            if not path_f:
                return self._send(404, {"error": "fil ikke funnet"})
            try:
                data, ctype = render_media_preview(path_f, cache_dir, page=page)
                return self._send(200, data, ctype)
            except Exception as e:
                return self._send(500, {"error": str(e)})

        if path == "/api/media/meta":
            p = get_project(params.get("id", ""))
            if not p:
                return self._send(404, {"error": "unknown project"})
            rel = params.get("file", "")
            path_f, cache_dir, entry = load_cache_entry(p["folders"], rel)
            if not path_f:
                return self._send(404, {"error": "fil ikke funnet"})
            pages = 1
            try:
                pages = media_page_count(path_f, cache_dir)
            except Exception:
                pages = 1
            budget = fc.figure_page_budget(rel, entry, pages) or (1 if path_f.suffix.lower() in fc.PHOTO_EXT else 0)
            show = max(budget, 1) if is_visual_file(path_f, entry) or path_f.suffix.lower() in fc.PHOTO_EXT else 0
            show = min(pages, max(show, 1)) if show else min(pages, 1)
            return self._send(200, {
                "file": rel, "pages": pages, "show_pages": show,
                "kind": "photo" if path_f.suffix.lower() in fc.PHOTO_EXT else "doc",
                "caption": (entry or {}).get("caption") or rel,
                "roles": (entry or {}).get("doc_role_hints") or [],
                "visual": is_visual_file(path_f, entry) or path_f.suffix.lower() in {".pdf", ".pptx"} | fc.PHOTO_EXT,
            })

        if path == "/api/bootstrap":
            caps = hub.load_capabilities()
            settings = hub.load_settings()
            version = (ROOT / "VERSION").read_text(encoding="utf-8").strip()
            web_meta = {"document": "FOLDOK-WEB-001", "revision": "C", "date": "auto",
                        "languages": ["no", "en"], "default_language": "en"}
            meta_path = ROOT / "site-meta.json"
            if meta_path.exists():
                try:
                    web_meta.update(json.loads(meta_path.read_text(encoding="utf-8")))
                except Exception:
                    pass
            if (web_meta.get("date") or "auto") == "auto":
                from datetime import date as _date
                web_meta["date"] = _date.today().isoformat()
            web_meta["version"] = version
            return self._send(200, {
                "projects": load_projects(),
                "templates": templates_list(),
                "capabilities": caps,
                "settings": {"projects_base_dir": settings.get("projects_base_dir") or ""},
                "key_set": KEY_SET,
                "version": version,
                "web_meta": web_meta,
                "learning_path": str(learning.path()),
                "account": acct.account_snapshot(),
            })

        if path == "/api/account":
            return self._send(200, acct.account_snapshot())

        if path == "/api/account/usage":
            tok = acct.device_token()
            if not tok:
                return self._send(401, {"error": "Ikke innlogget"})
            try:
                return self._send(200, acct.get_ledger().usage(tok))
            except acct.MeterDenied as e:
                return self._send(402, {"error": str(e), "code": e.code})

        if path == "/api/account/receipt":
            tok = acct.device_token()
            rid = params.get("id", "")
            if not tok or not rid:
                return self._send(400, {"error": "id påkrevd"})
            try:
                raw, rcpt = acct.get_ledger().get_receipt_pdf(tok, rid)
                self.send_response(200)
                self.send_header("Content-Type", "application/octet-stream")
                self.send_header("Content-Disposition",
                                 f'attachment; filename="{rcpt.get("doc_name", "export")}.bin"')
                self.send_header("Content-Length", str(len(raw)))
                self.end_headers()
                self.wfile.write(raw)
                return
            except acct.MeterDenied as e:
                return self._send(404, {"error": str(e), "code": e.code})

        if path == "/api/learning":
            return self._send(200, {"learning": learning.load(), "path": str(learning.path())})

        if path == "/api/browse":
            result = browse_dirs(unquote(params.get("path", "") or ""))
            if result.get("error"):
                return self._send(400, result)
            return self._send(200, result)

        if path == "/api/open-source":
            p = get_project(params.get("id", ""))
            if not p:
                return self._send(404, {"error": "unknown project"})
            rel = unquote(params.get("file", "") or "")
            if not rel:
                return self._send(400, {"error": "file er påkrevd"})
            path_f, _ = find_source(p["folders"], rel)
            if not path_f:
                return self._send(404, {"error": f"Filen finnes ikke: {rel}"})
            try:
                open_path_in_os(path_f, reveal=params.get("reveal") == "1")
                return self._send(200, {"ok": True, "path": str(path_f)})
            except Exception as e:
                return self._send(500, {"error": str(e)})

        if path == "/api/project":
            p = get_project(params.get("id", ""))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folders = list(p.get("folders") or [])
            # WORKORDER 0.61 — folder-less project: memory state, no IndexError
            if not folders or not any(Path(f).is_dir() for f in folders):
                state, _ = project_state_load(p, params.get("template"))
                view_tpl = (
                    params.get("template")
                    or state.get("active_template")
                    or state.get("template")
                    or p.get("preferred_template")
                )
                templates = templates_list(p)
                documents = list(state.get("documents") or [])
                # enrich status lightly
                by_file = {t["file"]: t for t in templates}
                documents = [
                    _enrich_doc(None, d, by_file, state) if d.get("template") else d
                    for d in documents
                ]
                gap_sum = ds.gaps_summary(state.get("gaps", []))
                return self._send(200, {
                    **p,
                    "folders": [],
                    "folderless": True,
                    "need_folder_banner": bool(state.get("need_folder_banner", True)),
                    "files": [],
                    "estimate_eur": 0,
                    "artifact": state.get("artifact"),
                    "confirmed": state.get("confirmed"),
                    "gaps": state.get("gaps", []),
                    "gap_summary": gap_sum,
                    "fillable_gaps": 0,
                    "index_file_count": 0,
                    "suggestions": [],
                    "supersede_map": {},
                    "dismissed_suggestions": [],
                    "dismissed": state.get("dismissed", []),
                    "excluded_figures": state.get("excluded_figures", []),
                    "excluded_sources": state.get("excluded_sources", []),
                    "source_selection": {},
                    "blocking_dismissed": ds.blocking_dismissed(state),
                    "reference_facts": ds.reference_facts(state),
                    "empty_folder": True,
                    "conversation": isolated_conversation(state, p.get("id")),
                    "template": state.get("template") or view_tpl,
                    "active_template": view_tpl or state.get("active_template"),
                    "documents": documents,
                    "doc": state.get("doc"),
                    "user_facts": state.get("user_facts", []),
                    "versions": state.get("versions", [])[:20],
                    "draft_exists": bool(documents),
                    "complete": False,
                    "product_repo_warn": False,
                    "capture": None,
                })
            missing = [f for f in folders if not Path(f).is_dir()]
            if missing:
                return self._send(200, {**p, "missing": True, "missing_folders": missing})
            rows = file_rows(folders)
            state = load_state(folders[0], project_id=p.get("id"))
            sel = prescan.normalize_source_selection(state.get("source_selection"))
            for r in rows:
                r["enabled"] = prescan.source_is_enabled(r.get("name") or "", sel)
            product_repo_warn = looks_like_product_repo(folders[0])
            view_tpl = params.get("template") or state.get("active_template") or state.get("template")
            vt = None
            if view_tpl:
                ensure_doc_for_template(folders[0], state, view_tpl)
                # Align active pointer with the document being viewed
                if state.get("active_template") != view_tpl:
                    state["active_template"] = view_tpl
                    state["template"] = view_tpl
                vt = load_template(view_tpl, folders[0])
                if vt and state.get("doc", {}).get("sections"):
                    # Refuse to compute gaps if doc still belongs to another template
                    doc_tf = (state.get("doc") or {}).get("template_file")
                    if doc_tf and doc_tf != view_tpl:
                        ensure_doc_for_template(folders[0], state, view_tpl)
                    # Figures are injected on generate / refresh-figures only — not every page view
                    if ensure_template_sections(state, vt):
                        save_state(folders[0], state)
                    bom_sec = (state.get("doc") or {}).get("sections", {}).get("bom")
                    if bom_sec is not None and not (bom_sec.get("md") or "").strip():
                        if refresh_bom_section(state, folders, view_tpl):
                            save_state(folders[0], state)
                    dc_sec = (state.get("doc") or {}).get("sections", {}).get("doc_control")
                    dc_md = (dc_sec or {}).get("md") or ""
                    if dc_sec is not None and (dc_md.count("[MANGLER: doc_no]") >= 2
                                               or not dc_sec.get("table")):
                        if refresh_doc_control_section(state, folders, view_tpl):
                            save_state(folders[0], state)
                    so_sec = (state.get("doc") or {}).get("sections", {}).get("spec_overview")
                    so_md = (so_sec or {}).get("md") or ""
                    if so_sec is not None and (so_md.count("[MANGLER: issuer]") >= 1
                                               or "[MANGLER: invoked_document]" in so_md
                                               or not so_sec.get("table")):
                        if refresh_spec_overview_section(state, folders, view_tpl):
                            save_state(folders[0], state)
                    if trim_excess_figures(state, vt):
                        save_state(folders[0], state)
                    # Sketch / empty docs: never block page load on live indexing
                    is_sketch = (
                        (vt or {}).get("document_species") == "sketch"
                        or view_tpl == "sketch_document.json"
                        or bool((state.get("doc") or {}).get("sketch", {}).get("mode"))
                    )
                    try:
                        index = [] if is_sketch else load_active_index(state, folders, "no", cache_only=True)
                        # Zero-token: model often wrote [MANGLER: criterion] for tilbygg m²
                        if index and fc.repair_miskeyed_area_mangler(state, index, state.get("artifact")):
                            save_state(folders[0], state)
                        # WORKORDER 0.58 §0 — document-scoped gaps only
                        documents_pre = list_documents(folders[0], state, templates_list(p))
                        state["gaps"] = ds.gaps_for_document(
                            state, vt, index, state.get("artifact"), fc,
                            fast=True, documents=documents_pre)
                    except Exception as e:
                        print(f"[project] gaps/index skipped: {e}", flush=True)
                        state["gaps"] = state.get("gaps") or []
            if state.get("doc") and not state.get("gaps"):
                # Only MD-scan the active document's own sections
                doc_tf = (state.get("doc") or {}).get("template_file")
                if not view_tpl or doc_tf == view_tpl:
                    state["gaps"] = []
                    for sk, sec in (state["doc"].get("sections") or {}).items():
                        state["gaps"].extend(ds.gaps_from_md(sec.get("md", ""), sk))
                    if view_tpl:
                        vt2 = load_template(view_tpl)
                        if vt2:
                            state["gaps"] = ds.filter_gaps_to_template(state["gaps"], vt2)
            templates = templates_list(p)
            documents = list_documents(folders[0], state, templates)
            active = state.get("active_template") or state.get("template")
            gap_sum = ds.gaps_summary(state.get("gaps", []))
            fillable = 0
            if view_tpl and state.get("gaps"):
                try:
                    idx_for_fill = load_active_index(state, folders, "no", cache_only=True)
                    fillable = fc.count_fillable_gaps(
                        {**state, "documents": documents}, idx_for_fill, state.get("artifact"))
                except Exception:
                    fillable = 0
            suggestions = []
            try:
                suggestions = active_suggestions(state, folders, view_tpl) if state.get("artifact") else []
            except Exception:
                suggestions = []
            full_index = []
            try:
                full_index = load_index(folders, "no", state.get("user_facts"), cache_only=True) if state.get("artifact") else []
            except Exception:
                full_index = []
            dismissed_sup = {d.get("name") for d in state.get("dismissed_suggestions", [])}
            supersede_map = (fc.superseded_files_map(full_index, ds.excluded_source_files(state), folders)
                             if full_index else {})
            supersede_map = {k: v for k, v in supersede_map.items()
                             if v.get("name") not in dismissed_sup}
            return self._send(200, {**p, "files": rows, "estimate_eur": estimate_cost(rows),
                                    "artifact": state.get("artifact"), "confirmed": state.get("confirmed"),
                                    "gaps": state.get("gaps", []), "gap_summary": gap_sum,
                                    "fillable_gaps": fillable,
                                    "index_file_count": len([f for f in rows if f.get("kind") != "skipped"]),
                                    "suggestions": suggestions,
                                    "supersede_map": supersede_map,
                                    "dismissed_suggestions": state.get("dismissed_suggestions", []),
                                    "dismissed": state.get("dismissed", []),
                                    "excluded_figures": state.get("excluded_figures", []),
                                    "excluded_sources": state.get("excluded_sources", []),
                                    "source_selection": sel,
                                    "cell_overrides": state.get("cell_overrides", []),
                                    "source_citation_warnings": state.get("source_citation_warnings", []),
                                    "illustration_section": illustration_section_key(vt) if view_tpl else None,
                                    "blocking_dismissed": ds.blocking_dismissed(state),
                                    "reference_facts": ds.reference_facts(state),
                                    "empty_folder": len([f for f in rows if f.get("kind") != "skipped"]) == 0,
                                    "conversation": isolated_conversation(state, p.get("id")),
                                    "chat_pending": state.get("chat_pending"),
                                    "assist_hint_shown": bool(state.get("assist_hint_shown")),
                                    "template": state.get("template"),
                                    "active_template": active, "documents": documents,
                                    "doc": state.get("doc"), "user_facts": state.get("user_facts", []),
                                    "versions": state.get("versions", [])[:20],
                                    "violations": state.get("violations", []),
                                    "draft_exists": bool(documents) or (Path(folders[0]) / "draft.md").exists(),
                                    "complete": gap_sum["blocking"] == 0 and gap_sum.get("warning", 0) == 0
                                                and bool(state.get("doc")),
                                    "product_repo_warn": product_repo_warn,
                                    "capture": capture_status_for_folder(folders[0])})

        if path == "/api/job":
            jid = params.get("id", "")
            job = JOBS.get(jid)
            if not job:
                return self._send(404, {"error": "unknown job"})
            with LOCK:
                j = JOBS.get(jid)
                if j and j.get("kind") == "index":
                    j["last_heartbeat"] = time.time()
            return self._send(200, job)
        if path == "/api/draft":
            p = get_project(params.get("id", ""))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = Path(p["folders"][0])
            tpl = params.get("template", "")
            state = load_state(folder, tpl or None)
            if not tpl:
                tpl = state.get("active_template") or state.get("template")
            template = load_template(tpl) if tpl else None
            if template and state.get("doc"):
                ensure_figures_in_doc(state, p["folders"], template)
                save_state(folder, state)
            content = get_draft_md(folder, state, template, tpl) if template else read_draft_content(folder, tpl)
            if not content:
                return self._send(404, {"error": "no draft"})
            return self._send(200, content.encode("utf-8"), "text/plain; charset=utf-8")

        if path == "/api/doc/versions":
            p = get_project(params.get("id", "") or params.get("project", ""))
            if not p:
                return self._send(404, {"error": "unknown project"})
            state = load_state(p["folders"][0])
            return self._send(200, {"versions": state.get("versions", [])})

        return self._send(404, {"error": "not found"})

    def do_POST(self):
        body = self._json_body()
        path = self.path.split("?")[0].rstrip("/") or "/"

        # ── Box layout (foldok_boxes · WO 0.73) ───────────────────────
        if path == "/api/layout/session":
            try:
                import box_sessions as bses

                return self._send(200, bses.create_session())
            except Exception as e:
                return self._send(500, {"error": str(e)})

        if path == "/api/layout/session/intent":
            try:
                import box_sessions as bses

                sid = (body.get("id") or body.get("session_id") or "").strip()
                if not sid:
                    return self._send(400, {"error": "session id required"})
                intent = body.get("intent") if isinstance(body.get("intent"), dict) else body
                return self._send(200, bses.apply(sid, intent))
            except KeyError:
                return self._send(404, {"error": "unknown layout session"})
            except Exception as e:
                return self._send(500, {"error": str(e)})

        # ── Diagram canvas (foldok_diagram · pins not geometry) ───────
        if path == "/api/diagram/save":
            try:
                import diagram_sessions as dses

                sid = (body.get("id") or body.get("session_id") or "").strip()
                if not sid:
                    return self._send(400, {"error": "session id required"})
                paths = dses.persist_session(sid)
                return self._send(200, {"ok": True, "paths": paths})
            except KeyError:
                return self._send(404, {"error": "unknown session"})
            except ValueError as e:
                return self._send(400, {"error": str(e)})
            except Exception as e:
                return self._send(500, {"error": str(e)})

        if path == "/api/diagram/session":
            try:
                import diagram_sessions as dses

                fixture = (body.get("fixture") or "water_heater").strip()
                profile = (body.get("profile") or "").strip() or None
                width = float(body.get("target_width_pt") or body.get("width") or 420)
                project_dir = None
                pid = (body.get("project_id") or body.get("id_project") or "").strip()
                if pid:
                    p = get_project(pid)
                    if p and (p.get("folders") or []):
                        project_dir = p["folders"][0]
                return self._send(
                    200,
                    dses.create_session(
                        fixture,
                        profile=profile,
                        target_width_pt=width,
                        project_dir=project_dir,
                        graph_id=body.get("graph_id"),
                    ),
                )
            except Exception as e:
                return self._send(500, {"error": str(e)})

        if path.startswith("/api/diagram/session/"):
            try:
                import diagram_sessions as dses

                action = path.split("/api/diagram/session/", 1)[-1].strip("/")
                sid = (body.get("id") or body.get("session_id") or "").strip()
                if not sid:
                    return self._send(400, {"error": "session id required"})
                if not dses.get_bundle(sid):
                    return self._send(404, {"error": "unknown session"})
                payload = dses.apply_action(sid, action, body)
                if body.get("save") and dses.get_bundle(sid).get("project_dir"):
                    payload["saved"] = dses.persist_session(sid)
                return self._send(200, payload)
            except ValueError as e:
                return self._send(400, {"error": str(e)})
            except KeyError:
                return self._send(404, {"error": "unknown session"})
            except Exception as e:
                return self._send(500, {"error": str(e)})

        if path == "/api/diagram/propose":
            try:
                import diagram_sessions as dses

                sid = (body.get("id") or body.get("session_id") or "").strip()
                if not sid:
                    return self._send(400, {"error": "session id required"})
                payload = dses.propose_ai_graph(
                    sid,
                    body.get("graph") or body,
                    ref=str(body.get("ref") or ""),
                )
                return self._send(200, payload)
            except KeyError:
                return self._send(404, {"error": "unknown session"})
            except Exception as e:
                return self._send(500, {"error": str(e)})

        if path == "/api/diagram/bind-project":
            try:
                import diagram_sessions as dses

                sid = (body.get("id") or body.get("session_id") or "").strip()
                pid = (body.get("project_id") or body.get("id_project") or "").strip()
                if not sid:
                    return self._send(400, {"error": "session id required"})
                if not pid:
                    return self._send(400, {"error": "project_id required"})
                p = get_project(pid)
                if not p or not (p.get("folders") or []):
                    return self._send(404, {"error": "unknown project or no folder"})
                payload = dses.bind_project(
                    sid,
                    p["folders"][0],
                    graph_id=body.get("graph_id"),
                )
                payload["project_id"] = pid
                payload["project_name"] = p.get("name")
                return self._send(200, payload)
            except KeyError:
                return self._send(404, {"error": "unknown session"})
            except Exception as e:
                return self._send(500, {"error": str(e)})

        if path == "/api/diagram/insert-into-doc":
            try:
                import diagram_sessions as dses
                import diagram_document as ddoc
                import template_lifecycle as tl

                sid = (body.get("id") or body.get("session_id") or "").strip()
                pid = (body.get("project_id") or body.get("id_project") or "").strip()
                if not sid:
                    return self._send(400, {"error": "session id required"})
                if not pid:
                    return self._send(400, {"error": "project_id required"})
                bundle = dses.get_bundle(sid)
                if not bundle:
                    return self._send(404, {"error": "unknown session"})
                payload = dses.session_payload(sid, show_handles=False)
                if payload.get("export_blocked") and not body.get("force"):
                    return self._send(400, {
                        "error": "export blocked — fix validation errors first",
                        "issues": payload.get("issues") or [],
                        "export_blocked": True,
                    })
                p = get_project(pid)
                if not p:
                    return self._send(404, {"error": "unknown project"})
                folder = primary_folder(p)
                if folder is None or not Path(folder).is_dir():
                    return self._send(400, {
                        "error": "project has no folder — choose a project folder first",
                        "code": "need_folder",
                    })
                dses.bind_project(sid, str(folder), graph_id=body.get("graph_id") or payload.get("graph_id"))
                paths = dses.persist_session(sid)
                # Also write SVG next to graph for citation
                try:
                    from diagram_store import diagrams_dir
                    gid = payload.get("graph_id") or "diagram"
                    safe = "".join(c if c.isalnum() or c in "-_" else "_" for c in gid)
                    svg_path = diagrams_dir(folder) / f"{safe}.svg"
                    svg_path.write_text(payload.get("export_svg") or payload.get("svg") or "", encoding="utf-8")
                    paths["svg"] = str(svg_path)
                except Exception:
                    pass

                target = ddoc.resolve_target(
                    body.get("document_type") or body.get("type") or "installation_guide",
                    body.get("section_key") or body.get("section"),
                )
                tf = target["template"]
                template = load_template(tf, folder)
                if not template:
                    # Fall back to registry materialisation as workbench shell
                    try:
                        template = dtr.materialise_template(
                            target["document_type"],
                            project_id=pid,
                            include="required+recommended",
                        )
                        tf = f"{target['document_type']}.json"
                    except LookupError:
                        return self._send(400, {"error": f"template not found: {target['template']}"})

                state, st_folder = project_state_load(p, tf)
                create_missing = body.get("create_if_missing", True)
                doc = state.get("doc") or {}
                if create_missing and (not doc or doc.get("template_file") != tf):
                    tl.create_document_shell(state, tf, template)
                ensure_doc_for_template(st_folder or folder, state, tf)
                ensure_template_sections(state, template)

                section_key = target["section"]
                # If preferred section missing from template, pick first preferred diagram-friendly section
                tpl_keys = {s.get("section_key") for s in (template.get("sections") or [])}
                if section_key not in tpl_keys and tpl_keys:
                    for cand in ("system_overview", "description", "installation",
                                 "work_description", "design_documentation", "findings"):
                        if cand in tpl_keys:
                            section_key = cand
                            break
                    else:
                        section_key = next(iter(tpl_keys))

                graph = payload.get("graph") or {}
                citation = paths.get("graph") or paths.get("svg") or ""
                md = ddoc.diagram_markdown(
                    title=payload.get("title") or graph.get("title") or section_key,
                    graph_id=str(graph.get("id") or payload.get("graph_id") or "diagram"),
                    profile=payload.get("profile") or "wiring",
                    jurisdiction=payload.get("jurisdiction") or "",
                    svg=payload.get("export_svg") or payload.get("svg") or "",
                    revision=str(payload.get("revision") or "A"),
                    citation=citation,
                    lang=body.get("lang") or "en",
                )
                sec = ddoc.insert_into_section(
                    state,
                    section_key=section_key,
                    md=md,
                    svg=payload.get("export_svg") or payload.get("svg") or "",
                    graph=graph,
                    paths=paths,
                    profile=payload.get("profile") or "wiring",
                    replace=bool(body.get("replace", True)),
                )
                project_state_save(p, state, st_folder or folder)
                return self._send(200, {
                    "ok": True,
                    "project_id": pid,
                    "document_type": target["document_type"],
                    "template": tf,
                    "section_key": section_key,
                    "paths": paths,
                    "diagram": sec.get("foldok_diagram"),
                    "open": f"/?project={pid}&template={tf}",
                    "disclaimer": (
                        "Diagram inserted as project evidence. "
                        "Foldok does not claim standard conformity from the figure alone."
                    ),
                })
            except KeyError:
                return self._send(404, {"error": "unknown session"})
            except ValueError as e:
                return self._send(400, {"error": str(e)})
            except Exception as e:
                return self._send(500, {"error": str(e)})

        # ── Capture app folder bridge (foldok_capture) ─────────────────
        if path in ("/api/capture/bind", "/api/capture/publish", "/api/capture/ingest"):
            pid = _pid(body) or body.get("id", "")
            p = get_project(pid)
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = primary_folder(p)
            if folder is None or not folder.is_dir():
                return self._send(400, {"error": "Capture bridge requires a project folder"})
            try:
                from local_app import capture_bridge as cbr

                state, folder = project_state_load(p)
                if path == "/api/capture/bind":
                    result = cbr.capture_bind(folder, p, state)
                elif path == "/api/capture/publish":
                    result = cbr.capture_publish(folder, p, state)
                else:
                    result = cbr.capture_ingest(folder, p, state, by=body.get("by") or "")
                project_state_save(p, state, folder)
                return self._send(200, result)
            except Exception as e:
                return self._send(500, {"error": str(e)})

        # ── WORKORDER 0.60 account / credits ───────────────────────────
        if path == "/api/account/magic-link":
            try:
                return self._send(200, acct.magic_link(body.get("email") or ""))
            except Exception as e:
                return self._send(400, {"error": str(e)})
        if path == "/api/account/verify":
            try:
                return self._send(200, acct.verify(body.get("email") or "", body.get("code") or ""))
            except Exception as e:
                code = getattr(e, "code", None) or "error"
                return self._send(400, {"error": str(e), "code": code})
        if path == "/api/account/guest":
            return self._send(200, acct.try_without_account())
        if path == "/api/account/sign-out":
            return self._send(200, acct.sign_out())
        if path == "/api/account/delete":
            return self._send(200, acct.delete_account())
        if path == "/api/account/topup":
            tok = acct.device_token()
            if not tok:
                return self._send(401, {"error": "Logg inn for å fylle på"})
            try:
                r = acct.get_ledger().topup(tok, float(body.get("amount_eur") or 0))
                return self._send(200, {**r, "account_snap": acct.account_snapshot()})
            except Exception as e:
                code = getattr(e, "code", None)
                return self._send(402 if code else 400, {"error": str(e), "code": code})
        if path == "/api/account/auto-refill":
            tok = acct.device_token()
            if not tok:
                return self._send(401, {"error": "Ikke innlogget"})
            try:
                acc = acct.get_ledger().set_auto_refill(tok, body or {})
                return self._send(200, {"account": acc, **acct.account_snapshot()})
            except Exception as e:
                return self._send(400, {"error": str(e)})
        if path == "/api/account/company":
            tok = acct.device_token()
            if not tok:
                return self._send(401, {"error": "Ikke innlogget"})
            try:
                acc = acct.get_ledger().update_company(tok, body.get("company") or body or {})
                return self._send(200, {"account": acc, **acct.account_snapshot()})
            except Exception as e:
                return self._send(400, {"error": str(e)})
        if path == "/api/account/profile":
            tok = acct.device_token()
            if not tok:
                return self._send(401, {"error": "Ikke innlogget"})
            try:
                acc = acct.get_ledger().update_profile(tok, name=body.get("name"))
                return self._send(200, {"account": acc, **acct.account_snapshot()})
            except Exception as e:
                return self._send(400, {"error": str(e)})

        if path == "/api/pick-folder":
            chosen = pick_folder_dialog((body.get("title") or "Velg prosjektmappe").strip())
            if not chosen:
                return self._send(200, {"folder": None, "cancelled": True})
            return self._send(200, {"folder": chosen, "name": Path(chosen).name})

        if path == "/api/projects":
            folder = (body.get("folder") or "").strip().strip('"')
            name = (body.get("name") or "").strip()
            if folder and not name:
                name = Path(folder).name
            if not name or not folder:
                return self._send(400, {"error": "Navn og mappesti er påkrevd"})
            if not Path(folder).is_dir():
                return self._send(400, {"error": f"Mappen finnes ikke: {folder}"})
            if looks_like_product_repo(folder) and not body.get("force"):
                return self._send(400, {
                    "error": "Dette ser ut som Foldok-produktmappen (ikke et kundeprosjekt). "
                             "UI-skjermbilder i capture/ blir da «kilder» og artefaktmodellen "
                             "beskriver mocken. Koble en kunde- eller jobbmappe. "
                             "Send force:true for å overstyre.",
                    "code": "product_repo",
                })
            projects = load_projects()
            if any(Path(folder) in [Path(f) for f in p["folders"]] for p in projects):
                return self._send(400, {"error": "Denne mappen er allerede i et prosjekt"})
            proj = {"id": uuid.uuid4().hex[:8], "name": name, "folders": [str(Path(folder))]}
            projects.append(proj)
            save_projects(projects)
            return self._send(200, proj)

        if path == "/api/hub/chat":
            msg = (body.get("message") or "").strip()
            caps = hub.load_capabilities()
            history = body.get("history") or []
            ask_fn = fc.ask if KEY_SET else None
            result = hub.hub_chat(msg, caps, history=history, ask_fn=ask_fn)

            # WORKORDER_0.21 A — execute create_project_with_skeleton from chat
            ex = (result.get("execute") or {})
            if ex.get("tool") == "create_project_from_staged":
                created = create_project_from_staged(
                    ex.get("token"),
                    name=ex.get("name") or result.get("project_name"),
                    template_key=ex.get("template_key") or result.get("template_key"),
                    base_dir=body.get("base_dir"),
                    conversation=history,
                )
                if created.get("code") == "need_base_dir":
                    lang = result.get("lang") or hub.detect_lang(msg)
                    result = {
                        "reply": (
                            "Velg basismappe først, så oppretter jeg prosjektet rundt filen."
                            if lang == "no" else
                            "Pick a base folder first, then I'll create the project around the file."
                        ),
                        "kind": "need_base_dir",
                        "code": "need_base_dir",
                        "lang": lang,
                        "execute": ex,
                        "model_called": False,
                        "actions": [{
                            "id": "create_project",
                            "label": "Opprett prosjekt →" if lang == "no" else "Create project →",
                            "token": ex.get("token"),
                        }],
                    }
                elif created.get("error"):
                    result = {"reply": created["error"], "kind": "error",
                              "model_called": False, "actions": []}
                else:
                    lang = result.get("lang") or hub.detect_lang(msg)
                    reply = (
                        f"Opprettet **{created.get('name')}** med "
                        f"{created.get('staged_file') or 'filen'} inne. "
                        f"Samme samtale fortsetter i prosjektet."
                        if lang == "no" else
                        f"Created **{created.get('name')}** with "
                        f"{created.get('staged_file') or 'the file'} inside. "
                        f"Same conversation continues in the project."
                    )
                    result = {
                        "reply": reply,
                        "kind": "created",
                        "lang": lang,
                        "model_called": False,
                        "template_key": created.get("template_key"),
                        "project": {
                            "id": created["id"],
                            "name": created["name"],
                            "folders": created["folders"],
                        },
                        "open_project_id": created["id"],
                        "actions": [],
                        "tool": {"tool": "create_project_from_staged", "ok": True,
                                 "id": created["id"]},
                    }

            if ex.get("tool") == "create_demo_project":
                created = create_demo_project(
                    ex.get("kind") or "contract",
                    base_dir=body.get("base_dir"),
                    name=ex.get("name") or result.get("project_name"),
                    template_key=ex.get("template_key") or result.get("template_key"),
                    conversation=history,
                )
                if created.get("code") == "need_base_dir":
                    lang = result.get("lang") or hub.detect_lang(msg)
                    result = {
                        "reply": (
                            "Velg basismappe først, så lager jeg demosaken."
                            if lang == "no" else
                            "Pick a base folder first, then I'll create the demo case."
                        ),
                        "kind": "need_base_dir",
                        "code": "need_base_dir",
                        "lang": lang,
                        "execute": ex,
                        "model_called": False,
                        "actions": [{"id": "create_demo",
                                     "label": "Lag demosak" if lang == "no" else "Create demo case",
                                     "kind": ex.get("kind") or "contract"}],
                    }
                elif created.get("error"):
                    result = {"reply": created["error"], "kind": "error",
                              "model_called": False, "actions": []}
                else:
                    lang = result.get("lang") or hub.detect_lang(msg)
                    nfiles = len(created.get("demo_files") or [])
                    if lang == "en":
                        reply = (
                            f"Created marked demo **{created.get('name')}** with {nfiles} "
                            f"synthetic files (DEMO_… + banner). Not exportable as a paid "
                            f"document — preview only. Open it to run contract_review."
                        )
                    else:
                        reply = (
                            f"Opprettet merket demosak **{created.get('name')}** med {nfiles} "
                            f"syntetiske filer (DEMO_… + banner). Ikke eksporterbar som betalt "
                            f"dokument — kun forhåndsvisning. Åpne den for kontraktsgjennomgang."
                        )
                    reply = hub.enforce_reply_budget(reply)
                    result = {
                        "reply": reply,
                        "kind": "created_demo",
                        "lang": lang,
                        "model_called": False,
                        "template_key": created.get("template_key"),
                        "project": {
                            "id": created["id"],
                            "name": created["name"],
                            "folders": created["folders"],
                            "demo": True,
                        },
                        "open_project_id": created["id"],
                        "actions": [],
                    }

            if ex.get("tool") == "create_project_with_skeleton":
                # WORKORDER 0.61 C — create document immediately without inventing a folder
                tkey = ex.get("template_key") or result.get("template_key")
                caps = hub.load_capabilities()
                caps_entry = next(
                    (t for t in (caps.get("templates") or []) if t.get("key") == tkey),
                    None,
                ) if tkey else None
                tfile = (caps_entry or {}).get("file") if caps_entry else None
                pname = pio.provisional_name_from_request(
                    msg,
                    (caps_entry or {}).get("name_no") or ex.get("name") or "Nytt dokument",
                )
                if tfile or tkey:
                    created = create_folderless_project(
                        pname,
                        tfile or f"{tkey}.json",
                        output_format=body.get("output_format") or "pdf",
                    )
                    lang = result.get("lang") or hub.detect_lang(msg)
                    reply = (
                        f"Opprettet **{created.get('name_no') or pname}**. "
                        f"Ingen mappe valgt ennå — velg mappe når du vil hente innhold fra filene. "
                        f"Du kan tegne og redigere med en gang."
                        if lang == "no" else
                        f"Created **{created.get('name_no') or pname}**. "
                        f"No folder yet — pick one when you want content from files."
                    )
                    result = {
                        "reply": reply,
                        "kind": "created",
                        "lang": lang,
                        "model_called": False,
                        "template_key": tkey,
                        "template": created.get("template") or tfile,
                        "project": {
                            "id": created["id"],
                            "name": created["name"],
                            "folders": [],
                            "folderless": True,
                        },
                        "open_project_id": created["id"],
                        "open_template": created.get("template") or tfile,
                        "need_folder": True,
                        "actions": [],
                    }
                else:
                    created = create_project_with_skeleton(
                        ex.get("name") or result.get("project_name") or "Prosjekt",
                        tkey,
                        base_dir=body.get("base_dir"),
                        conversation=history,
                    )
                    if created.get("code") == "need_base_dir":
                        # Still folderless shell without template
                        created = create_folderless_project(
                            ex.get("name") or result.get("project_name") or "Prosjekt",
                            None,
                        )
                        result = {
                            "reply": (
                                f"Opprettet prosjektet **{created['name']}** uten mappe. "
                                f"Velg mappe når du er klar."
                            ),
                            "kind": "created",
                            "model_called": False,
                            "project": {
                                "id": created["id"],
                                "name": created["name"],
                                "folders": [],
                                "folderless": True,
                            },
                            "open_project_id": created["id"],
                            "need_folder": True,
                            "actions": [],
                        }
                    elif created.get("error"):
                        result = {
                            "reply": created["error"],
                            "kind": "error",
                            "model_called": False,
                            "actions": [],
                        }
                    else:
                        lang = result.get("lang") or hub.detect_lang(msg)
                        reply = hub.created_folder_reply(
                            created.get("name") or ex.get("name"),
                            created.get("caps_entry"),
                            lang=lang,
                        )
                        reply = hub.enforce_reply_budget(reply, max_words=120)
                        try:
                            st = load_state(created["folders"][0])
                            edchat.append_turn(st, "user", msg)
                            edchat.append_turn(st, "bot", reply)
                            save_state(created["folders"][0], st)
                        except Exception:
                            pass
                        result = {
                            "reply": reply,
                            "kind": "created",
                            "lang": lang,
                            "model_called": False,
                            "template_key": created.get("template_key"),
                            "project": {
                                "id": created["id"],
                                "name": created["name"],
                                "folders": created["folders"],
                            },
                            "open_project_id": created["id"],
                            "checklist_count": created.get("checklist_count"),
                            "actions": [],
                        }

            print(
                f"[hub/chat] model_called={result.get('model_called')} "
                f"kind={result.get('kind')} key_set={KEY_SET} "
                f"msg={msg[:60]!r}",
                flush=True,
            )
            return self._send(200, result)

        if path == "/api/settings":
            settings = hub.load_settings()
            if "projects_base_dir" in body:
                base = (body.get("projects_base_dir") or "").strip().strip('"')
                if base and not Path(base).is_dir():
                    return self._send(400, {"error": f"Basismappe finnes ikke: {base}"})
                settings["projects_base_dir"] = base
                hub.save_settings(settings)
            return self._send(200, {"projects_base_dir": settings.get("projects_base_dir") or ""})

        if path == "/api/project/create-with-skeleton":
            name = (body.get("name") or "").strip()
            if not name:
                return self._send(400, {"error": "Gi prosjektet et navn"})
            created = create_project_with_skeleton(
                name,
                body.get("template_key"),
                base_dir=body.get("base_dir"),
                conversation=body.get("conversation"),
            )
            if created.get("error"):
                code = 400
                payload = {"error": created["error"]}
                if created.get("code"):
                    payload["code"] = created["code"]
                return self._send(code, payload)
            # Don't leak caps_entry blob to client
            created.pop("caps_entry", None)
            return self._send(200, created)

        if path == "/api/add_folder":
            p = get_project(body.get("id", ""))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = (body.get("folder") or "").strip().strip('"')
            if not Path(folder).is_dir():
                return self._send(400, {"error": f"Mappen finnes ikke: {folder}"})
            if Path(folder) in [Path(f) for f in p["folders"]]:
                return self._send(400, {"error": "Mappen er allerede i prosjektet"})
            projects = load_projects()
            for rec in projects:
                if rec["id"] == p["id"]:
                    rec["folders"].append(str(Path(folder)))
            save_projects(projects)
            # WORKORDER 0.55 A3 — auto pre-scan on attach (zero tokens)
            folders = [str(Path(f)) for f in (get_project(p["id"]) or p)["folders"]]
            st = load_state(folders[0])
            thr = st.get("index_last_throughput")
            report = prescan.scan_folders(
                folders, skip_dir_names=SKIP_DIR_NAMES,
                last_throughput=thr, check_cache=True,
            )
            return self._send(200, {"ok": True, "prescan": report,
                                    "card_text": prescan.format_decision_card_no(report)})

        if path == "/api/sources/toggle":
            # Project-level on/off for folders and files (before + after index)
            pid = _pid(body) or body.get("id", "") or body.get("project_id", "")
            p = get_project(pid)
            if not p:
                return self._send(404, {"error": "unknown project"})
            kind = (body.get("kind") or "file").strip().lower()
            if kind not in ("file", "folder"):
                return self._send(400, {"error": "kind must be file or folder"})
            target = (body.get("path") or body.get("file") or body.get("folder") or "").strip()
            if not target:
                return self._send(400, {"error": "path er påkrevd"})
            on = body.get("on")
            if on is None:
                return self._send(400, {"error": "on (true/false) er påkrevd"})
            st = load_state(p["folders"][0], project_id=p.get("id"))
            sel = prescan.toggle_source_selection(
                st.get("source_selection"), kind=kind, path=target, on=bool(on))
            st["source_selection"] = sel
            # Keep document exclude list in sync for files (agent / draft)
            if kind == "file":
                ds.toggle_source(st, target, on=bool(on), full_index=None, artifact=st.get("artifact"))
            save_state(p["folders"][0], st)
            return self._send(200, {
                "ok": True,
                "kind": kind,
                "path": target,
                "on": bool(on),
                "source_selection": sel,
                "excluded_sources": st.get("excluded_sources", []),
            })

        if path == "/api/index/prescan":
            pid = _pid(body) or body.get("id", "") or body.get("project_id", "")
            p = get_project(pid)
            if not p:
                return self._send(404, {"error": "unknown project"})
            st = load_state(p["folders"][0])
            report = prescan.scan_folders(
                p["folders"], skip_dir_names=SKIP_DIR_NAMES,
                last_throughput=st.get("index_last_throughput"),
                check_cache=not bool(body.get("fast")),
            )
            return self._send(200, {
                "ok": True,
                "prescan": report,
                "card_text": prescan.format_decision_card_no(report),
                "index_budget_eur": st.get("index_budget_eur", DEFAULT_INDEX_BUDGET_EUR),
            })

        if path == "/api/index/cancel":
            job_id = body.get("job_id") or body.get("id") or ""
            if not request_cancel(job_id):
                return self._send(404, {"error": "unknown job"})
            return self._send(200, {"ok": True, "job_id": job_id, "cancel_requested": True})

        if path == "/api/index/heartbeat":
            job_id = body.get("job_id") or body.get("id") or ""
            if not heartbeat(job_id):
                return self._send(404, {"error": "unknown job"})
            return self._send(200, {"ok": True, "job_id": job_id})

        if path == "/api/index/resume":
            pid = _pid(body) or body.get("id", "") or body.get("project_id", "")
            p = get_project(pid)
            if not p:
                return self._send(404, {"error": "unknown project"})
            if not KEY_SET:
                return self._send(503, {"error": "ANTHROPIC_API_KEY er ikke satt"})
            st = load_state(p["folders"][0])
            scope = body.get("scope") or st.get("index_scope") or {"mode": "all"}
            if body.get("budget_eur") is not None:
                scope = {**scope, "budget_eur": float(body["budget_eur"])}
            # Extend budget if continuing after ceiling
            if body.get("add_budget_eur"):
                cur = float(st.get("index_budget_eur", DEFAULT_INDEX_BUDGET_EUR) or DEFAULT_INDEX_BUDGET_EUR)
                new_b = cur + float(body["add_budget_eur"])
                st["index_budget_eur"] = new_b
                scope = {**scope, "budget_eur": new_b}
                save_state(p["folders"][0], st)
            sel = prescan.normalize_source_selection(st.get("source_selection"))
            scope = {
                **scope,
                "source_selection": sel,
                "disabled_files": sel["disabled_files"],
                "disabled_folders": sel["disabled_folders"],
            }
            job_id = start_job(run_index, p["folders"], body.get("lang", "no"), scope)
            return self._send(200, {"job": job_id, "scope": scope})

        if path == "/api/index" or path == "/api/reindex":
            # ENGINE_TOOLS `reindex` + WORKORDER 0.55 decision card / scope
            pid = _pid(body) or body.get("id", "") or body.get("project_id", "")
            p = get_project(pid)
            if not p:
                return self._send(404, {"error": "unknown project"})
            if not KEY_SET:
                return self._send(503, {"error": "ANTHROPIC_API_KEY er ikke satt — start serveren med nøkkel"})
            confirm = bool(body.get("confirm"))
            scope = body.get("scope")
            # Explicit scope modes from decision card
            if body.get("mode"):
                scope = {
                    "mode": body.get("mode") or "all",
                    "subfolders": body.get("subfolders"),
                    "newest_n": body.get("newest_n"),
                    "budget_eur": body.get("budget_eur"),
                }
            st = load_state(p["folders"][0])
            report = body.get("prescan")
            if not report:
                report = prescan.scan_folders(
                    p["folders"], skip_dir_names=SKIP_DIR_NAMES,
                    last_throughput=st.get("index_last_throughput"),
                    check_cache=True,
                )
            # Gate: >200 indexable without chosen scope → decision card
            if report.get("needs_decision_card") and not scope and not body.get("force_all"):
                return self._send(200, {
                    "needs_decision_card": True,
                    "prescan": report,
                    "card_text": prescan.format_decision_card_no(report),
                    "index_budget_eur": st.get("index_budget_eur", DEFAULT_INDEX_BUDGET_EUR),
                    "message": "Stor mappe — velg scope før indeksering.",
                })
            plan = idxtools.reindex_plan(
                p["folders"][0], p["folders"], fc, source_files, confirm=confirm)
            if plan["needs_confirm"] and not scope:
                return self._send(200, {
                    "needs_confirm": True,
                    "confirm_threshold": plan["confirm_threshold"],
                    "delta_count": plan["delta_count"],
                    **plan["names"],
                    "total_files": plan["total_files"],
                    "unindexed_count": plan["unindexed_count"],
                    "index_version": plan["index_version"],
                    "prescan": report,
                    "message": (
                        f"{plan['delta_count']} filer endret (terskel "
                        f"{plan['confirm_threshold']}). Bekreft med confirm=true."
                    ),
                })
            scope = scope or st.get("index_scope") or {"mode": "all"}
            if body.get("budget_eur") is not None:
                scope = {**scope, "budget_eur": float(body["budget_eur"])}
            elif "budget_eur" not in (scope or {}):
                scope = {**(scope or {}), "budget_eur": st.get(
                    "index_budget_eur", DEFAULT_INDEX_BUDGET_EUR)}
            sel = prescan.normalize_source_selection(st.get("source_selection"))
            scope = {
                **scope,
                "source_selection": sel,
                "disabled_files": sel["disabled_files"],
                "disabled_folders": sel["disabled_folders"],
            }
            st["index_scope"] = scope
            save_state(p["folders"][0], st)
            job_id = start_job(
                run_index, p["folders"], body.get("lang", "no"), scope)
            return self._send(200, {
                "job": job_id,
                "needs_confirm": False,
                "needs_decision_card": False,
                "scope": scope,
                "prescan": report,
                **plan["names"],
                "total_files": plan["total_files"],
                "unindexed_count": plan["unindexed_count"],
                "index_version": plan["index_version"],
            })

        if path == "/api/diff-index":
            pid = _pid(body) or body.get("id", "") or body.get("project_id", "")
            p = get_project(pid)
            if not p:
                return self._send(404, {"error": "unknown project"})
            result = idxtools.diff_index(
                p["folders"][0], p["folders"], fc, source_files,
                since_version=body.get("since_version"))
            return self._send(200, result)

        if path == "/api/doc/update-from-sources":
            pid = _pid(body) or body.get("id", "") or body.get("project_id", "")
            p = get_project(pid)
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            tf = body.get("template") or body.get("document_id") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            template = load_template(tf)
            if not template or not state.get("doc"):
                return self._send(400, {"error": "Ingen dokument — generer først. Bruk update_document_from_sources bare på eksisterende AST."})
            documents = list_documents(folder, state, templates_list())
            result = idxtools.update_document_from_sources(
                state, template, p["folders"], tf, fc,
                load_index_fn=load_index,
                refresh_code_tables_fn=refresh_code_tables,
                refresh_bom_fn=refresh_bom_section,
                persist_helpers={
                    "source_files": source_files,
                    "project_name": p.get("name"),
                },
                source_ids=body.get("source_ids"),
                mode=body.get("mode") or "merge",
                documents=documents,
            )
            persist_doc(folder, state, tf)
            return self._send(200, result)

        if path == "/api/registry/list":
            return self._send(200, {
                "types": dtr.list_document_types(
                    body.get("industry"),
                    region=body.get("region"),
                    domain=body.get("domain"),
                ),
            })

        if path == "/api/registry/get":
            tid = body.get("type_id") or body.get("id") or ""
            definition = dtr.get_document_type(tid)
            if not definition:
                return self._send(404, {"error": f"Ukjent dokumenttype: {tid}"})
            return self._send(200, {"type": definition})

        if path == "/api/registry/materialise":
            tid = body.get("type_id") or body.get("id") or ""
            try:
                template = dtr.materialise_template(
                    tid,
                    project_id=body.get("project_id") or _pid(body) or None,
                    overrides=body.get("overrides") or {},
                    include=body.get("include") or "required+recommended",
                )
            except LookupError as e:
                return self._send(404, {"error": str(e)})
            return self._send(200, {"template": template})

        if path == "/api/registry/match":
            q = (body.get("query") or body.get("message") or "").strip()
            matches = dtr.match_document_types(q, limit=int(body.get("limit") or 5))
            return self._send(200, {"matches": matches})

        if path == "/api/compliance/frameworks":
            import compliance_engine as ceng
            return self._send(200, {
                "kind": "structural_profiles",
                "disclaimer": ceng.DISCLAIMER,
                "legal_compliance_claimed": False,
                "frameworks": ceng.list_frameworks(
                    region=body.get("region"),
                    domain=body.get("domain"),
                ),
                "regions": list(ceng.REGIONS),
                "domains": list(ceng.DOMAINS),
                "evidence_types": list(ceng.EVIDENCE_TYPES),
            })

        if path == "/api/compliance/suggest":
            import compliance_engine as ceng
            regions = body.get("regions") or []
            domains = body.get("domains") or []
            suggested = ceng.suggest_frameworks(regions, domains)
            return self._send(200, {
                "kind": "structural_profile_suggestions",
                "disclaimer": ceng.DISCLAIMER,
                "legal_compliance_claimed": False,
                "suggested_frameworks": suggested,
                "frameworks": [ceng.get_framework(x) for x in suggested if ceng.get_framework(x)],
            })

        if path == "/api/compliance/gaps":
            import compliance_engine as ceng
            pid = _pid(body) or body.get("id", "")
            p = get_project(pid) if pid else None
            if not p:
                # ephemeral: body carries compliance + optional index
                state = {"compliance": body.get("compliance") or ceng.default_compliance()}
                fw = body.get("frameworks") or state["compliance"].get("frameworks") or []
                status = ceng.package_status(
                    framework_ids=fw,
                    index=body.get("index") or [],
                    state=state,
                )
                return self._send(200, {
                    "gaps": status["gaps"],
                    "package": status,
                    "suggested": False,
                    "disclaimer": ceng.DISCLAIMER,
                    "legal_compliance_claimed": False,
                })
            folder = p["folders"][0]
            state = load_state(folder)
            ceng.ensure_compliance(state)
            if body.get("regions") is not None:
                state["compliance"]["regions"] = list(body.get("regions") or [])
            if body.get("domains") is not None:
                state["compliance"]["domains"] = list(body.get("domains") or [])
            if body.get("frameworks") is not None:
                state["compliance"]["frameworks"] = list(body.get("frameworks") or [])
                state["compliance"]["confirmed"] = bool(body.get("confirmed", True))
            if body.get("suggest"):
                state["compliance"]["suggested_frameworks"] = ceng.suggest_frameworks(
                    state["compliance"].get("regions"),
                    state["compliance"].get("domains"),
                )
            index = []
            try:
                index = load_index(
                    p.get("folders") or [folder], "no",
                    user_facts=state.get("user_facts"),
                    project_name=p.get("name"),
                    cache_only=True,
                ) or []
            except Exception:
                index = []
            package = ceng.project_package_status(state, index)
            save_state(folder, state)
            return self._send(200, {
                "gaps": package["gaps"],
                "package": package,
                "compliance": state.get("compliance"),
                "disclaimer": ceng.DISCLAIMER,
                "legal_compliance_claimed": False,
            })

        # ── Calculation Engine (library formulas + user confirm) ─────
        if path == "/api/calculations/profiles":
            import calculation_engine as calceng
            return self._send(200, {
                "kind": "formula_library",
                "disclaimer": calceng.DISCLAIMER,
                "certified_result": False,
                "profiles": calceng.list_profiles(domain=body.get("domain")),
            })

        if path == "/api/calculations/suggest":
            import calculation_engine as calceng
            suggested = calceng.suggest_profiles(
                domains=body.get("domains") or [],
                intent=body.get("intent") or body.get("q") or "",
            )
            return self._send(200, {
                "disclaimer": calceng.DISCLAIMER,
                "suggested": suggested,
                "profiles": [calceng.get_profile(x) for x in suggested if calceng.get_profile(x)],
            })

        if path == "/api/calculations/propose":
            import calculation_engine as calceng
            profile_id = body.get("profile") or body.get("profile_id")
            if not profile_id:
                return self._send(400, {"error": "profile required"})
            pid = _pid(body) or body.get("id", "")
            p = get_project(pid) if pid else None
            index = body.get("index") or []
            state = {}
            if p:
                folder = p["folders"][0]
                state = load_state(folder)
                if not index:
                    try:
                        index = load_index(
                            p.get("folders") or [folder], "no",
                            user_facts=state.get("user_facts"),
                            project_name=p.get("name"),
                            cache_only=True,
                        ) or []
                    except Exception:
                        index = []
            try:
                calc = calceng.create_calculation(
                    profile_id,
                    index=index,
                    state=state,
                    name=body.get("name"),
                    user_inputs=body.get("inputs") or body.get("user_inputs"),
                    material_id=body.get("material_id") or body.get("material"),
                    section_id=body.get("section_id") or body.get("section"),
                    material_overrides=body.get("material_overrides") or body.get("properties"),
                )
            except ValueError as e:
                return self._send(400, {"error": str(e)})
            if p:
                calceng.upsert_calculation(state, calc)
                save_state(p["folders"][0], state)
            return self._send(200, {
                "calculation": calc,
                "block": calceng.calculation_to_block(calc),
                "text": calceng.render_calculation_text(calc),
                "disclaimer": calceng.DISCLAIMER,
                "certified_result": False,
            })

        if path == "/api/materials/list":
            import materials_engine as meng
            return self._send(200, {
                "disclaimer": meng.DISCLAIMER,
                "code_compliance_claimed": False,
                "materials": meng.list_materials(
                    family=body.get("family"),
                    pack=body.get("pack"),
                    include_templates=body.get("include_templates", True),
                ),
            })

        if path == "/api/materials/get":
            import materials_engine as meng
            mid = body.get("material_id") or body.get("id")
            m = meng.get_material(mid) if mid else None
            if not m:
                return self._send(404, {"error": "unknown material"})
            return self._send(200, {
                "material": m,
                "disclaimer": meng.DISCLAIMER,
                "code_compliance_claimed": False,
            })

        if path == "/api/materials/suggest":
            import materials_engine as meng
            suggested = meng.suggest_material(
                index=body.get("index"),
                state=body.get("state"),
                text=body.get("text") or body.get("q"),
            )
            return self._send(200, {
                "suggested": suggested,
                "materials": [meng.get_material(x) for x in suggested if meng.get_material(x)],
                "disclaimer": meng.DISCLAIMER,
            })

        if path == "/api/sections/list":
            import materials_engine as meng
            return self._send(200, {
                "disclaimer": meng.DISCLAIMER,
                "sections": meng.list_sections(
                    pack=body.get("pack"),
                    series=body.get("series"),
                ),
            })

        if path == "/api/sections/get":
            import materials_engine as meng
            sid = body.get("section_id") or body.get("id")
            s = meng.get_section(sid) if sid else None
            if not s:
                return self._send(404, {"error": "unknown section"})
            return self._send(200, {
                "section": s,
                "disclaimer": meng.DISCLAIMER,
            })

        if path == "/api/knowledge-packs/list":
            import knowledge_registry as kreg
            return self._send(200, {
                "disclaimer": kreg.DISCLAIMER,
                "legal_compliance_claimed": False,
                "packs": kreg.list_packs(),
            })

        if path == "/api/knowledge-packs/get":
            import knowledge_registry as kreg
            pid = body.get("pack_id") or body.get("id")
            pack = kreg.get_pack(pid) if pid else None
            if not pack:
                return self._send(404, {"error": "unknown knowledge pack"})
            return self._send(200, {
                "pack": pack,
                "disclaimer": kreg.DISCLAIMER,
                "legal_compliance_claimed": False,
            })

        if path == "/api/knowledge-packs/gaps":
            import knowledge_registry as kreg
            pid = body.get("pack_id") or body.get("id") or "corrosion_materials"
            gaps = kreg.blocking_gaps(pid, body.get("context") or body.get("facts") or {})
            return self._send(200, {
                "pack_id": pid,
                "gaps": gaps,
                "disclaimer": kreg.DISCLAIMER,
                "legal_compliance_claimed": False,
            })

        if path == "/api/knowledge-packs/render-note":
            import knowledge_registry as kreg
            pid = body.get("pack_id") or body.get("id") or "corrosion_materials"
            note = body.get("note") or body
            text = kreg.render_report_block(pid, note)
            return self._send(200, {
                "pack_id": pid,
                "text": text,
                "disclaimer": kreg.DISCLAIMER,
                "legal_compliance_claimed": False,
            })

        if path == "/api/calculations/set-input":
            import calculation_engine as calceng
            pid = _pid(body) or body.get("id", "")
            calc_id = body.get("calculation_id") or body.get("calc_id")
            key = body.get("key")
            if not key:
                return self._send(400, {"error": "key required"})
            p = get_project(pid) if pid else None
            calc = body.get("calculation")
            state = {}
            if p:
                state = load_state(p["folders"][0])
                calc = calceng.get_calculation(state, calc_id) if calc_id else calc
            if not calc:
                return self._send(404, {"error": "calculation not found"})
            try:
                calc = calceng.set_input(
                    calc, key, body.get("value"),
                    unit=body.get("unit"),
                    source=body.get("source") or "user_entry",
                )
            except ValueError as e:
                return self._send(400, {"error": str(e)})
            if p:
                calceng.upsert_calculation(state, calc)
                save_state(p["folders"][0], state)
            return self._send(200, {
                "calculation": calc,
                "block": calceng.calculation_to_block(calc),
                "disclaimer": calceng.DISCLAIMER,
            })

        if path == "/api/calculations/confirm":
            import calculation_engine as calceng
            pid = _pid(body) or body.get("id", "")
            calc_id = body.get("calculation_id") or body.get("calc_id")
            p = get_project(pid) if pid else None
            calc = body.get("calculation")
            state = {}
            if p:
                state = load_state(p["folders"][0])
                calc = calceng.get_calculation(state, calc_id) if calc_id else calc
            if not calc:
                return self._send(404, {"error": "calculation not found"})
            try:
                calc = calceng.confirm_calculation(
                    calc, confirmed_by=body.get("confirmed_by") or "user",
                )
            except ValueError as e:
                return self._send(400, {"error": str(e)})
            if p:
                calceng.upsert_calculation(state, calc)
                save_state(p["folders"][0], state)
            return self._send(200, {
                "calculation": calc,
                "block": calceng.calculation_to_block(calc),
                "text": calceng.render_calculation_text(calc),
                "disclaimer": calceng.DISCLAIMER,
                "certified_result": False,
            })

        # ── Hybrid Knowledge (project_findings.xlsx) ─────────────────
        if path.startswith("/api/knowledge/"):
            pid = _pid(body) or body.get("id", "") or body.get("project_id", "")
            p = get_project(pid)
            if not p:
                return self._send(404, {"error": "unknown project"})
            primary = p["folders"][0]
            try:
                eng = open_knowledge_engine(primary, enable_vectors=body.get("enable_vectors", True))
            except ImportError as e:
                return self._send(503, {"error": str(e)})

            if path == "/api/knowledge/index":
                result = eng.index_project(force_rebuild=bool(body.get("force_rebuild")))
                return self._send(200, {"ok": True, **result})

            if path == "/api/knowledge/findings":
                rows = eng.get_findings(
                    component=body.get("component"),
                    property_name=body.get("property_name") or body.get("property"),
                    source_file=body.get("source_file"),
                )
                return self._send(200, {"ok": True, "findings": rows, "count": len(rows)})

            if path == "/api/knowledge/update":
                finding = body.get("finding") or body
                if not (finding.get("citation") or finding.get("source_file")):
                    return self._send(400, {"error": "citation or source_file required"})
                fid = eng.update_finding(finding)
                return self._send(200, {"ok": True, "finding_id": fid})

            if path == "/api/knowledge/search":
                q = (body.get("query") or body.get("message") or "").strip()
                rows = eng.semantic_search(q, limit=int(body.get("limit") or 10))
                return self._send(200, {"ok": True, "findings": rows, "count": len(rows)})

            if path == "/api/knowledge/rebuild":
                result = eng.rebuild_index()
                return self._send(200, result)

            if path == "/api/knowledge/import-index":
                index = load_index(p["folders"], body.get("lang", "no"),
                                   load_state(primary).get("user_facts"),
                                   project_name=p.get("name"))
                ids = eng.import_from_index_facts(index)
                return self._send(200, {"ok": True, "imported": len(ids), "ids": ids})

            if path == "/api/knowledge/location":
                return self._send(200, {"ok": True, "location": eng.get_location()})

            if path == "/api/knowledge/set-location":
                loc = eng.set_location(
                    address=body.get("address") or "",
                    municipality=body.get("municipality"),
                    postal_code=body.get("postal_code"),
                    latitude=body.get("latitude"),
                    longitude=body.get("longitude"),
                    location_type=body.get("location_type") or "project_site",
                    citation=body.get("citation"),
                    geocode_if_needed=body.get("geocode_if_needed", True),
                    map_style=body.get("map_style"),
                )
                return self._send(200, {"ok": True, "location": loc})

            if path == "/api/knowledge/generate-map":
                try:
                    rel = eng.generate_location_map(
                        style=body.get("style") or "default",
                        width=int(body.get("width") or 1200),
                        height=int(body.get("height") or 800),
                        zoom=int(body.get("zoom") or 16),
                        color_overrides=body.get("color_overrides"),
                        output_format=body.get("output_format") or "png",
                    )
                except Exception as e:
                    return self._send(400, {"error": str(e)})
                return self._send(200, {"ok": True, "map_image_path": rel,
                                        "location": eng.get_location()})

            if path == "/api/knowledge/propose-map":
                try:
                    proposal = eng.propose_location_map(
                        style=body.get("style") or "default",
                        width=int(body.get("width") or 1200),
                        height=int(body.get("height") or 800),
                        zoom=int(body.get("zoom") or 16),
                        color_overrides=body.get("color_overrides"),
                        output_format=body.get("output_format") or "png",
                        caption=body.get("caption"),
                    )
                except Exception as e:
                    return self._send(400, {"error": str(e)})
                return self._send(200, {"ok": True, **proposal})

            return self._send(404, {"error": f"unknown knowledge path: {path}"})

        if path == "/api/artifact":
            p = get_project(body.get("id", ""))
            if not p:
                return self._send(404, {"error": "unknown project"})
            if not KEY_SET:
                return self._send(503, {"error": "ANTHROPIC_API_KEY er ikke satt"})
            return self._send(200, {"job": start_job(run_artifact, p["folders"], body.get("lang", "no"))})

        if path == "/api/confirm":
            p = get_project(body.get("id", ""))
            if not p:
                return self._send(404, {"error": "unknown project"})
            state = load_state(p["folders"][0])
            if body.get("artifact"):
                state["artifact"] = body["artifact"]  # user-edited — sovereign user
            state["confirmed"] = True
            save_state(p["folders"][0], state)
            return self._send(200, {"ok": True, "artifact": state["artifact"]})

        if path == "/api/artifact/save":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            state = load_state(p["folders"][0])
            art = body.get("artifact")
            if not art or not isinstance(art, dict):
                return self._send(400, {"error": "artifact mangler"})
            state["artifact"] = art
            if body.get("confirm"):
                state["confirmed"] = True
            elif body.get("unconfirm"):
                state["confirmed"] = False
            save_state(p["folders"][0], state)
            return self._send(200, {"ok": True, "artifact": art, "confirmed": state.get("confirmed")})

        if path == "/api/artifact/assist":
            pid = _pid(body)
            try:
                ctx = build_artifact_assist_sources(pid, body.get("lang", "no"))
            except ValueError as e:
                return self._send(400, {"error": str(e)})
            except LookupError as e:
                return self._send(404, {"error": str(e)})
            except IsolationError as e:
                return self._send(500, {"error": str(e)})
            if not KEY_SET:
                return self._send(503, {"error": "ANTHROPIC_API_KEY er ikke satt"})
            msg = (body.get("message") or "").strip()
            if not msg:
                return self._send(400, {"error": "Tom melding"})
            p = ctx["project"]
            state = load_state(ctx["primary"])
            # Prefer server state artifact; body artifact is an edit draft for THIS project only
            art = body.get("artifact") or state.get("artifact") or {}
            # Rebuild context with the draft artifact so the model sees the live model
            if art and art is not state.get("artifact"):
                live_state = {**state, "artifact": art}
                ctx_pack = build_project_chat_context(
                    p, ctx["folders"], ctx["primary"], live_state, ctx["index"],
                    lang=body.get("lang", "no"))
                chat_block = ctx_pack["text"]
                file_count = ctx_pack["file_count"]
            else:
                chat_block = ctx.get("chat_context") or build_project_chat_context(
                    p, ctx["folders"], ctx["primary"], state, ctx["index"],
                    lang=body.get("lang", "no"))["text"]
                file_count = (ctx.get("chat_context_meta") or {}).get("file_count")
            extras = chat_turn_extras(msg, ctx["index"], art, file_count)
            offer_line = ""
            if extras["open_ended"]:
                offer_line = (
                    f"\nThis is an open-ended create ask. End by offering to build a document "
                    f"from what is already here (~€{extras['estimate_eur']:.2f}). "
                    f"At most two questions, only for facts the index cannot contain.\n"
                )
            # WORKORDER_0.20 A — code-first grounded reply (never "helt nytt")
            lang = hub.detect_lang(msg)
            if extras["open_ended"] and (file_count or 0) > 0:
                prose = edchat.open_ended_grounded_reply(
                    project_name=p.get("name") or "",
                    brief=extras["corpus_brief"],
                    artifact=art,
                    known_block=extras["known_block"],
                    estimate_eur=extras["estimate_eur"],
                    lang=lang,
                )
                prose = edchat.scrub_chat_voice(prose)
                edchat.append_turn(state, "user", msg, project_id=p.get("id"))
                edchat.append_turn(state, "bot", prose, project_id=p.get("id"))
                save_state(ctx["primary"], state)
                return self._send(200, {
                    "reply": prose, "patch": None, "note": None, "cost_eur": 0,
                    "project_id": p["id"], "folder": ctx["primary"],
                    "conversation": isolated_conversation(state, p.get("id")),
                    "offer_document": True, "estimate_eur": extras["estimate_eur"],
                    "corpus_brief": extras["corpus_brief"], "kind": "open_ended_grounded",
                })
            prompt = (
                f"{chat_block}\n\n"
                f"{extras['known_block']}\n\n"
                f"{extras['policy']}\n"
                f"{offer_line}\n"
                f"{ctx['banner']}\n"
                f"You are Foldok's ONE project assistant (Checkpoint A / continuous thread).\n"
                f"Ground answers in PROJECT CHAT CONTEXT and ALREADY KNOWN FROM INDEX. "
                f"If SOURCES are empty, CONTEXT is still authoritative — never pretend the project is unknown. "
                f"NEVER invent measurements not in CONTEXT/SOURCES/user. NEVER use other projects.\n"
                f"When the user clarifies the artifact, include a JSON patch. Merge intelligently.\n"
                f"When answering about drawings/sources, cite file names from SOURCES.\n\n"
                f"SOURCES (indexed captions for THIS project only):\n{ctx['captions'] or '(none)'}\n\n"
                f"USER: {msg}\n\n"
                f"Reply with helpful prose first (ground in sentence 1). When the model should change, "
                f"append a fenced JSON block:\n"
                f'```json\n{{"patch":{{"name":"...","purpose":"...","main_components":[{{"name":"...","seen_in":[]}}],'
                f'"hazards":[{{"hazard":"...","source":""}}],"lifecycle_stages":["install","operate"],'
                f'"confidence":0.0-1.0}},"note":"short"}}\n```\n'
                f"Only include patch keys that change. Omit the JSON block if nothing should change."
            )
            raw = fc.ask("artifact_assist", fc.HAIKU, [{"role": "user", "content": prompt}],
                         max_tokens=2000)
            cost = round(fc.LEDGER[-1]["eur"], 4) if fc.LEDGER else 0
            patch = None
            note = None
            m = re.search(r"```json\s*(\{.*?\})\s*```", raw, re.S)
            if m:
                try:
                    data = json.loads(m.group(1))
                    patch = data.get("patch")
                    note = data.get("note")
                except json.JSONDecodeError:
                    pass
            prose = re.sub(r"```json\s*\{.*?\}\s*```", "", raw, flags=re.S).strip()
            prose = edchat.scrub_chat_voice(prose)
            if edchat.reply_violates_policy(prose):
                raw2 = fc.ask("artifact_assist", fc.HAIKU, [{"role": "user", "content":
                    prompt + "\n\nPREVIOUS REPLY VIOLATED POLICY (emoji / 'helt nytt' / Kult!). "
                    "Rewrite: ground first, ≤2 questions, no emoji, end with € document offer."
                }], max_tokens=1200)
                prose = edchat.scrub_chat_voice(
                    re.sub(r"```json\s*\{.*?\}\s*```", "", raw2 or "", flags=re.S).strip())
                cost = round(sum(l["eur"] for l in fc.LEDGER[-2:]), 4) if fc.LEDGER else cost
            edchat.append_turn(state, "user", msg, project_id=p.get("id"))
            edchat.append_turn(state, "bot", prose, project_id=p.get("id"))
            save_state(ctx["primary"], state)
            return self._send(200, {
                "reply": prose, "patch": patch, "note": note, "cost_eur": cost,
                "project_id": p["id"], "folder": ctx["primary"],
                "conversation": isolated_conversation(state, p.get("id")),
                "open_ended": extras["open_ended"],
                "offer_document": extras["open_ended"],
                "estimate_eur": extras["estimate_eur"] if extras["open_ended"] else None,
                "corpus_brief": extras["corpus_brief"],
            })

        if path == "/api/template/intent":
            pid = _pid(body)
            try:
                p, folders, primary = resolve_project(pid)
            except ValueError as e:
                return self._send(400, {"error": str(e)})
            except LookupError as e:
                return self._send(404, {"error": str(e)})
            except IsolationError as e:
                return self._send(500, {"error": str(e)})
            if not KEY_SET:
                return self._send(503, {"error": "ANTHROPIC_API_KEY er ikke satt"})
            story = (body.get("story") or "").strip()
            if not story:
                return self._send(400, {"error": "Fortell hva du trenger — med egne ord"})
            state = load_state(primary)
            art = state.get("artifact") or {"name": p.get("name"), "purpose": "", "confidence": 0.2}
            # PATCH 0.54 FIX 4 — curated template before generic Haiku intent
            # (editor_chat already does this; project /api/template/intent did not)
            import template_lifecycle as tl
            caps = hub.load_capabilities()
            curated = tl.match_curated_template(story, caps)
            if curated:
                return self._send(200, {
                    "choice": curated.get("key"),
                    "file": curated.get("file"),
                    "name_no": curated.get("name_no") or curated.get("name"),
                    "confidence": 1.0,
                    "why_no": "Kuratert mal — treff på dokumenttype i spørsmålet.",
                    "alternatives": [],
                    "no_fit": False,
                    "cost_eur": 0,
                    "project_id": p["id"],
                    "curated": True,
                })
            chat_block = build_project_chat_context(
                p, folders, primary, state, lang=body.get("lang", "no"))["text"]
            result = fc.resolve_template_intent(
                story, art, body.get("lang", "no"), project_context=chat_block)
            result["project_id"] = p["id"]
            return self._send(200, result)

        if path == "/api/template/draft":
            pid = _pid(body)
            try:
                p, folders, primary = resolve_project(pid)
            except ValueError as e:
                return self._send(400, {"error": str(e)})
            except LookupError as e:
                return self._send(404, {"error": str(e)})
            except IsolationError as e:
                return self._send(500, {"error": str(e)})
            if not KEY_SET:
                return self._send(503, {"error": "ANTHROPIC_API_KEY er ikke satt"})
            story = (body.get("story") or "").strip()
            if not story:
                return self._send(400, {"error": "Beskriv hva dokumentet skal gjøre"})
            state = load_state(primary)
            art = state.get("artifact") or {"name": p.get("name"), "purpose": story, "confidence": 0.3}
            chat_block = build_project_chat_context(
                p, folders, primary, state, lang=body.get("lang", "no"))["text"]
            try:
                drafted = fc.draft_template(story, art, body.get("lang", "no"),
                                            project_context=chat_block)
            except Exception as e:
                return self._send(500, {"error": str(e)})
            fname = save_drafted_template(primary, drafted)
            n, first = fc._count_first_section_required_facts(drafted)
            return self._send(200, {
                "ok": True,
                "file": fname,
                "name_no": drafted.get("name_no"),
                "template_key": drafted.get("template_key"),
                "badge": drafted.get("badge"),
                "identity_facts": n,
                "first_section": (first or {}).get("section_key"),
                "required_keys": [rf.get("key") for rf in ((first or {}).get("required_facts") or [])],
                "cost_eur": round(fc.LEDGER[-1]["eur"], 4) if fc.LEDGER else 0,
            })

        if path == "/api/doc/create":
            # WORKORDER 0.59/0.61 — always create a document shell (no chat, no toast)
            p = get_project(body.get("id", ""))
            if not p:
                return self._send(404, {"error": "unknown project"})
            tf = body.get("template") or "sketch_document.json"
            fmt = (body.get("output_format") or "pdf").lower()
            if fmt not in pio.OUTPUT_FORMATS:
                fmt = "pdf"
            template = load_template(tf, primary_folder(p))
            if not template:
                return self._send(400, {"error": f"Mal ikke funnet: {tf}"})
            state, folder = project_state_load(p, tf)
            import template_lifecycle as tl
            created = tl.create_document_shell(state, tf, template)
            if state.get("doc"):
                state["doc"]["output_format"] = fmt
            for d in state.get("documents") or []:
                if d.get("template") == tf:
                    d["output_format"] = fmt
            if (template.get("document_species") == "sketch"
                    or tf == "sketch_document.json"
                    or body.get("sketch")):
                state["doc"]["sketch"] = state["doc"].get("sketch") or {
                    "placeholders": [],
                    "mode": True,
                }
                state["doc"]["sketch"]["mode"] = True
                state["doc"]["document_species"] = "sketch"
                if folder:
                    try:
                        drafts = drafts_dir(folder)
                        drafts.mkdir(exist_ok=True)
                        stub = drafts / f"{template_stem(tf)}.md"
                        if not stub.exists():
                            stub.write_text("# Tomt dokument (skisse)\n\n", encoding="utf-8")
                    except Exception:
                        pass
            art = state.get("artifact") or {}
            conf = float(art.get("confidence") or 0)
            banner = None
            if state.get("folderless") or not folder:
                banner = "Ingen mappe valgt ennå — velg mappe for å hente innhold fra filene dine"
                state["need_folder_banner"] = True
            elif conf < 0.7 or not state.get("confirmed"):
                pct = int(round(conf * 100))
                banner = (
                    f"Forstår prosjektet {pct} % — bekreft modellen for bedre innhold"
                )
            generate = bool(body.get("generate")) and conf >= 0.7 and state.get("confirmed") and bool(folder)
            project_state_save(p, state, folder)
            job = None
            if generate and KEY_SET and not fm_is_form(template) and folder:
                job = start_job(run_generate, p["folders"], tf, body.get("lang", "no"))
            return self._send(200, {
                "ok": True,
                "template": tf,
                "name_no": created.get("name_no"),
                "banner": banner,
                "output_format": fmt,
                "confidence": conf,
                "confirmed": bool(state.get("confirmed")),
                "sketch": bool((state.get("doc") or {}).get("sketch")),
                "folderless": not bool(folder),
                "need_folder": bool(state.get("need_folder_banner")),
                "job": job,
                "generate_started": bool(job),
                "document": next(
                    (d for d in (state.get("documents") or []) if d.get("template") == tf),
                    {"template": tf, "name_no": created.get("name_no")},
                ),
            })

        # WORKORDER 0.61 D — dedicated zero-token sketch geometry APIs
        if path in ("/api/doc/sketch/upsert", "/api/doc/sketch/move", "/api/doc/sketch/delete"):
            p = get_project(body.get("id") or body.get("doc") or "")
            if not p:
                return self._send(404, {"error": "unknown project"})
            tf = body.get("template") or body.get("doc_template")
            state, folder = project_state_load(p, tf)
            import sketch_recognize as sk
            import foldok_compile as fc
            ledger_before = len(fc.LEDGER)
            sketch = (state.get("doc") or {}).setdefault("sketch", {"placeholders": [], "mode": True})
            phs = sketch.setdefault("placeholders", [])
            pid = body.get("block_id") or body.get("placeholder_id")
            ph = None
            if path.endswith("/delete"):
                sketch["placeholders"] = [x for x in phs if x.get("id") != pid]
                phs = sketch["placeholders"]
            elif path.endswith("/move"):
                for i, x in enumerate(phs):
                    if x.get("id") == pid:
                        if "x" in body:
                            x["x"] = float(body["x"])
                        if "y" in body:
                            x["y"] = float(body["y"])
                        phs[i] = sk.recognize_placeholder(x)
                        ph = phs[i]
                        break
                if not ph:
                    return self._send(404, {"error": "placeholder not found"})
            else:  # upsert
                existing = next((x for x in phs if x.get("id") == pid), None) if pid else None
                if existing:
                    if "label" in body:
                        existing["label"] = body.get("label") or ""
                    if body.get("type"):
                        existing["type"] = body["type"]
                    for k in ("x", "y", "w", "h"):
                        if k in body and body[k] is not None:
                            existing[k] = float(body[k])
                    ph = sk.recognize_placeholder(existing)
                    for i, x in enumerate(phs):
                        if x.get("id") == ph.get("id"):
                            phs[i] = ph
                else:
                    ph = sk.new_placeholder(
                        block_type=body.get("type") or "text",
                        x=float(body.get("x") or 45),
                        y=float(body.get("y") or 80),
                        w=float(body.get("w") or 200),
                        h=float(body.get("h") or 80),
                        label=body.get("label") or "",
                        page=int(body.get("page") or 0),
                    )
                    if pid:
                        ph["id"] = pid
                    phs.append(ph)
            _materialize_sketch_sections(state, sk)
            project_state_save(p, state, folder)
            if len(fc.LEDGER) != ledger_before:
                return self._send(500, {"error": "sketch geometry must be zero-token", "token_delta": len(fc.LEDGER) - ledger_before})
            return self._send(200, {
                "ok": True,
                "zero_token": True,
                "token_delta": 0,
                "placeholder": ph,
                "placeholders": sk.sort_placeholders(phs),
            })

        if path == "/api/project/bind-folder":
            p = get_project(body.get("id", ""))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = (body.get("folder") or "").strip().strip('"')
            if not folder:
                return self._send(400, {"error": "Velg en mappe"})
            try:
                pio.bind_folder_to_project(
                    p, folder,
                    load_state_fn=load_state,
                    save_state_fn=save_state,
                    load_projects=load_projects,
                    save_projects=save_projects,
                )
            except Exception as e:
                return self._send(400, {"error": str(e)})
            # refresh project from disk
            p = get_project(body.get("id", ""))
            return self._send(200, {
                "ok": True,
                "project": p,
                "offer_index": True,
                "reply": "Mappe koblet — klar for indeksering.",
            })

        if path == "/api/doc/output-format":
            p = get_project(body.get("id", ""))
            if not p:
                return self._send(404, {"error": "unknown project"})
            tf = body.get("template")
            fmt = (body.get("output_format") or "pdf").lower()
            if fmt not in pio.OUTPUT_FORMATS:
                return self._send(400, {"error": f"Ukjent format: {fmt}"})
            state, folder = project_state_load(p, tf)
            if state.get("doc"):
                state["doc"]["output_format"] = fmt
            for d in state.get("documents") or []:
                if not tf or d.get("template") == tf:
                    d["output_format"] = fmt
            project_state_save(p, state, folder)
            return self._send(200, {"ok": True, "output_format": fmt, "recharged": False})

        if path == "/api/sketch/placeholder":
            p = get_project(body.get("id", ""))
            if not p:
                return self._send(404, {"error": "unknown project"})
            tf = body.get("template")
            state, folder = project_state_load(p, tf)
            import sketch_recognize as sk
            action = body.get("action") or "add"
            sketch = (state.get("doc") or {}).setdefault("sketch", {"placeholders": [], "mode": True})
            phs = sketch.setdefault("placeholders", [])
            if action == "add":
                ph = sk.new_placeholder(
                    block_type=body.get("type") or "text",
                    x=float(body.get("x") or 45),
                    y=float(body.get("y") or 80),
                    w=float(body.get("w") or 200),
                    h=float(body.get("h") or 80),
                    label=body.get("label") or "",
                    page=int(body.get("page") or 0),
                )
                phs.append(ph)
            elif action == "update":
                pid = body.get("placeholder_id")
                for i, ph in enumerate(phs):
                    if ph.get("id") == pid:
                        if "label" in body:
                            ph["label"] = body.get("label") or ""
                        if "type" in body and body["type"]:
                            ph["type"] = body["type"]
                        if "x" in body:
                            ph["x"] = float(body["x"])
                        if "y" in body:
                            ph["y"] = float(body["y"])
                        if "w" in body:
                            ph["w"] = float(body["w"])
                        if "h" in body:
                            ph["h"] = float(body["h"])
                        phs[i] = sk.recognize_placeholder(ph)
                        ph = phs[i]
                        break
                else:
                    return self._send(404, {"error": "placeholder not found"})
            elif action == "fill":
                pid = body.get("placeholder_id")
                import foldok_compile as fc
                folders = p.get("folders") or []
                if not folders:
                    return self._send(400, {
                        "error": "Ingen mappe valgt — velg mappe og indekser før «Fyll fra kilder».",
                        "code": "need_folder",
                    })
                index = load_index(folders, body.get("lang", "no"), state.get("user_facts"))
                for i, ph in enumerate(phs):
                    if ph.get("id") == pid:
                        ledger_before = len(fc.LEDGER)
                        phs[i] = sk.fill_placeholder_from_index(
                            ph, index, state.get("artifact"), lang=body.get("lang", "no"),
                        )
                        ledger_after = len(fc.LEDGER)
                        ph = phs[i]
                        _materialize_sketch_sections(state, sk)
                        project_state_save(p, state, folder)
                        return self._send(200, {
                            "ok": True,
                            "placeholder": ph,
                            "placeholders": sk.sort_placeholders(phs),
                            "token_delta": ledger_after - ledger_before,
                        })
                return self._send(404, {"error": "placeholder not found"})
            elif action == "delete":
                pid = body.get("placeholder_id")
                sketch["placeholders"] = [ph for ph in phs if ph.get("id") != pid]
                phs = sketch["placeholders"]
                ph = None
            else:
                return self._send(400, {"error": f"unknown action: {action}"})
            if action == "add":
                ph = phs[-1] if phs else None
            _materialize_sketch_sections(state, sk)
            project_state_save(p, state, folder)
            return self._send(200, {
                "ok": True,
                "placeholders": sk.sort_placeholders(phs),
                "placeholder": ph if action != "delete" else None,
                "token_delta": 0,
                "zero_token": True,
            })

        if path == "/api/sketch/save-template":
            p = get_project(body.get("id", ""))
            if not p:
                return self._send(404, {"error": "unknown project"})
            tf = body.get("template")
            folder = p["folders"][0]
            state = load_state(folder, tf)
            import sketch_recognize as sk
            phs = ((state.get("doc") or {}).get("sketch") or {}).get("placeholders") or []
            labelled = [ph for ph in phs if (ph.get("label") or "").strip()]
            if len(labelled) < 3:
                return self._send(400, {"error": "Merk minst 3 blokker før du lagrer mal"})
            owned = sk.sketch_to_owned_template(
                labelled,
                name=body.get("name") or "Skisset mal",
                artifact_hint=state.get("artifact"),
            )
            fname = save_company_template(owned)
            regen_capabilities()
            return self._send(200, {
                "ok": True,
                "file": fname,
                "origin": "sketched",
                "badge": "Egen mal",
                "name_no": owned.get("name_no"),
            })

        if path == "/api/generate":
            p = get_project(body.get("id", ""))
            if not p:
                return self._send(404, {"error": "unknown project"})
            if not KEY_SET:
                return self._send(503, {"error": "ANTHROPIC_API_KEY er ikke satt"})
            tf = body.get("template")
            folder = p["folders"][0]
            if not tf or not load_template(tf, folder):
                return self._send(400, {"error": "Velg en mal"})
            return self._send(200, {"job": start_job(run_generate, p["folders"], tf, body.get("lang", "no"))})

        if path == "/api/export":
            p = get_project(body.get("id", ""))
            if not p:
                return self._send(404, {"error": "unknown project"})
            tf = body.get("template")
            if not tf:
                return self._send(400, {"error": "Velg et dokument"})
            state, folder = project_state_load(p, tf)
            fmt = (body.get("output_format")
                   or (state.get("doc") or {}).get("output_format")
                   or "pdf").lower()
            if fmt not in pio.OUTPUT_FORMATS:
                fmt = "pdf"
            # WORKORDER 0.59 D2 — unlabelled sketch placeholders block export
            try:
                import sketch_recognize as sk
                blockers = sk.export_blocking_placeholders(state)
                if blockers:
                    n = len(blockers)
                    return self._send(400, {
                        "error": f"{n} blokker mangler innhold",
                        "blocking": blockers,
                        "reply": f"Kan ikke eksportere — {n} blokker mangler innhold (merk dem først).",
                    })
            except Exception:
                pass
            # WORKORDER_0.23 C2 — demo projects: watermark only, not paid export
            if state.get("demo") or state.get("export_paid_blocked"):
                import demo_project as demo
                template = load_template(tf)
                content = get_draft_md(folder, state, template, tf) if template and folder else (
                    (state.get("doc") and ds.assemble_draft(state, template, state.get("artifact"))) or ""
                )
                if not content:
                    return self._send(404, {"error": "Dokumentet finnes ikke — generer på nytt"})
                content = demo.stamp_demo_watermark(content)
                if not folder:
                    import base64
                    return self._send(200, {
                        "ok": True, "demo": True, "paid_blocked": True,
                        "export_name": "DEMO.md",
                        "download_base64": base64.b64encode(content.encode("utf-8")).decode("ascii"),
                        "reply": "Demosak — vannmerke. Ingen mappe koblet.",
                    })
                out_path, display = write_report_export(folder, template or {"name_no": "DEMO"}, content)
                return self._send(200, {
                    "ok": True, "demo": True, "paid_blocked": True,
                    "export_name": display, "export_path": str(out_path),
                    "reply": ("Dette er en demosak — eksportert med DEMO-vannmerke. "
                              "Betalt eksport er ikke tilgjengelig for syntetisk materiale."),
                })
            template = load_template(tf)
            if template and state.get("doc") and folder and (p.get("folders") or []):
                ensure_figures_in_doc(state, p["folders"], template)
            if folder:
                content = get_draft_md(folder, state, template, tf) if template else read_draft_content(folder, tf)
            else:
                content = ds.assemble_draft(state, template, state.get("artifact")) if template else ""
            if content is None:
                content = ""
            if not str(content).strip() and not ((state.get("doc") or {}).get("sketch") or {}).get("placeholders"):
                return self._send(404, {"error": "Dokumentet finnes ikke — generer på nytt"})
            if folder:
                content = materialize_export_figures(folder, content)
            refs = ds.reference_facts(state)
            if refs:
                lines = ["\n---\n", "## Ansvarserklæring — referanseverdier\n"]
                for f in refs:
                    u = f" {f['unit']}" if f.get("unit") else ""
                    lines.append(f"- **{f['key']}**: {f['value']}{u} ~\n")
                lines.append(
                    "\nVerdier merket ~ er AI-foreslåtte referanseverdier som ikke er "
                    "verifisert mot prosjektets kilder; ansvaret for verifisering ligger "
                    "hos undertegnede.\n"
                )
                content = content.rstrip() + "\n" + "".join(lines)

            def _finish_format_write(content_md, display_hint, extra=None):
                import base64
                extra = extra or {}
                display = display_hint or report_display_name(template or {"name_no": "Dokument"})
                if fmt in ("html", "pptx", "docx"):
                    path, name, notices, raw = expfmt.write_format_export(
                        folder, state, template, fmt=fmt, display_name=display, md_content=content_md)
                    out = {
                        "ok": True, "output_format": fmt, "export_name": name,
                        "notices": notices, "recharged": False, **extra,
                    }
                    if path:
                        out["export_path"] = str(path)
                    if raw is not None:
                        out["download_base64"] = base64.b64encode(raw).decode("ascii")
                    return out
                if not folder:
                    return {
                        "ok": True, "output_format": "pdf", "export_name": f"{display}.md",
                        "download_base64": base64.b64encode((content_md or "").encode("utf-8")).decode("ascii"),
                        "notices": [], "recharged": False, **extra,
                    }
                out_path, display = write_report_export(folder, template, content_md)
                return {
                    "ok": True, "output_format": "pdf", "export_name": display,
                    "export_path": str(out_path), "notices": [], "recharged": False, **extra,
                }

            # WORKORDER 0.60 — charge export against balance (Path B)
            caps = hub.load_capabilities()
            tier, price = acct.export_price_for_template(template or {}, caps)
            docs = state.get("documents", [])
            doc_entry = next((d for d in docs if d.get("template") == tf), None) or {}
            ent = acct.export_entitlement(doc_entry, state)
            snap = acct.account_snapshot()
            watermarked = False
            receipt_meta = None
            pay_mode = (body.get("pay") or "auto").lower()  # auto | charge | preview

            if ent.get("charge") and pay_mode != "preview":
                if not snap.get("signed_in"):
                    # Guest / signed-out: watermarked preview only unless they log in
                    content = acct.stamp_utkast_watermark(content)
                    watermarked = True
                else:
                    bal = float((snap.get("account") or {}).get("balance_eur") or 0)
                    if bal < price:
                        return self._send(402, {
                            "error": f"Saldo €{bal:.2f} er for lav — trenger €{price}.",
                            "code": "insufficient_balance",
                            "price_eur": price,
                            "balance_eur": bal,
                            "need_topup": True,
                        })
                    # Render first; charge; refund on failure
                    try:
                        finished = _finish_format_write(
                            content, report_display_name(template or {}))
                    except Exception as e:
                        return self._send(500, {"error": f"Eksport feilet før trekk: {e}"})
                    pdf_bytes = b""
                    if finished.get("export_path"):
                        pdf_bytes = Path(finished["export_path"]).read_bytes()
                    elif finished.get("download_base64"):
                        import base64
                        pdf_bytes = base64.b64decode(finished["download_base64"])
                    pdf_sha = acct.content_sha256(pdf_bytes or content.encode("utf-8"))
                    block_snap = {
                        "template": tf,
                        "content_sha256": acct.doc_content_fingerprint(state),
                        "sections": list(((state.get("doc") or {}).get("sections") or {}).keys()),
                        "output_format": fmt,
                    }
                    try:
                        if not acct.device_token():
                            return self._send(401, {"error": "Logg inn for betalt eksport"})
                        charged = acct.get_ledger().charge_export(
                            acct.device_token(),
                            tier=tier,
                            project_id=p.get("id") or "",
                            project_name=p.get("name") or "",
                            doc_name=finished.get("export_name") or "export",
                            template=tf,
                            revision=ent.get("revision") or "A",
                            pdf_sha256=pdf_sha,
                            block_snapshot=block_snap,
                            pdf_bytes=pdf_bytes or None,
                        )
                    except acct.MeterDenied as e:
                        return self._send(402, {
                            "error": str(e), "code": e.code,
                            "price_eur": price, "need_topup": True,
                        })
                    except Exception as e:
                        return self._send(500, {"error": str(e)})
                    receipt_meta = charged.get("receipt")
                    acct.mark_document_paid(
                        state, tf,
                        price_eur=price,
                        revision=ent.get("revision") or "A",
                        receipt_id=(receipt_meta or {}).get("id") or "",
                        pdf_sha256=pdf_sha,
                    )
                    for d in state.get("documents", []):
                        if d.get("template") == tf:
                            d["export_name"] = finished.get("export_name")
                            if finished.get("export_path"):
                                d["export_path"] = finished["export_path"]
                            d["output_format"] = fmt
                    project_state_save(p, state, folder)
                    return self._send(200, {
                        **finished,
                        "paid": True, "watermarked": False,
                        "price_eur": price, "revision": ent.get("revision") or "A",
                        "receipt": receipt_meta,
                        "account": acct.account_snapshot(),
                        "blocking_dismissed": ds.blocking_dismissed(state),
                        "reference_facts": refs,
                        "refund_note": "Hvis eksporten feiler, reverseres trekket automatisk.",
                    })

            if pay_mode == "preview" or watermarked:
                if not watermarked:
                    content = acct.stamp_utkast_watermark(content)
                    watermarked = True

            finished = _finish_format_write(
                content, report_display_name(template or {}))
            for d in state.get("documents", []) or []:
                if d.get("template") == tf:
                    d["export_name"] = finished.get("export_name")
                    if finished.get("export_path"):
                        d["export_path"] = finished["export_path"]
                    d["output_format"] = fmt
            project_state_save(p, state, folder)
            return self._send(200, {
                **finished,
                "paid": not watermarked and not ent.get("charge"),
                "reexport_free": not ent.get("charge"),
                "watermarked": watermarked,
                "price_eur": price,
                "blocking_dismissed": ds.blocking_dismissed(state),
                "reference_facts": refs,
                "account": acct.account_snapshot(),
            })

        if path == "/api/doc/refresh-figures":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = primary_folder(p)
            if not folder:
                return self._send(400, {"error": "Ingen mappe valgt — velg mappe først", "code": "need_folder"})
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            template = load_template(tf)
            if not template or not state.get("doc"):
                return self._send(400, {"error": "Ingen dokument å illustrere — generer først"})
            ensure_figures_in_doc(state, p["folders"], template)
            persist_doc(folder, state, tf)
            n = sum(1 for sec in state["doc"]["sections"].values()
                    for m in fc.FIGURE_MARK.finditer(sec.get("md") or ""))
            return self._send(200, {"ok": True, "figures": n, "gaps": state.get("gaps", []),
                                    "gap_summary": ds.gaps_summary(state.get("gaps", []))})

        if path == "/api/doc/resolve-mangler":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            key = (body.get("key") or "").strip()
            value = (body.get("value") or "").strip()
            unit = (body.get("unit") or "").strip() or None
            if not key or not value:
                return self._send(400, {"error": "Nøkkel og verdi er påkrevd"})
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            template = load_template(tf, folder)
            artifact = state.get("artifact")
            index = load_index(p["folders"], body.get("lang", "no"), state.get("user_facts"),
                               project_name=p.get("name"))
            provenance = (body.get("provenance") or "user").strip()
            if provenance == "reference" and not fc.allows_reference_suggest(key):
                return self._send(400, {"error": "Referanseverdi tillates ikke for denne nøkkelen"})
            result = ds.resolve_mangler(state, key, value, unit, template, index, artifact, fc,
                                        provenance=provenance, section=body.get("section"))
            persist_doc(folder, state, tf)
            return self._send(200, {**result, "gap_summary": ds.gaps_summary(result["gaps"])})

        if path == "/api/doc/dismiss":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            key = (body.get("key") or "").strip()
            section = body.get("section", "")
            reason = body.get("reason", "")
            severity = body.get("severity", "warning")
            if not key:
                return self._send(400, {"error": "Nøkkel er påkrevd"})
            if severity == "blocking" and not (reason or "").strip():
                return self._send(400, {"error": "Blokkerende avvisning krever begrunnelse"})
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            template = load_template(tf)
            index = load_index(p["folders"], body.get("lang", "no"), state.get("user_facts"))
            result = ds.dismiss_mangler(state, key, section, reason, severity, template,
                                          state.get("artifact"), index, fc)
            persist_doc(folder, state, tf)
            return self._send(200, {**result, "gap_summary": ds.gaps_summary(result["gaps"]),
                                    "blocking_dismissed": ds.blocking_dismissed(state)})

        if path == "/api/doc/undismiss":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            key = body.get("key")
            section = body.get("section", "")
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            template = load_template(tf)
            index = load_index(p["folders"], "no", state.get("user_facts"))
            result = ds.undismiss_mangler(state, key, section, template, index, state.get("artifact"), fc)
            persist_doc(folder, state, tf)
            return self._send(200, {**result, "gap_summary": ds.gaps_summary(result["gaps"])})

        if path == "/api/doc/exclude-figure":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            file = (body.get("file") or "").strip()
            section = (body.get("section") or "").strip()
            page = int(body.get("page") or 0)
            if not file or not section:
                return self._send(400, {"error": "file og section er påkrevd"})
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            result = ds.exclude_figure(state, section, file, page)
            ds.add_version(state, "user", "figure", f"Fjernet illustrasjon {Path(file).name}", section=section)
            persist_doc(folder, state, tf)
            n = sum(1 for sec in (state.get("doc") or {}).get("sections", {}).values()
                    for _ in fc.FIGURE_MARK.finditer(sec.get("md") or ""))
            return self._send(200, {**result, "figures": n,
                                    "excluded_figures": state.get("excluded_figures", [])})

        if path == "/api/doc/restore-figure":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            file = (body.get("file") or "").strip()
            page = int(body.get("page") or 0)
            if not file:
                return self._send(400, {"error": "file er påkrevd"})
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            template = load_template(tf)
            result = ds.restore_figure(state, file, page)
            ensure_figures_in_doc(state, p["folders"], template)
            persist_doc(folder, state, tf)
            n = sum(1 for sec in (state.get("doc") or {}).get("sections", {}).values()
                    for _ in fc.FIGURE_MARK.finditer(sec.get("md") or ""))
            return self._send(200, {**result, "figures": n,
                                    "excluded_figures": state.get("excluded_figures", [])})

        if path == "/api/doc/toggle-figure":
            # Unified media toggle (SOURCE_INTERACTION S5): on=false → exclude, on=true → restore
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            file = (body.get("file") or "").strip()
            page = int(body.get("page") or 0)
            on = bool(body.get("on"))
            if not file:
                return self._send(400, {"error": "file er påkrevd"})
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            template = load_template(tf)
            if on:
                result = ds.restore_figure(state, file, page)
                ensure_figures_in_doc(state, p["folders"], template)
            else:
                section = (body.get("section") or illustration_section_key(template) or "").strip()
                result = ds.exclude_figure(state, section, file, page)
                ds.add_version(state, "user", "figure",
                               f"Fjernet illustrasjon {Path(file).name}", section=section)
            persist_doc(folder, state, tf)
            return self._send(200, {**result, "on": on,
                                    "excluded_figures": state.get("excluded_figures", [])})

        if path == "/api/doc/toggle-source":
            # S1: source file participates in / is excluded from THIS document. Index unchanged.
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            file = (body.get("file") or "").strip()
            if not file:
                return self._send(400, {"error": "file er påkrevd"})
            on = bool(body.get("on"))
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            template = load_template(tf)
            result = ds.toggle_source(state, file, on,
                                      full_index=load_index(p["folders"], "no", state.get("user_facts")),
                                      artifact=state.get("artifact"))
            ds.add_version(state, "user", "source",
                           ("Aktiverte kilde " if on else "Slo av kilde ") + Path(file).name)
            refresh_code_tables(state, p["folders"], tf, body.get("lang", "no"))
            if template:
                index = load_active_index(state, p["folders"], "no")
                state["gaps"] = ds.gaps_for_document(state, template, index,
                                                    state.get("artifact"), fc, fast=True)
            persist_doc(folder, state, tf)
            return self._send(200, {**result, "gaps": state.get("gaps", []),
                                    "gap_summary": ds.gaps_summary(state.get("gaps", [])),
                                    "source_citation_warnings": state.get("source_citation_warnings", []),
                                    "impacts": result.get("impacts", [])})

        if path == "/api/doc/cell-candidates":
            # S4: facts eligible for one table cell — zero tokens
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            section = body.get("section") or ""
            column = body.get("column") or ""
            row_key = body.get("row_key") or ""
            cols = fc.TABLE_COLUMNS.get(section) or []
            col = next((c for c in cols if c["id"] == column), None)
            if not col:
                return self._send(400, {"error": f"Ukjent kolonne {column} i {section}"})
            if not col.get("editable"):
                return self._send(400, {"error": "Kolonnen kan ikke redigeres"})
            fact_key = col.get("key")
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            index = load_active_index(state, p["folders"], "no")
            source_file = row_key.split("|", 1)[1] if "|" in row_key else None
            cands = fc.search_fact_candidates(index, fact_key) if fact_key else []
            # facts from the row's own source file first
            if source_file:
                cands.sort(key=lambda c: 0 if c.get("file") == source_file else 1)
            return self._send(200, {"column": column, "fact_key": fact_key,
                                    "source_file": source_file, "candidates": cands[:8]})

        if path == "/api/doc/edit-cell":
            # S4: inline table edit — cite existing fact or create verified user fact
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            section = body.get("section") or ""
            row_key = body.get("row_key") or ""
            column = body.get("column") or ""
            value = (body.get("value") or "").strip()
            fact_id = body.get("fact_id")
            if not section or not row_key or not column:
                return self._send(400, {"error": "section, row_key og column er påkrevd"})
            clear = bool(body.get("clear"))
            if not value and not fact_id and not clear:
                return self._send(400, {"error": "value, fact_id eller clear er påkrevd"})
            cols = fc.TABLE_COLUMNS.get(section) or []
            col = next((c for c in cols if c["id"] == column), None)
            if not col or not col.get("editable"):
                return self._send(400, {"error": "Kolonnen kan ikke redigeres"})
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            template = load_template(tf)
            if not template or not state.get("doc"):
                return self._send(400, {"error": "Ingen dokument — generer først"})
            index = load_active_index(state, p["folders"], "no")
            fact_key = col.get("key")
            src_note = row_key.split("|", 1)[1] if "|" in row_key else "tabell"
            if clear:
                state["cell_overrides"] = [
                    o for o in state.get("cell_overrides", [])
                    if not (o.get("section") == section and o.get("row_key") == row_key
                            and o.get("column") == column)]
                ds.add_version(state, "user", "cell",
                               f"Tilbakestilte {col['label'].lower()} til kildeverdi", section=section)
                refresh_code_tables(state, p["folders"], tf, body.get("lang", "no"))
                state["gaps"] = ds.gaps_for_document(state, template, index,
                                                    state.get("artifact"), fc, fast=True)
                persist_doc(folder, state, tf)
                return self._send(200, {"ok": True, "cleared": True,
                                        "gaps": state.get("gaps", []),
                                        "gap_summary": ds.gaps_summary(state.get("gaps", []))})
            if fact_id:
                fact = None
                for e in index:
                    for f in e.get("facts", []):
                        if f["id"] == fact_id:
                            fact = {**f, "source_location": f.get("source_location") or e["file"]}
                if not fact:
                    return self._send(404, {"error": "Faktum ikke funnet"})
                ds.set_cell_override(state, section, row_key, column,
                                     fact["value"], fact_id=fact_id)
                summary = f"Bekreftet {col['label'].lower()} fra {Path(fact['source_location']).name}"
                # L2 — remember column/label → fact key when user corrects mapping
                if fact_key and fact.get("key") and fact["key"] != fact_key:
                    learning.record_alias(fact_key, fact["key"], kind="column")
                lab = (col.get("label") or "").strip().lower()
                if lab and fact.get("key"):
                    learning.record_alias(lab, fact["key"], kind="alias")
            else:
                # verified user fact — never silently invents (hard rule)
                fact = {"id": ds.next_user_fact_id(state), "key": fact_key or column,
                        "fact_type": "spec", "value": value, "unit": None,
                        "verified_by_user": True,
                        "source_location": f"oppgitt manuelt av bruker ({src_note})",
                        "created": ds.iso_now()}
                state.setdefault("user_facts", []).append(fact)
                ds.set_cell_override(state, section, row_key, column, value, verified=True)
                summary = f"Skrev inn {col['label'].lower()}: {value}"
            ds.add_version(state, "user", "cell", summary, section=section)
            # re-render the table from stored structure (single-cell swap, no recompile)
            sec = (state["doc"].get("sections") or {}).get(section) or {}
            if sec.get("table"):
                _store_table_section(state, section, sec["table"], body.get("lang", "no"))
            else:
                refresh_code_tables(state, p["folders"], tf, body.get("lang", "no"))
            state["gaps"] = ds.gaps_for_document(state, template, index,
                                                state.get("artifact"), fc, fast=True)
            persist_doc(folder, state, tf)
            return self._send(200, {"ok": True, "override": True,
                                    "gaps": state.get("gaps", []),
                                    "gap_summary": ds.gaps_summary(state.get("gaps", []))})

        if path == "/api/doc/refresh-spec-overview":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            template = load_template(tf)
            if not template or not state.get("doc"):
                return self._send(400, {"error": "Ingen dokument — generer først"})
            if not refresh_spec_overview_section(state, p["folders"], tf, body.get("lang", "no")):
                return self._send(400, {"error": "Malen har ingen spesifikasjonsoversikt"})
            index = load_index(p["folders"], "no", state.get("user_facts"))
            state["gaps"] = ds.gaps_for_document(state, template, index, state.get("artifact"), fc, fast=True)
            persist_doc(folder, state, tf)
            return self._send(200, {"ok": True, "gaps": state.get("gaps", []),
                                    "gap_summary": ds.gaps_summary(state.get("gaps", []))})

        if path == "/api/doc/refresh-doc-control":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            template = load_template(tf)
            if not template or not state.get("doc"):
                return self._send(400, {"error": "Ingen dokument — generer først"})
            if not refresh_doc_control_section(state, p["folders"], tf, body.get("lang", "no")):
                return self._send(400, {"error": "Malen har ingen dokumentkontroll-seksjon"})
            index = load_index(p["folders"], "no", state.get("user_facts"))
            state["gaps"] = ds.gaps_for_document(state, template, index, state.get("artifact"), fc, fast=True)
            persist_doc(folder, state, tf)
            return self._send(200, {"ok": True, "gaps": state.get("gaps", []),
                                    "gap_summary": ds.gaps_summary(state.get("gaps", []))})

        if path == "/api/doc/refresh-bom":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            template = load_template(tf)
            if not template or not state.get("doc"):
                return self._send(400, {"error": "Ingen dokument — generer først"})
            if not refresh_bom_section(state, p["folders"], tf, body.get("lang", "no")):
                return self._send(400, {"error": "Malen har ingen BOM-seksjon"})
            index = load_index(p["folders"], "no", state.get("user_facts"))
            state["gaps"] = ds.gaps_for_document(state, template, index, state.get("artifact"), fc, fast=True)
            persist_doc(folder, state, tf)
            n = len(be.aggregate_bom(index))
            return self._send(200, {"ok": True, "bom_lines": n,
                                    "gaps": state.get("gaps", []),
                                    "gap_summary": ds.gaps_summary(state.get("gaps", []))})

        if path == "/api/suggestions/dismiss":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            name = (body.get("name") or "").strip()
            if not name:
                return self._send(400, {"error": "name er påkrevd"})
            state = load_state(folder)
            dismissed = state.setdefault("dismissed_suggestions", [])
            if not any(d.get("name") == name for d in dismissed):
                dismissed.append({"name": name, "at": ds.iso_now(),
                                  "reason": body.get("reason", "")})
            save_state(folder, state)
            tf = state.get("active_template") or state.get("template")
            return self._send(200, {"ok": True,
                                    "suggestions": active_suggestions(state, p["folders"], tf)})

        if path == "/api/suggestions/create":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            tkey = body.get("template_key")
            tf = body.get("template") or template_file_for_key(tkey)
            if not tf:
                return self._send(400, {"error": "Ukjent dokumentmal"})
            if not load_state(p["folders"][0]).get("confirmed"):
                return self._send(400, {"error": "Bekreft artefaktmodellen først (steg 2)"})
            job_id = start_job(run_generate, p["folders"], tf, body.get("lang", "no"))
            return self._send(200, {"job_id": job_id, "template": tf})

        if path == "/api/doc/extract-targeted":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            if not KEY_SET:
                return self._send(503, {"error": "ANTHROPIC_API_KEY er ikke satt"})
            key = body.get("key")
            rel = body.get("file")
            hint = body.get("hint", "")
            if not key or not rel:
                return self._send(400, {"error": "key og file er påkrevd"})
            path_f, cache_dir = find_source(p["folders"], rel)
            if not path_f:
                return self._send(404, {"error": "Filen finnes ikke i prosjektet"})
            import hashlib
            sha = hashlib.sha256(path_f.read_bytes()).hexdigest()
            cache = cache_dir / f"{sha}.json"
            entry = fc.read_json_file(cache) if cache.exists() else fc.index_file(
                path_f, body.get("lang", "no"), cache_dir, rel_name=rel)
            fact, cost = fc.extract_targeted(path_f, rel, entry, cache, key, hint, body.get("lang", "no"))
            if not fact:
                return self._send(200, {"found": False, "cost_eur": round(cost, 4),
                                        "message": f"Fant ikke {key} i {rel}"})
            fact["provenance"] = "extracted_targeted"
            folder = p["folders"][0]
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            template = load_template(tf)
            index = load_index(p["folders"], "no", state.get("user_facts"))
            result = ds.apply_cited_fact(state, key, fact, template, index, state.get("artifact"), fc)
            if body.get("section") == "bom" or "bom" in (state.get("doc") or {}).get("sections", {}):
                refresh_bom_section(state, p["folders"], tf, body.get("lang", "no"))
            persist_doc(folder, state, tf)
            return self._send(200, {"found": True, "cost_eur": round(cost, 4), **result,
                                    "gap_summary": ds.gaps_summary(result["gaps"])})

        if path == "/api/doc/apply-cited":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            key = body.get("key")
            fact_id = body.get("fact_id")
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            template = load_template(tf)
            index = load_index(p["folders"], "no", state.get("user_facts"))
            fact = None
            for e in index:
                for f in e.get("facts", []):
                    if f["id"] == fact_id:
                        fact = {**f, "source_location": f.get("source_location") or e["file"]}
            if not fact:
                return self._send(404, {"error": "Faktum ikke funnet"})
            result = ds.apply_cited_fact(state, key or fact["key"], fact, template, index,
                                         state.get("artifact"), fc)
            persist_doc(folder, state, tf)
            return self._send(200, {**result, "gap_summary": ds.gaps_summary(result["gaps"])})

        if path == "/api/doc/gap-search":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            key = body.get("key")
            if not key:
                return self._send(400, {"error": "key er påkrevd"})
            folder = p["folders"][0]
            state = load_state(folder)
            index = load_index(p["folders"], "no", state.get("user_facts"))
            cands = fc.search_fact_candidates(index, key)
            return self._send(200, {"candidates": cands, "found": len(cands) > 0})

        if path == "/api/doc/gap-guide":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            key = body.get("key")
            if not key:
                return self._send(400, {"error": "key er påkrevd"})
            folder = p["folders"][0]
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            index = load_index(p["folders"], "no", state.get("user_facts"))
            templates = templates_list(p)
            documents = list_documents(folder, state, templates)
            guide = fc.gap_guide(key, body.get("section", ""), index,
                                 state.get("artifact"), documents)
            return self._send(200, guide)

        if path == "/api/doc/reference-suggest":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            key = (body.get("key") or "").strip()
            if not key:
                return self._send(400, {"error": "key er påkrevd"})
            if not fc.allows_reference_suggest(key, body.get("severity")):
                return self._send(200, {"ok": False, "allowed": False})
            if not KEY_SET:
                return self._send(503, {"error": "ANTHROPIC_API_KEY er ikke satt"})
            folder = p["folders"][0]
            state = load_state(folder)
            folders = p.get("folders") or [folder]
            index = load_index(folders, "no", state.get("user_facts"), project_name=p.get("name"))
            chat_block = build_project_chat_context(
                p, folders, folder, state, index, lang=body.get("lang", "no"))["text"]
            sug = fc.reference_suggest(key, state.get("artifact") or {"name": p.get("name")},
                                       body.get("lang", "no"), project_context=chat_block)
            if not sug:
                return self._send(200, {"ok": False, "allowed": True})
            cost = round(fc.LEDGER[-1]["eur"], 4) if fc.LEDGER else 0
            return self._send(200, {"ok": True, "allowed": True, **sug, "cost_eur": cost})

        if path == "/api/doc/fill-known-gaps":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            template = load_template(tf)
            if not template or not state.get("doc"):
                return self._send(400, {"error": "Ingen dokument — generer først"})
            index = load_index(p["folders"], "no", state.get("user_facts"))
            templates = templates_list(p)
            documents = list_documents(folder, state, templates)
            keys = body.get("keys")
            guide_refresh = None
            if keys and len(keys) == 1:
                g0 = fc.gap_guide(keys[0], body.get("section", ""), index,
                                   state.get("artifact"), documents)
                guide_refresh = g0.get("refresh_section")
            if guide_refresh == "spec_overview":
                refresh_spec_overview_section(state, p["folders"], tf, body.get("lang", "no"))
            elif guide_refresh == "doc_control":
                refresh_doc_control_section(state, p["folders"], tf, body.get("lang", "no"))
            elif guide_refresh == "drawings_register":
                from foldok_compile import compile_drawings_register
                sec = (state.get("doc") or {}).get("sections", {}).get("drawings_register")
                if sec is not None:
                    sec["md"] = compile_drawings_register(index, body.get("lang", "no"))
            result = fc.fill_known_gaps(state, template, index, state.get("artifact"), fc,
                                        keys_only=keys, documents=documents)
            persist_doc(folder, state, tf)
            return self._send(200, {**result, "gap_summary": ds.gaps_summary(result["gaps"])})

        if path == "/api/doc/chat":
            """ONE_AGENT_SPEC — continuous editor assistant with tools."""
            pid = _pid(body)
            try:
                p, folders, primary = resolve_project(pid)
            except ValueError as e:
                return self._send(400, {"error": str(e)})
            except LookupError as e:
                return self._send(404, {"error": str(e)})
            except IsolationError as e:
                return self._send(500, {"error": str(e)})
            msg = (body.get("message") or "").strip()
            if not msg:
                return self._send(400, {"error": "Tom melding"})
            tf = body.get("template") or load_state(primary).get("active_template")
            state = load_state(primary, tf, project_id=p.get("id"))
            template = load_template(tf, primary) if tf else None
            gaps = state.get("gaps") or []
            scope = body.get("section")
            annotations = body.get("annotations") or []
            annot_execute = bool(body.get("annot_execute"))
            edchat.append_turn(state, "user", msg, project_id=p.get("id"))

            # Accept pending reference offer
            pend = state.get("chat_pending") or {}
            if pend.get("action") == "accept_reference" and re.search(
                    r"bruk referansen|bruk den|godta|ok\b|ja\b", msg, re.I):
                index = load_index(folders, "no", state.get("user_facts"), project_name=p.get("name"))
                result = ds.resolve_mangler(
                    state, pend["key"], pend["value"], pend.get("unit"),
                    template, index, state.get("artifact"), fc,
                    provenance="reference", section=pend.get("section"))
                if tf:
                    persist_doc(primary, state, tf)
                else:
                    save_state(primary, state)
                reply = f"Lagt inn som ubekreftet referanse ~ for `{pend['key']}`."
                state["chat_pending"] = None
                tool_result = {"tool": "resolve_mangler", "provenance": "reference",
                               "gap_summary": ds.gaps_summary(result["gaps"])}
                edchat.append_turn(state, "bot", reply, meta=tool_result, project_id=p.get("id"))
                save_state(primary, state)
                return self._send(200, {
                    "reply": reply, "kind": "tool", "tool": tool_result,
                    "conversation": isolated_conversation(state, p.get("id")),
                    "chat_pending": None,
                    "gaps": state.get("gaps") or [],
                    "gap_summary": ds.gaps_summary(state.get("gaps") or []),
                })

            route = edchat.route_editor_message(
                msg, state, gaps, scope_section=scope, template=template,
                caps=hub.load_capabilities(), annotations=annotations)

            # WORKORDER 0.56 §C5 — hard stop create-doc tools while marks pending
            _CREATE_TOOLS = {
                "run_generate", "create_document", "draft_template_rung3",
                "materialise_template", "recreate_form", "accept_drafted_template",
                "propose_form_template",
            }
            ex = (route or {}).get("execute") or {}
            _named = edchat.explicitly_names_document(msg, hub.load_capabilities())
            if (annotations and not _named and (
                    ex.get("tool") in _CREATE_TOOLS
                    or (route or {}).get("kind") in (
                        "propose_generate", "propose_generate_reask", "dispatch_generate")
            )):
                labels = ", ".join(
                    (a.get("chip") or a.get("kind") or "?") for a in annotations[:4])
                route = {
                    "reply": (
                        f"Du har {len(annotations)} merke(r) som venter ({labels}). "
                        "Jeg oppretter ikke dokument mens merker er åpne — "
                        "med mindre du eksplisitt navngir dokumentet/malen."
                    ),
                    "kind": "annot_blocks_create",
                    "annotations": annotations,
                }

            if annot_execute and annotations and not route.get("execute"):
                # Model-path after tegnelag Utfør / chat Enter (C5)
                ctx = edchat.format_annotations_context(annotations)
                has_note = any((a.get("note") or "").strip() for a in annotations)
                if has_note and KEY_SET and route.get("kind") not in (
                        "annot_blocks_create",):
                    route = {
                        "kind": "annot_execute",
                        "need_model": True,
                        "annotations": annotations,
                        "_annot_ctx": ctx,
                    }
                else:
                    route = {
                        "reply": (
                            f"Forstått — tegnelag:\n"
                            + "\n".join(
                                f"• {a.get('chip') or a.get('kind')}"
                                + (f" — «{(a.get('note') or '')[:80]}»"
                                   if a.get("note") else "")
                                for a in annotations
                            )
                            + "\nJeg behandler «dette/her» som disse merkene."
                        ),
                        "kind": "annot_execute",
                        "annotations": annotations,
                        "_annot_ctx": ctx,
                    }

            if route.get("set_pending") is not None:
                state["chat_pending"] = route["set_pending"]
            if route.get("clear_pending"):
                state["chat_pending"] = None

            reply = route.get("reply") or ""
            tool_result = None
            execute = route.get("execute")
            chat_extras = None

            if execute and execute.get("tool") == "list_gaps":
                reply = edchat.format_gaps_reply(gaps, scope)
                tool_result = {"tool": "list_gaps", "ok": True, "count": len(gaps)}

            elif execute and execute.get("tool") == "diff_index":
                lang = hub.detect_lang(msg)
                diff = idxtools.diff_index(
                    primary, folders, fc, source_files,
                    since_version=execute.get("since_version"))
                reply = idxtools.format_diff_reply(diff, lang=lang)
                if diff.get("live_unindexed"):
                    reply += (
                        f"\n\n{diff['live_unindexed']} filer er ikke indeksert ennå — si «reindekser» for å oppdatere."
                        if lang != "en" else
                        f"\n\n{diff['live_unindexed']} files not indexed yet — say «reindex» to update."
                    )
                tool_result = {"tool": "diff_index", "ok": True, **diff}

            elif execute and execute.get("tool") == "reindex":
                lang = hub.detect_lang(msg)
                confirm = bool(execute.get("confirm"))
                plan = idxtools.reindex_plan(
                    primary, folders, fc, source_files, confirm=confirm)
                if plan["needs_confirm"]:
                    n = plan["delta_count"]
                    reply = (
                        f"Jeg fant **{n}** nye/endrede/fjernede filer (terskel "
                        f"{plan['confirm_threshold']}). Skal jeg reindeksere nå?"
                        if lang != "en" else
                        f"Found **{n}** added/changed/removed files (threshold "
                        f"{plan['confirm_threshold']}). Reindex now?"
                    )
                    reply += "\n\n" + idxtools.format_diff_reply(plan["diff"], lang=lang)
                    state["chat_pending"] = {
                        "action": "reindex", "confirm": True,
                        "delta_count": n,
                    }
                    tool_result = {
                        "tool": "reindex", "ok": False, "needs_confirm": True,
                        **plan["names"], "total_files": plan["total_files"],
                        "index_version": plan["index_version"],
                        "actions": [
                            {"id": "confirm_generate", "label": "Ja — reindekser" if lang != "en" else "Yes — reindex"},
                        ],
                    }
                else:
                    if not KEY_SET:
                        reply = ("API-nøkkel mangler — kan ikke indeksere."
                                 if lang != "en" else "API key missing — cannot index.")
                        tool_result = {"tool": "reindex", "ok": False}
                    else:
                        job_id = start_job(run_index, folders, lang if lang else "no")
                        names = plan["names"]
                        bits = []
                        if names["added"]:
                            bits.append(f"+{len(names['added'])}")
                        if names["changed"]:
                            bits.append(f"~{len(names['changed'])}")
                        if names["removed"]:
                            bits.append(f"-{len(names['removed'])}")
                        delta = " · ".join(bits) if bits else "ingen delta"
                        reply = (
                            f"Reindekserer… ({delta}, {plan['total_files']} filer totalt)."
                            if lang != "en" else
                            f"Reindexing… ({delta}, {plan['total_files']} files total)."
                        )
                        tool_result = {
                            "tool": "reindex", "ok": True, "job_id": job_id,
                            **names, "total_files": plan["total_files"],
                            "index_version": plan["index_version"],
                        }

            elif execute and execute.get("tool") == "update_document_from_sources":
                lang = hub.detect_lang(msg)
                if not template or not state.get("doc"):
                    reply = (
                        "Ingen åpent dokument å oppdatere — lag dokumentet først, "
                        "eller bruk reindeks + diff før du ber om oppdatering."
                        if lang != "en" else
                        "No open document to update — create one first."
                    )
                    tool_result = {"tool": "update_document_from_sources", "ok": False}
                else:
                    documents = list_documents(primary, state, templates_list())
                    result = idxtools.update_document_from_sources(
                        state, template, folders, tf, fc,
                        load_index_fn=load_index,
                        refresh_code_tables_fn=refresh_code_tables,
                        refresh_bom_fn=refresh_bom_section,
                        persist_helpers={
                            "source_files": source_files,
                            "project_name": p.get("name"),
                        },
                        source_ids=execute.get("source_ids"),
                        mode=execute.get("mode") or "merge",
                        documents=documents,
                    )
                    persist_doc(primary, state, tf)
                    gaps = state.get("gaps") or []
                    reply = result.get("change_summary") or "Dokument oppdatert fra kilder."
                    rem = len(result.get("remaining_gaps") or [])
                    if rem:
                        reply += (
                            f" {rem} mangler står fortsatt åpne."
                            if lang != "en" else
                            f" {rem} gaps remain open."
                        )
                    tool_result = {"tool": "update_document_from_sources", "ok": True, **result}

            elif execute and execute.get("tool") == "list_document_types":
                lang = hub.detect_lang(msg)
                types = dtr.list_document_types(execute.get("industry"))
                lines = [f"- `{t['id']}` — {t['name']}" for t in types]
                reply = (
                    ("Dokumenttyper i registeret:\n" if lang != "en" else "Document types in registry:\n")
                    + "\n".join(lines)
                )
                tool_result = {"tool": "list_document_types", "ok": True, "types": types}

            elif execute and execute.get("tool") == "get_document_type":
                lang = hub.detect_lang(msg)
                tid = execute.get("type_id") or execute.get("id") or ""
                definition = dtr.get_document_type(tid)
                if not definition:
                    reply = (f"Ukjent dokumenttype: `{tid}`." if lang != "en"
                             else f"Unknown document type: `{tid}`.")
                    tool_result = {"tool": "get_document_type", "ok": False}
                else:
                    reply = dtr.format_type_brief(definition)
                    skills = definition.get("skills") or {}
                    if skills.get("primary"):
                        reply += (
                            f"\nNeste: last skill `{skills['primary'][0]}`, deretter "
                            f"`materialise_template` → `compose_document`."
                            if lang != "en" else
                            f"\nNext: load skill `{skills['primary'][0]}`, then "
                            f"`materialise_template` → `compose_document`."
                        )
                    tool_result = {"tool": "get_document_type", "ok": True, "type": definition}

            elif execute and execute.get("tool") == "materialise_template":
                lang = hub.detect_lang(msg)
                tid = execute.get("type_id") or execute.get("id") or ""
                try:
                    mat = dtr.materialise_template(
                        tid,
                        project_id=p.get("id"),
                        overrides=execute.get("overrides") or {},
                        include=execute.get("include") or "required+recommended",
                    )
                except LookupError:
                    reply = (f"Ukjent dokumenttype: `{tid}`." if lang != "en"
                             else f"Unknown document type: `{tid}`.")
                    tool_result = {"tool": "materialise_template", "ok": False}
                else:
                    nsec = len(mat.get("sections") or [])
                    wb = mat.get("workbench_template")
                    reply = (
                        f"Mal materialisert: **{mat.get('name')}** · {nsec} seksjoner"
                        + (f" · workbench `{wb}`" if wb else "")
                        + (". Klar for compose." if lang != "en" else ". Ready to compose.")
                    )
                    tool_result = {"tool": "materialise_template", "ok": True, "template": mat}

            elif execute and execute.get("tool") == "match_document_type":
                lang = hub.detect_lang(msg)
                q = execute.get("query") or msg
                matches = dtr.match_document_types(q, limit=5)
                if not matches:
                    reply = ("Ingen treff i dokumenttyperegisteret — si hvilken type du trenger."
                             if lang != "en" else
                             "No registry match — say which document type you need.")
                    tool_result = {"tool": "match_document_type", "ok": True, "matches": []}
                elif len(matches) == 1:
                    definition = dtr.get_document_type(matches[0]["id"])
                    reply = dtr.format_type_brief(definition or matches[0])
                    tool_result = {
                        "tool": "match_document_type", "ok": True,
                        "matches": matches, "type": definition,
                        "execute_next": {"tool": "materialise_template", "type_id": matches[0]["id"]},
                    }
                else:
                    opts = ", ".join(f"`{m['id']}` ({m['name']})" for m in matches[:4])
                    reply = (
                        f"Flere mulige typer: {opts}. Hvilken mener du?"
                        if lang != "en" else
                        f"Several possible types: {opts}. Which one?"
                    )
                    tool_result = {"tool": "match_document_type", "ok": True, "matches": matches}

            elif execute and execute.get("tool") == "knowledge_index_project":
                try:
                    eng = open_knowledge_engine(primary)
                except ImportError as e:
                    reply = str(e)
                    tool_result = {"tool": "knowledge_index_project", "ok": False}
                else:
                    result = eng.index_project(force_rebuild=bool(execute.get("force_rebuild")))
                    reply = f"Findings-register: {result['rows']} rader · `{Path(result['registry_path']).name}`"
                    tool_result = {"tool": "knowledge_index_project", "ok": True, **result}

            elif execute and execute.get("tool") == "knowledge_get_findings":
                try:
                    eng = open_knowledge_engine(primary)
                except ImportError as e:
                    reply = str(e)
                    tool_result = {"tool": "knowledge_get_findings", "ok": False}
                else:
                    component = execute.get("component")
                    prop = execute.get("property_name")
                    q = execute.get("query") or msg
                    if not component:
                        m = re.search(r"\b(?:om|about|for)\s+([A-Za-z0-9_\-ÅÆØåæø]+)", q, re.I)
                        if m:
                            component = m.group(1)
                    rows = eng.get_findings(component=component, property_name=prop,
                                            source_file=execute.get("source_file"))
                    if not rows and q:
                        rows = eng.semantic_search(q, limit=8)
                    if not rows:
                        reply = "Ingen funn i project_findings.xlsx ennå — indekser kilder og si «importer fakta»."
                    else:
                        lines = []
                        for r in rows[:8]:
                            lines.append(
                                f"- {r.get('component') or '?'} · {r.get('property') or '?'}: "
                                f"**{r.get('value')}** {r.get('unit') or ''} "
                                f"({r.get('citation') or r.get('source_file') or '—'})"
                            )
                        reply = f"{len(rows)} funn" + (f" (viser {len(lines)})" if len(rows) > 8 else "") + ":\n" + "\n".join(lines)
                    tool_result = {"tool": "knowledge_get_findings", "ok": True,
                                   "count": len(rows), "findings": rows[:20]}

            elif execute and execute.get("tool") == "knowledge_update_finding":
                try:
                    eng = open_knowledge_engine(primary)
                except ImportError as e:
                    reply = str(e)
                    tool_result = {"tool": "knowledge_update_finding", "ok": False}
                else:
                    finding = execute.get("finding") or {}
                    if not finding.get("citation") and not finding.get("source_file"):
                        reply = "Trenger citation eller source_file for å lagre funn."
                        tool_result = {"tool": "knowledge_update_finding", "ok": False}
                    else:
                        fid = eng.update_finding(finding)
                        reply = f"Lagret funn `{fid}` i project_findings.xlsx."
                        tool_result = {"tool": "knowledge_update_finding", "ok": True, "finding_id": fid}

            elif execute and execute.get("tool") == "knowledge_semantic_search":
                try:
                    eng = open_knowledge_engine(primary)
                except ImportError as e:
                    reply = str(e)
                    tool_result = {"tool": "knowledge_semantic_search", "ok": False}
                else:
                    rows = eng.semantic_search(execute.get("query") or msg,
                                               limit=int(execute.get("limit") or 10))
                    if not rows:
                        reply = "Ingen semantiske treff i funnregisteret."
                    else:
                        lines = [f"- {r.get('component')} / {r.get('property')}: {r.get('value')} {r.get('unit') or ''}"
                                 for r in rows[:8]]
                        reply = "Treff:\n" + "\n".join(lines)
                    tool_result = {"tool": "knowledge_semantic_search", "ok": True,
                                   "count": len(rows), "findings": rows}

            elif execute and execute.get("tool") == "knowledge_rebuild_index":
                try:
                    eng = open_knowledge_engine(primary)
                except ImportError as e:
                    reply = str(e)
                    tool_result = {"tool": "knowledge_rebuild_index", "ok": False}
                else:
                    result = eng.rebuild_index()
                    reply = f"Vector-indeks bygget på nytt · {result.get('rows_indexed', 0)} rader."
                    tool_result = {"tool": "knowledge_rebuild_index", "ok": True, **result}

            elif execute and execute.get("tool") == "knowledge_import_index_facts":
                try:
                    eng = open_knowledge_engine(primary)
                except ImportError as e:
                    reply = str(e)
                    tool_result = {"tool": "knowledge_import_index_facts", "ok": False}
                else:
                    index = load_index(folders, "no", state.get("user_facts"), project_name=p.get("name"))
                    ids = eng.import_from_index_facts(index)
                    reply = f"Importerte {len(ids)} fakta til project_findings.xlsx."
                    tool_result = {"tool": "knowledge_import_index_facts", "ok": True,
                                   "imported": len(ids), "ids": ids[:50]}

            elif execute and execute.get("tool") == "get_location":
                try:
                    eng = open_knowledge_engine(primary)
                    loc = eng.get_location()
                except ImportError as e:
                    reply = str(e)
                    tool_result = {"tool": "get_location", "ok": False}
                else:
                    if not loc:
                        reply = "Ingen lokasjon i project_findings.xlsx — oppgi adresse."
                    else:
                        reply = (
                            f"Lokasjon: {loc.get('address') or '—'}, "
                            f"{loc.get('municipality') or '—'} "
                            f"({loc.get('latitude')}, {loc.get('longitude')})"
                            + (f" · kart: `{loc.get('map_image_path')}`" if loc.get("map_image_path") else "")
                        )
                    tool_result = {"tool": "get_location", "ok": True, "location": loc}

            elif execute and execute.get("tool") == "set_location":
                try:
                    eng = open_knowledge_engine(primary)
                    loc = eng.set_location(
                        address=execute.get("address") or "",
                        municipality=execute.get("municipality"),
                        postal_code=execute.get("postal_code"),
                        latitude=execute.get("latitude"),
                        longitude=execute.get("longitude"),
                        location_type=execute.get("location_type") or "project_site",
                        citation=execute.get("citation"),
                        map_style=execute.get("map_style"),
                    )
                    reply = f"Lagret lokasjon: {loc.get('address')} · {loc.get('municipality') or ''}"
                    tool_result = {"tool": "set_location", "ok": True, "location": loc}
                except Exception as e:
                    reply = str(e)
                    tool_result = {"tool": "set_location", "ok": False}

            elif execute and execute.get("tool") == "generate_location_map":
                try:
                    eng = open_knowledge_engine(primary)
                    rel = eng.generate_location_map(
                        style=execute.get("style") or "default",
                        width=int(execute.get("width") or 1200),
                        height=int(execute.get("height") or 800),
                        zoom=int(execute.get("zoom") or 16),
                        color_overrides=execute.get("color_overrides"),
                        output_format=execute.get("output_format") or "png",
                    )
                    reply = f"Kart generert: `{rel}` (ligger i prosjektmappen)."
                    tool_result = {"tool": "generate_location_map", "ok": True,
                                   "map_image_path": rel}
                except Exception as e:
                    reply = str(e)
                    tool_result = {"tool": "generate_location_map", "ok": False}

            elif execute and execute.get("tool") == "propose_location_map":
                try:
                    eng = open_knowledge_engine(primary)
                    proposal = eng.propose_location_map(
                        style=execute.get("style") or "default",
                        width=int(execute.get("width") or 1200),
                        height=int(execute.get("height") or 800),
                        zoom=int(execute.get("zoom") or 16),
                        color_overrides=execute.get("color_overrides"),
                        output_format=execute.get("output_format") or "png",
                        caption=execute.get("caption"),
                    )
                    reply = (
                        f"{proposal.get('message')} "
                        f"Fil: `{proposal['image']['path']}`. Si **ja** for å bekrefte innsetting."
                    )
                    state["chat_pending"] = {
                        "action": "confirm_location_map",
                        "proposal": proposal,
                    }
                    tool_result = {
                        "tool": "propose_location_map", "ok": True, **proposal,
                        "actions": [{"id": "confirm_generate", "label": "Ja — sett inn kart"}],
                    }
                except Exception as e:
                    reply = str(e)
                    tool_result = {"tool": "propose_location_map", "ok": False}

            elif execute and execute.get("tool") == "confirm_location_map":
                prop = execute.get("proposal") or (state.get("chat_pending") or {}).get("proposal") or {}
                img = (prop.get("image") or {})
                path = img.get("path") or ""
                if path:
                    reply = (
                        f"Kart bekreftet: `{path}`. Sett inn ImageBlock i aktuell seksjon "
                        f"(bildetekst: {img.get('caption') or '—'})."
                    )
                    tool_result = {
                        "tool": "confirm_location_map", "ok": True,
                        "image_path": path, "proposal": prop,
                    }
                else:
                    reply = "Ingen kart-proposal å bekrefte."
                    tool_result = {"tool": "confirm_location_map", "ok": False}
                state["chat_pending"] = None

            elif execute and execute.get("tool") == "ground_photo":
                import agent_truth as atruth
                index = load_index(folders, "no", state.get("user_facts"), project_name=p.get("name"))
                lang = hub.detect_lang(msg)
                last = state.get("last_indexed_media") or {}
                entry = None
                if last.get("file"):
                    entry = next((e for e in index if e.get("file") == last["file"]), None)
                    if not entry:
                        entry = {"file": last.get("file"), "caption": last.get("caption"),
                                 "facts": [], "kind": "photo"}
                if not entry:
                    photos = atruth.photo_entries(index)
                    entry = photos[-1] if photos else None
                bom_hypos = [c.get("part_no") for c in (state.get("bom_components") or [])
                             if c.get("part_no")]
                grounded = atruth.ground_photo_reply(entry, bom_hypotheses=bom_hypos, lang=lang)
                reply = grounded.get("reply") or ""
                if grounded.get("set_pending") is not None:
                    state["chat_pending"] = grounded["set_pending"]
                nested = grounded.get("execute")
                tool_result = {"tool": "ground_photo", "ok": True}
                if nested and nested.get("tool") == "add_bom_component":
                    add_bom_component(
                        state,
                        part_no=nested.get("part_no"),
                        file=nested.get("file"),
                        caption=nested.get("caption"),
                        confidence=nested.get("confidence") or 0,
                        fact_id=nested.get("fact_id"),
                        status=nested.get("status") or "ok",
                        verified_by_user=bool(nested.get("verified_by_user")),
                    )
                    if tf:
                        refresh_bom_section(state, folders, tf, lang)
                        persist_doc(primary, state, tf)
                    else:
                        save_state(primary, state)
                    tool_result = {"tool": "add_bom_component", "ok": True,
                                   "part_no": nested.get("part_no"), "file": nested.get("file")}
                if grounded.get("actions"):
                    tool_result["actions"] = grounded["actions"]
                ok_p, reply, _ = atruth.validate_perception(
                    reply, [entry] if entry else [], lang=lang)
                if not ok_p:
                    tool_result["perception_scrub"] = True

            elif execute and execute.get("tool") == "add_bom_component":
                import agent_truth as atruth
                lang = hub.detect_lang(msg)
                row = add_bom_component(
                    state,
                    part_no=execute.get("part_no"),
                    file=execute.get("file"),
                    caption=execute.get("caption"),
                    confidence=execute.get("confidence") or 1.0,
                    fact_id=execute.get("fact_id"),
                    status=execute.get("status") or "ok",
                    verified_by_user=bool(execute.get("verified_by_user")),
                )
                if tf:
                    refresh_bom_section(state, folders, tf, lang)
                    persist_doc(primary, state, tf)
                else:
                    save_state(primary, state)
                pn = row.get("part_no") or "?"
                reply = (f"Lagt inn **{pn}** i BOM med bildereferanse ✓ — Indeksert som: "
                         f"{row.get('caption') or Path(row.get('file') or '').name}.")
                tool_result = {"tool": "add_bom_component", "ok": True, "part_no": pn}

            elif execute and execute.get("tool") == "run_generate":
                # WORKORDER_0.25 B — «ja» / imperative dispatches generate in this turn
                lang = hub.detect_lang(msg)
                tkey = execute.get("template_key")
                gen_tf = execute.get("template") or template_file_for_key(tkey) or tf
                if not gen_tf:
                    reply = ("Ukjent dokumentmal — åpne Contract Review først."
                             if lang != "en" else
                             "Unknown template — open Contract Review first.")
                    tool_result = {"tool": "run_generate", "ok": False}
                elif not state.get("confirmed"):
                    reply = ("Bekreft artefaktmodellen først (steg 2), så starter jeg genereringen."
                             if lang != "en" else
                             "Confirm the artifact model first (step 2), then I'll start generation.")
                    tool_result = {"tool": "run_generate", "ok": False}
                else:
                    job_id = start_job(run_generate, folders, gen_tf, lang)
                    state["chat_pending"] = None
                    # System event joins conversation (A1)
                    edchat.append_turn(
                        state, "system",
                        f"[job_started] generate job_id={job_id} template={gen_tf}",
                        meta={"kind": "job_started", "job_id": job_id},
                        project_id=p.get("id"))
                    reply = (
                        f"Starter Contract Review — jobb `{job_id}` kjører nå."
                        if lang != "en" else
                        f"Starting Contract Review — job `{job_id}` is running."
                    )
                    tool_result = {
                        "tool": "run_generate", "ok": True,
                        "job_id": job_id, "template": gen_tf,
                    }

            elif execute and execute.get("tool") == "propose_connection_spec":
                import connection_diagram as cdiag
                import sys
                if str(ROOT) not in sys.path:
                    sys.path.insert(0, str(ROOT))
                lang = hub.detect_lang(msg)
                index = load_index(folders, "no", state.get("user_facts"),
                                   project_name=p.get("name"))
                # Circuit schematic ask → honest boundary, still propose block diagram
                preamble = ""
                if cdiag.is_circuit_schematic_ask(msg):
                    preamble = cdiag.circuit_boundary_reply(lang) + "\n\n"
                ask = None
                if KEY_SET:
                    try:
                        ask = fc.ask
                    except Exception:
                        ask = None
                spec = cdiag.propose_connection_spec(
                    artifact=state.get("artifact"),
                    bom_components=state.get("bom_components"),
                    index=index,
                    ask_fn=ask,
                    lang=lang,
                )
                # Process / funksjonsdiagram with no electronics → renseanlegg fixture
                if (cdiag.is_process_diagram_ask(msg)
                        and len(spec.get("components") or []) < 2):
                    spec = cdiag.process_fixture_spec(lang=lang)
                # Enrich excavator-style demos: if only Pi + excavator, add std companions
                labels = " ".join(c.get("label", "") for c in spec.get("components") or [])
                if re.search(r"raspberry\s*pi", labels, re.I) and len(spec.get("components") or []) < 4:
                    extra_names = ["PCA9685", "5V buck converter", "Battery / supply"]
                    existing = {c["id"] for c in spec["components"]}
                    for name in extra_names:
                        std = cdiag.match_standard(name)
                        if std and std["id"] not in existing:
                            spec["components"].append({
                                "id": std["id"], "label": std["label"],
                                "pins": list(std["pins"]), "role": std["role"],
                                "fact_id": None, "image": None,
                            })
                            existing.add(std["id"])
                    spec = cdiag.propose_connection_spec(
                        components=spec["components"], lang=lang)
                reply = preamble + cdiag.format_confirm_table(spec, lang=lang)
                state["chat_pending"] = {
                    "action": "confirm_connection_spec",
                    "spec": spec,
                }
                state["connection_spec_proposed"] = spec
                save_state(primary, state)
                tool_result = {
                    "tool": "propose_connection_spec", "ok": True,
                    "edges": len(spec.get("connections") or []),
                    "components": len(spec.get("components") or []),
                    "cost_eur": cdiag.DEMO_EUR if spec.get("model_enriched") else 0,
                    "actions": [{
                        "id": "confirm_connection_all",
                        "label": "Confirm all" if lang == "en" else "Bekreft alle",
                    }],
                }

            elif execute and execute.get("tool") in (
                    "confirm_connection_spec", "create_diagram"):
                import connection_diagram as cdiag
                import sys
                if str(ROOT) not in sys.path:
                    sys.path.insert(0, str(ROOT))
                lang = hub.detect_lang(msg)
                pending_spec = (state.get("chat_pending") or {}).get("spec") \
                    or state.get("connection_spec_proposed") \
                    or execute.get("spec")
                if not pending_spec:
                    reply = ("No proposed connection spec — ask for a block diagram first."
                             if lang == "en" else
                             "Ingen foreslått tilkoblingsspesifikasjon — be om blokkskjema først.")
                    tool_result = {"tool": "create_diagram", "ok": False}
                else:
                    confirmed = cdiag.apply_edge_decisions(
                        pending_spec,
                        accept_all=bool(execute.get("accept_all")
                                        or execute.get("tool") == "create_diagram"),
                        drop_rows=execute.get("drop_rows"),
                        keep_rows=execute.get("keep_rows"),
                    )
                    title = None
                    if confirmed.get("kind") == "process_flow":
                        title = "Funksjonsdiagram" if lang != "en" else "Process flow"
                    svg = cdiag.render_block_diagram(confirmed, title=title)
                    section = execute.get("section") or "connection_diagram"
                    store_connection_diagram(state, confirmed, svg, lang=lang)
                    try:
                        media = reports_dir(primary) / "media"
                        media.mkdir(parents=True, exist_ok=True)
                        (media / "block_diagram.svg").write_text(svg, encoding="utf-8")
                    except Exception:
                        pass
                    if tf:
                        persist_doc(primary, state, tf)
                    else:
                        save_state(primary, state)
                    reply = cdiag.diagram_created_reply(
                        confirmed, section=section, lang=lang)
                    tool_result = {
                        "tool": "create_diagram", "ok": True,
                        "edges": len(confirmed.get("connections") or []),
                        "components": len(confirmed.get("components") or []),
                        "section": section,
                        "svg_bytes": len(svg.encode("utf-8")),
                    }

            elif execute and execute.get("tool") == "write_checklist":
                import agent_truth as atruth
                lang = hub.detect_lang(msg)
                caps = hub.load_capabilities()
                tkey = (template or {}).get("template_key") if template else None
                if not tkey and tf:
                    tkey = Path(tf).stem
                caps_entry = next(
                    (t for t in (caps.get("templates") or [])
                     if t.get("key") == tkey or t.get("file") == tf),
                    None,
                )
                hub.write_checklist(Path(primary), caps_entry)
                path = Path(primary) / "SJEKKLISTE.txt"
                n_items = 0
                if caps_entry and caps_entry.get("checklist"):
                    n_items = len(caps_entry["checklist"])
                reply = atruth.checklist_created_reply(
                    str(path), n_items=n_items, lang=lang)
                tool_result = {
                    "tool": "write_checklist", "ok": True,
                    "path": str(path), "items": n_items,
                }

            elif execute and execute.get("tool") == "recreate_form":
                # TEMPLATE_STANDARD — default flexible inspection_checklist;
                # sample_multipoint only when source explicitly names the fixture
                import form_model as fm
                import form_engine as fe
                import template_lifecycle as tl
                lang = hub.detect_lang(msg)
                source = (execute.get("source") or "inspection_checklist").strip()
                if source == "sample_multipoint":
                    tpl_obj = fe.fixture_as_template()
                    tpl_obj = fm.validate_form_template(tpl_obj) if hasattr(fm, "validate_form_template") else tpl_obj
                    fname = save_company_template(tpl_obj)
                    try:
                        regen_capabilities()
                    except Exception:
                        pass
                    # Keep canonical JSON under fixtures/ only — never ship in templates/
                    fpath = ROOT / "fixtures" / "sample_multipoint" / "sample_multipoint.json"
                    if not fpath.exists():
                        fpath.parent.mkdir(parents=True, exist_ok=True)
                        fpath.write_text(json.dumps(tpl_obj, indent=2, ensure_ascii=False),
                                         encoding="utf-8")
                    tf = fname  # owned company template file
                    loaded = load_template(tf, primary) or tpl_obj
                else:
                    tf = "inspection_checklist.json"
                    loaded = load_template(tf, primary)
                    if not loaded:
                        raise ValueError("inspection_checklist template missing")
                    fname = tf
                index = load_index(folders, lang, state.get("user_facts"),
                                   project_name=p.get("name"))
                tl.create_document_shell(state, tf, loaded)
                pref = fm.prefill_form(state, loaded, index)
                state["gaps"] = ds.gaps_for_document(
                    state, loaded, index, state.get("artifact") or {}, fc)
                content = ds.assemble_draft(state, loaded, state.get("artifact"))
                sync_draft_files(primary, state, loaded, tf, content)
                state["chat_pending"] = None
                state["active_template"] = tf
                html_path = None
                try:
                    html_path = write_form_html_export(
                        primary, loaded, state, report_display_name(loaded))
                except Exception:
                    pass
                n_fields = sum(len(s.get("fields") or []) for s in loaded.get("sections") or [])
                reply = (
                    f"Opprettet mal **{loaded.get('name_no') or loaded.get('name')}** "
                    f"({n_fields} felt) og dokument i prosjektet"
                    + (f" — {pref.get('prefilled', 0)} felt forhåndsutfylt" if pref.get("prefilled") else "")
                    + (f". HTML: `{Path(str(html_path)).name}`" if html_path else "")
                    + "."
                    if lang != "en" else
                    f"Created template **{loaded.get('name')}** ({n_fields} fields) "
                    f"and a document in the project."
                )
                save_state(primary, state)
                template = loaded
                tool_result = {
                    "tool": "recreate_form", "ok": True,
                    "template": tf,
                    "template_key": loaded.get("template_key") or Path(tf).stem,
                    "document_species": "form_fill",
                    "fields": n_fields,
                    "prefilled": pref.get("prefilled", 0),
                    "export_html": str(html_path) if html_path else None,
                    "model_calls": 0,
                    "cost_eur": 0,
                }

            elif execute and execute.get("tool") == "propose_form_template":
                # WORKORDER_0.30 — chat: «create service skjema.jpg as a template»
                import form_model as fm
                import chat_attach as chattach
                lang = hub.detect_lang(msg)
                fname = execute.get("file") or chattach.mentioned_filename(msg)
                found = chattach.find_project_file(folders, fname) if fname else None
                if not found:
                    # Prefer form-named images already in project
                    for folder in folders:
                        for p in Path(folder).rglob("*"):
                            if p.is_file() and chattach.FORM_NAME.search(p.name):
                                found = p
                                break
                        if found:
                            break
                if not found:
                    reply = (
                        "Fant ikke skjemafilen — last opp eller oppgi filnavnet "
                        "(f.eks. skjema.jpg)."
                        if lang != "en" else
                        "Couldn't find that form file — upload it or name the file "
                        "(e.g. skjema.jpg)."
                    )
                    tool_result = {"tool": "propose_form_template", "ok": False}
                else:
                    raw = found.read_bytes()
                    name = found.name
                    ATTACH_STAGING.mkdir(parents=True, exist_ok=True)
                    token = uuid.uuid4().hex[:12]
                    (ATTACH_STAGING / f"{token}_{name}").write_bytes(raw)
                    peek = chattach.peek_text_bytes(raw, name)
                    det = fm.detect_form_shaped(peek, name)
                    if not det.get("form_shaped") and chattach.FORM_NAME.search(name):
                        det = {**det, "form_shaped": True}
                    draft = fm.offline_extract_form_structure(peek, name)
                    summary = fm.form_summary_for_offer(draft)
                    if not peek and Path(name).suffix.lower() in chattach.PHOTO_EXT:
                        summary = (
                            f"bildet «{name}» (skjema)"
                            if lang != "en" else
                            f"image «{name}» (form)"
                        )
                    offer = fm.form_propose_reply(
                        summary, filled=det.get("filled"), lang=lang)
                    reply = offer["reply"]
                    acts = [
                        {**a, "token": token}
                        for a in offer["actions"]
                    ] + [
                        {"id": "as_project", "label": "Prosjektfil"
                         if lang != "en" else "Project file", "token": token}
                    ]
                    tool_result = {
                        "tool": "propose_form_template", "ok": True,
                        "kind": "form_propose",
                        "token": token, "name": name,
                        "summary": summary,
                        "actions": acts,
                    }

            elif execute and execute.get("tool") == "create_document":
                import template_lifecycle as tl
                import form_model as fm
                lang = hub.detect_lang(msg)
                caps = hub.load_capabilities()
                tkey = execute.get("template_key") or "installation_manual"
                gen_tf = execute.get("template") or template_file_for_key(tkey)
                tpl_obj = load_template(gen_tf, primary) if gen_tf else None
                if not gen_tf or not tpl_obj:
                    reply = ("Mal ikke funnet." if lang != "en" else "Template not found.")
                    tool_result = {"tool": "create_document", "ok": False}
                elif fm.is_form_fill(tpl_obj):
                    # WORKORDER_0.29 — create + prefill, zero Sonnet, no € generate ask
                    index = load_index(folders, lang, state.get("user_facts"),
                                       project_name=p.get("name"))
                    tl.create_document_shell(state, gen_tf, tpl_obj)
                    pref = fm.prefill_form(state, tpl_obj, index)
                    state["gaps"] = ds.gaps_for_document(
                        state, tpl_obj, index, state.get("artifact") or {}, fc)
                    content = ds.assemble_draft(state, tpl_obj, state.get("artifact"))
                    sync_draft_files(primary, state, tpl_obj, gen_tf, content)
                    state["chat_pending"] = None
                    reply = tl.document_created_reply(tpl_obj, lang=lang)
                    if pref.get("prefilled"):
                        reply += (f" ({pref['prefilled']} felt forhåndsutfylt)."
                                  if lang != "en" else
                                  f" ({pref['prefilled']} fields prefilled).")
                    save_state(primary, state)
                    tf = gen_tf
                    template = tpl_obj
                    tool_result = {
                        "tool": "create_document", "ok": True,
                        "template": gen_tf, "template_key": tkey,
                        "document_species": "form_fill",
                        "prefilled": pref.get("prefilled", 0),
                        "model_calls": 0,
                    }
                else:
                    tl.create_document_shell(state, gen_tf, tpl_obj)
                    tier = tl.tier_eur_for_template(tpl_obj, caps)
                    reply = tl.document_created_reply(tpl_obj, lang=lang, tier_eur=tier)
                    state["chat_pending"] = {
                        "action": "run_generate",
                        "template_key": tkey,
                        "template": gen_tf,
                    }
                    save_state(primary, state)
                    tf = gen_tf
                    template = tpl_obj
                    tool_result = {
                        "tool": "create_document", "ok": True,
                        "template": gen_tf, "template_key": tkey,
                        "actions": [{"id": "confirm_generate", "label": "Ja — generer"}],
                    }

            elif execute and execute.get("tool") == "draft_template_rung3":
                import template_lifecycle as tl
                import telemetry as tele
                lang = hub.detect_lang(msg)
                story = execute.get("story") or msg
                art = state.get("artifact") or {"name": p.get("name"), "purpose": story}
                chat_block = build_project_chat_context(
                    p, folders, primary, state, lang=lang)["text"]
                if KEY_SET:
                    try:
                        drafted = fc.draft_template(story, art, lang, project_context=chat_block)
                    except Exception as e:
                        reply = f"Kunne ikke utforme mal: {e}"
                        tool_result = {"tool": "draft_template_rung3", "ok": False}
                        drafted = None
                else:
                    drafted = tl.offline_stub_commissioning_template(story, lang)
                if drafted:
                    tele.log_rung3_request(drafted)
                    state["chat_pending"] = {
                        "action": "accept_drafted_template",
                        "draft": drafted,
                    }
                    reply = tl.format_draft_structure_card(drafted, lang=lang)
                    tool_result = {
                        "tool": "draft_template_rung3", "ok": True,
                        "template_key": drafted.get("template_key"),
                        "actions": tl.accept_draft_actions(lang),
                    }

            elif execute and execute.get("tool") == "accept_drafted_template":
                import template_lifecycle as tl
                lang = hub.detect_lang(msg)
                caps = hub.load_capabilities()
                drafted = execute.get("draft") or (state.get("chat_pending") or {}).get("draft")
                if not drafted:
                    reply = ("Ingen utkast å godta." if lang != "en" else "No draft to accept.")
                    tool_result = {"tool": "accept_drafted_template", "ok": False}
                else:
                    fname = save_drafted_template(primary, drafted)
                    tpl_obj = load_template(fname, primary)
                    tl.create_document_shell(state, fname, tpl_obj)
                    tier = tl.tier_eur_for_template(tpl_obj, caps)
                    name = tpl_obj.get("name_no") if lang == "no" else tpl_obj.get("name")
                    reply = (
                        f"Lagret **{name}** som egen mal ({tpl_obj.get('badge')}) og opprettet dokumentet. "
                        f"Generering ~€{tier} — si **ja** for å starte."
                        if lang != "en" else
                        f"Saved **{name}** to your library ({tpl_obj.get('badge')}) and created the document. "
                        f"Generation ~€{tier} — say **yes** to start."
                    )
                    state["chat_pending"] = {
                        "action": "run_generate",
                        "template_key": tpl_obj.get("template_key"),
                        "template": fname,
                    }
                    tf = fname
                    template = tpl_obj
                    save_state(primary, state)
                    tool_result = {
                        "tool": "accept_drafted_template", "ok": True,
                        "template": fname, "badge": tpl_obj.get("badge"),
                        "actions": [{"id": "confirm_generate", "label": "Ja — generer"}],
                    }

            elif execute and execute.get("tool") in (
                    "move_section", "toggle_section", "set_block_layout", "add_section"):
                import doc_structure as dstruct
                lang = hub.detect_lang(msg)
                tool = execute.get("tool")
                try:
                    if tool == "move_section":
                        result = dstruct.move_section(
                            state, template, execute["key"],
                            position=execute.get("position"),
                            after=execute.get("after"),
                            before=execute.get("before"),
                        )
                        reply = dstruct.format_move_reply(execute["key"], template, state, lang=lang)
                    elif tool == "toggle_section":
                        result = dstruct.toggle_section(
                            state, template, execute["key"],
                            enabled=bool(execute.get("enabled", True)),
                        )
                        reply = dstruct.format_move_reply(execute["key"], template, state, lang=lang)
                        if result.get("warning") == "traceability_reduced":
                            reply += dstruct.toggle_warning_reply(lang=lang)
                    elif tool == "set_block_layout":
                        result = dstruct.set_block_layout(
                            state, template, execute["key"], execute.get("layout") or "table")
                        reply = (
                            f"Layout **{result['layout']}** på **{execute['key']}** — logget, ingen regenerering."
                            if lang != "en" else
                            f"Layout **{result['layout']}** on **{execute['key']}** — logged, no regeneration."
                        )
                    else:
                        result = dstruct.add_section(
                            state, template,
                            execute.get("title") or "Ny seksjon",
                            after=execute.get("after"),
                            rules=execute.get("rules"),
                        )
                        reply = (
                            f"La til seksjon **{result['title']}** — logget, ingen regenerering."
                            if lang != "en" else
                            f"Added section **{result['title']}** — logged, no regeneration."
                        )
                    if tf:
                        persist_doc(primary, state, tf)
                    else:
                        save_state(primary, state)
                    tool_result = {"tool": tool, "ok": True, **result}
                    if dstruct.maybe_save_template_offer(state):
                        reply += " " + dstruct.save_template_offer_reply(lang=lang)
                        tool_result["actions"] = [
                            {"id": "save_as_template", "label": "Lagre som mal"}]
                        state["chat_pending"] = {"action": "save_as_template"}
                except ValueError as e:
                    reply = str(e)
                    tool_result = {"tool": tool, "ok": False}

            elif execute and execute.get("tool") == "save_as_template":
                import doc_structure as dstruct
                lang = hub.detect_lang(msg)
                if not template:
                    reply = ("Åpne et dokument først." if lang != "en" else "Open a document first.")
                    tool_result = {"tool": "save_as_template", "ok": False}
                else:
                    owned = dstruct.export_user_template(state, template, story=msg)
                    fname = save_company_template(owned)
                    regen_capabilities()
                    name = owned.get("name_no") if lang == "no" else owned.get("name")
                    reply = (
                        f"Lagret **{name}** som egen mal ({owned.get('badge')})."
                        if lang == "no" else
                        f"Saved **{name}** as your template ({owned.get('badge')})."
                    )
                    tool_result = {"tool": "save_as_template", "ok": True, "file": fname}

            elif execute and execute.get("tool") == "scan_components_offer":
                import agent_truth as atruth
                index = load_index(folders, "no", state.get("user_facts"), project_name=p.get("name"))
                lang = hub.detect_lang(msg)
                offer = atruth.scan_offer_reply(index, lang=lang)
                reply = offer.get("reply") or ""
                if offer.get("set_pending") is not None:
                    state["chat_pending"] = offer["set_pending"]
                if offer.get("execute", {}).get("tool") == "refresh_bom_from_components":
                    if tf:
                        refresh_bom_section(state, folders, tf, lang)
                        persist_doc(primary, state, tf)
                    tool_result = {"tool": "refresh_bom_from_components", "ok": True}
                else:
                    tool_result = {"tool": "scan_components_offer", "ok": True,
                                   "pending": offer.get("pending"),
                                   "estimate_eur": offer.get("estimate_eur")}
                # Surface Skann action for UI if present
                if offer.get("actions"):
                    tool_result["actions"] = offer["actions"]

            elif execute and execute.get("tool") == "scan_components_run":
                import agent_truth as atruth
                index = load_index(folders, "no", state.get("user_facts"), project_name=p.get("name"))
                lang = hub.detect_lang(msg)
                results = run_component_scan(folders, index, state, lang=lang)
                if tf:
                    refresh_bom_section(state, folders, tf, lang)
                    persist_doc(primary, state, tf)
                else:
                    save_state(primary, state)
                reply = atruth.scan_complete_reply(results, lang=lang)
                tool_result = {"tool": "scan_components_run", "ok": True, **results}

            elif execute and execute.get("tool") == "set_cover":
                # WORKORDER_0.20 B — execute forside; zero permission-seeking
                index = load_index(folders, "no", state.get("user_facts"), project_name=p.get("name"))
                rel, cap = edchat.pick_cover_image(state, index, msg)
                lang = hub.detect_lang(msg)
                if not rel:
                    reply = ("Jeg finner ikke et bilde å bruke — legg til et bilde først, "
                             "så setter jeg det på forsiden.") if lang == "no" else (
                        "I can't find an image to use — add one first, then I'll set it as cover.")
                    tool_result = {"tool": "set_cover", "ok": False}
                elif not template:
                    reply = ("Åpne et dokument først, så setter jeg bildet på forsiden.") if lang == "no" else (
                        "Open a document first, then I'll set the image as cover.")
                    tool_result = {"tool": "set_cover", "ok": False}
                else:
                    result = apply_cover_image(state, folders, template, rel, cap)
                    if tf:
                        persist_doc(primary, state, tf)
                    else:
                        save_state(primary, state)
                    doc_name = (template.get("name_no") if lang == "no" else template.get("name")) \
                        or template.get("name_no") or template.get("template_key") or "dokumentet"
                    # Optional follow-up if other drafted docs exist
                    other = []
                    for d in state.get("documents") or []:
                        if d.get("template") and d.get("template") != tf:
                            other.append(d.get("name_no") or d.get("name") or d.get("template"))
                    reply = edchat.format_cover_reply(
                        doc_name=doc_name, caption=result.get("caption") or cap,
                        rel=result.get("file") or rel, other_docs=other[:1], lang=lang)
                    tool_result = {"tool": "set_cover", "ok": True, "file": result.get("file"),
                                   "section": result.get("section")}

            elif execute and execute.get("tool") == "resolve_mangler":
                key = execute.get("key")
                value = execute.get("value")
                unit = execute.get("unit")
                index = load_index(folders, "no", state.get("user_facts"), project_name=p.get("name"))
                result = ds.resolve_mangler(state, key, value, unit, template, index,
                                            state.get("artifact"), fc, provenance="user",
                                            section=execute.get("section"))
                if tf:
                    persist_doc(primary, state, tf)
                else:
                    save_state(primary, state)
                gs = ds.gaps_summary(result["gaps"])
                reply = (f"Lagt inn **{value}** ✓ — merket som oppgitt av deg. "
                         f"{gs.get('blocking', 0)} blokkerende igjen.")
                tool_result = {"tool": "resolve_mangler", "key": key, "gap_summary": gs}
                state["chat_pending"] = None

            elif execute and execute.get("tool") == "suggest_reference":
                key = execute.get("key")
                if not fc.allows_reference_suggest(key):
                    reply = f"Referanseverdi tillates ikke for `{key}` (compliance/sikkerhet)."
                elif not KEY_SET:
                    reply = "API-nøkkel mangler for referanseforslag."
                else:
                    index = load_index(folders, "no", state.get("user_facts"), project_name=p.get("name"))
                    chat_block = build_project_chat_context(
                        p, folders, primary, state, index)["text"]
                    sug = fc.reference_suggest(key, state.get("artifact") or {"name": p.get("name")},
                                               project_context=chat_block)
                    if not sug:
                        reply = f"Jeg er ikke trygg nok på en typisk verdi for `{key}`."
                    else:
                        u = f" {sug['unit']}" if sug.get("unit") else ""
                        reply = (f"Foreslått referanse for `{key}`: **{sug['value']}{u}** — {sug.get('basis') or ''}\n"
                                 f"(AI-foreslått — ikke verifisert. Si «bruk referansen» for å merke som ubekreftet.)")
                        state["chat_pending"] = {"action": "accept_reference", "key": key,
                                                 "value": sug["value"], "unit": sug.get("unit")}
                tool_result = {"tool": "suggest_reference", "key": key}

            elif execute and execute.get("tool") == "regenerate_section":
                sec = execute.get("section")
                # Resolve friendly names against current doc sections
                doc_secs = ((state.get("doc") or {}).get("sections") or {})
                if sec and sec not in doc_secs:
                    aliases = {
                        "summary": ("summary", "executive_summary", "sammendrag"),
                        "identification": ("identification", "identifikasjon", "id"),
                    }
                    for cand in aliases.get(sec, (sec,)):
                        if cand in doc_secs:
                            sec = cand
                            break
                    else:
                        # fuzzy: any section key/title containing the word
                        needle = (sec or "").lower()
                        for sk in doc_secs:
                            if needle and needle in sk.lower():
                                sec = sk
                                break
                if route.get("need_section") or not sec:
                    reply = "Velg en seksjon i dokumentet først (klikk på den), så regenererer jeg med forhåndsvisning."
                else:
                    reply = (f"Jeg foreslår å regenerere seksjonen **{sec}** "
                             f"med instruksen din. Bekreft i diff-kortet (Godta) — jeg endrer ikke teksten alene.")
                    tool_result = {"tool": "regenerate_section", "section": sec,
                                   "instruction": execute.get("instruction"), "needs_confirm": True}

            elif route.get("kind") == "ask_value" or route.get("kind") == "clarify":
                reply = route.get("reply") or reply

            elif route.get("kind") == "confirm_part_no":
                reply = route.get("reply") or reply

            elif route.get("need_model"):
                # Same agent — full project chat context + §7 policy
                index = load_index(folders, "no", state.get("user_facts"), project_name=p.get("name"))
                ctx_pack = build_project_chat_context(
                    p, folders, primary, state, index, lang=body.get("lang", "no"))
                chat_block = ctx_pack["text"]
                extras = chat_turn_extras(msg, index, state.get("artifact") or {}, ctx_pack["file_count"])
                chat_extras = extras
                lang = hub.detect_lang(msg)
                annot_ctx = (route.get("_annot_ctx")
                             or edchat.format_annotations_context(annotations))
                # WORKORDER 0.56 §C5 — no open-ended create-offer while marks pending
                if annotations and not edchat.explicitly_names_document(
                        msg, hub.load_capabilities()):
                    extras = dict(extras)
                    extras["open_ended"] = False
                    chat_extras = extras
                # Attach last photo caption for perception discipline
                last = state.get("last_indexed_media") or {}
                photo_hint = ""
                if last.get("file"):
                    photo_hint = (
                        f"\nLAST INDEXED MEDIA: file={last.get('file')} "
                        f"caption={last.get('caption') or '(none)'} — "
                        f"quote only this; never invent visual details.\n"
                    )
                if extras["open_ended"] and (ctx_pack["file_count"] or 0) > 0:
                    reply = edchat.open_ended_grounded_reply(
                        project_name=p.get("name") or "",
                        brief=extras["corpus_brief"],
                        artifact=state.get("artifact") or {},
                        known_block=extras["known_block"],
                        estimate_eur=extras["estimate_eur"],
                        lang=lang,
                    )
                    reply = edchat.scrub_chat_voice(reply)
                else:
                    offer_line = ""
                    if extras["open_ended"]:
                        offer_line = (
                            f"\nOpen-ended create ask: end with document offer (~€{extras['estimate_eur']:.2f}). "
                            f"At most two questions; never ask what the index already knows.\n"
                        )
                    annot_line = f"\n{annot_ctx}\n" if annot_ctx else ""
                    if KEY_SET:
                        prompt = (
                            f"{chat_block}\n\n"
                            f"{extras['known_block']}\n\n"
                            f"{extras['policy']}\n"
                            f"{photo_hint}"
                            f"{annot_line}"
                            f"{offer_line}\n"
                            f"You are Foldok's ONE project assistant (same agent as checkpoint A).\n"
                            f"Scope focus (UI chip only — does NOT change identity): {scope or 'document'}.\n"
                            f"Ground in sentence 1. Never invent numbers. NEVER use other projects.\n"
                            f"Never claim file writes without tools. Never invent part numbers.\n"
                            f"Deictics (denne/her/dette) resolve to PENDING ANNOTATIONS when present — "
                            f"never invent or create a document unless the user explicitly names one.\n\n"
                            f"USER: {msg}\n\n"
                            f"Do NOT output a feature menu. Do NOT list capabilities unless asked."
                        )
                        raw = fc.ask("chat_edit", fc.HAIKU, [{"role": "user", "content": prompt}], max_tokens=500)
                        reply = edchat.scrub_chat_voice((raw or "").strip())
                        if edchat.reply_violates_policy(reply):
                            raw2 = fc.ask("chat_edit", fc.HAIKU, [{"role": "user", "content":
                                prompt + "\n\nRewrite without emoji / 'helt nytt' / Kult! / fictional writes."
                            }], max_tokens=400)
                            reply = edchat.scrub_chat_voice((raw2 or "").strip())
                    else:
                        reply = ("Jeg forsto ikke helt — prøv «hva mangler?», beskriv mangelen "
                                 "(f.eks. registreringsnummer), eller velg en seksjon og si «skriv om strengere».")


            if not reply:
                reply = "OK."
            # WORKORDER_0.25 B — model propose «Skal jeg …?» → store pending (no fiction)
            if (not execute and route.get("kind") not in ("propose_generate", "dispatch_generate",
                                                          "dispatch_pending", "recreate_form",
                                                          "annot_blocks_create", "annot_execute")
                    and not annotations  # C5: never invent create-pending while marks wait
                    and re.search(r"\bskal\s+jeg\b|\bshall\s+i\b|\bwant\s+me\s+to\b",
                                  reply or "", re.I)):
                # Form / mal proposals must never land as .txt checklist
                recent = " ".join(
                    (t.get("text") or "") for t in (state.get("conversation") or [])[-8:])
                formish = bool(re.search(
                    r"skjema|form|multipoint|mal\b|template|inspeksjon",
                    recent + " " + (reply or ""), re.I))
                if formish or re.search(r"skjema|form|mal|multipoint", reply or "", re.I):
                    src = (
                        "sample_multipoint"
                        if re.search(r"\bsample[_\s-]?multipoint\b", recent + " " + (reply or ""), re.I)
                        else "inspection_checklist"
                    )
                    state["chat_pending"] = {
                        "action": "recreate_form", "source": src,
                        "redirect_form": True,
                    }
                    tool_result = tool_result or {
                        "tool": "propose_recreate_form", "ok": True,
                        "actions": [{"id": "confirm_generate", "label": "Ja — opprett mal"}],
                    }
                elif re.search(r"skal\s+jeg\s+kj[øo]re\s+contract\s*review|"
                               r"shall\s+i\s+run\s+contract\s*review", reply or "", re.I):
                    pend_fp = "run_generate:contract_review"
                    asked = state.setdefault("asked_actions", [])
                    if pend_fp in asked:
                        reply = ("Bekreft med **ja** — jeg spør ikke om det samme to ganger."
                                 if hub.detect_lang(msg) != "en" else
                                 "Confirm with **yes** — I won't ask the same thing twice.")
                    else:
                        state["chat_pending"] = {
                            "action": "run_generate", "template_key": "contract_review"}
                        asked.append(pend_fp)
                        tool_result = tool_result or {
                            "tool": "propose_generate", "ok": True,
                            "actions": [{"id": "confirm_generate", "label": "Ja — kjør"}],
                        }
            if route.get("kind") == "propose_generate":
                fp = "run_generate:" + str(
                    (route.get("set_pending") or {}).get("template_key") or "contract_review")
                asked = state.setdefault("asked_actions", [])
                if fp in asked:
                    reply = ("Bekreft med **ja** — jeg spør ikke om det samme to ganger."
                             if hub.detect_lang(msg) != "en" else
                             "Confirm with **yes** — I won't ask the same thing twice.")
                    state["chat_pending"] = route.get("set_pending") or state.get("chat_pending")
                else:
                    asked.append(fp)
                tool_result = tool_result or {
                    "tool": "propose_generate", "ok": True,
                    "actions": route.get("actions") or [
                        {"id": "confirm_generate", "label": "Ja — kjør"}],
                }
            if route.get("kind") == "propose_generate_reask":
                tool_result = tool_result or {
                    "tool": "propose_generate", "ok": True,
                    "actions": route.get("actions") or [
                        {"id": "confirm_generate", "label": "Ja — kjør"}],
                }
            # WORKORDER_0.22 B2 + 0.25 C — completion/progress claims require receipts
            import agent_truth as atruth
            import manifest_claims as mc
            lang = hub.detect_lang(msg)
            tools_log = [tool_result] if tool_result else []
            ok_claim, reply, reason = atruth.validate_completion_claims(
                reply, tools_log, lang=lang)
            if not ok_claim and KEY_SET and route.get("need_model"):
                # one retry naming the violation
                try:
                    raw3 = fc.ask("chat_edit", fc.HAIKU, [{"role": "user", "content":
                        f"Your previous reply claimed a completed action without a tool receipt "
                        f"({reason}). Rewrite honestly: either report only what tools did, or say "
                        f"you lack the tool. USER: {msg}"
                    }], max_tokens=300)
                    reply2 = edchat.scrub_chat_voice((raw3 or "").strip())
                    ok2, reply2, _ = atruth.validate_completion_claims(reply2, tools_log, lang=lang)
                    reply = reply2 if ok2 else atruth.honest_fallback(lang)
                except Exception:
                    reply = atruth.honest_fallback(lang)
            elif not ok_claim:
                reply = atruth.honest_fallback(lang)
            # WORKORDER_0.23 A2 — every € amount must be manifest or tool receipt
            caps_money = hub.load_capabilities()
            extra_eur = []
            if tool_result:
                for k in ("cost_eur", "estimate_eur"):
                    if tool_result.get(k) is not None:
                        try:
                            extra_eur.append(float(tool_result[k]))
                        except (TypeError, ValueError):
                            pass
            if chat_extras and chat_extras.get("estimate_eur") is not None:
                try:
                    extra_eur.append(float(chat_extras["estimate_eur"]))
                except (TypeError, ValueError):
                    pass
            ok_money, _, reason_m = mc.validate_money_claims(
                reply, caps_money, extra_allowed=extra_eur)
            if not ok_money:
                reply = mc.money_fallback(caps_money, lang)
            # WORKORDER_0.26 A — no SVG/tables/lists/dumps in chat
            ok_art, reply_art, reason_art = atruth.validate_chat_artifacts(
                reply, user_msg=msg, lang=lang, enforce_prose_cap=bool(route.get("need_model")))
            if not ok_art:
                if KEY_SET and route.get("need_model"):
                    try:
                        raw4 = fc.ask("chat_edit", fc.HAIKU, [{"role": "user", "content":
                            f"Previous reply violated artifact-in-chat rule ({reason_art}). "
                            f"Do NOT paste SVG/HTML/tables/code/long lists. Reference the "
                            f"document artifact in ≤3 lines, or say which tool is missing. "
                            f"USER: {msg}"
                        }], max_tokens=250)
                        reply2 = edchat.scrub_chat_voice((raw4 or "").strip())
                        ok2, reply2, _ = atruth.validate_chat_artifacts(
                            reply2, user_msg=msg, lang=lang, enforce_prose_cap=True)
                        reply = reply2 if ok2 else atruth.honest_fallback(lang)
                    except Exception:
                        reply = atruth.honest_fallback(lang)
                else:
                    reply = reply_art
            # Perception: if last media present, scrub invented part numbers
            last = state.get("last_indexed_media") or {}
            if last.get("file") and ("bilde" in msg.lower() or "photo" in msg.lower()
                                     or "scan" in msg.lower() or "skann" in msg.lower()):
                index = load_index(folders, "no", state.get("user_facts"), project_name=p.get("name"))
                entry = next((e for e in index if e.get("file") == last["file"]), last)
                ok_p, reply_p, _ = atruth.validate_perception(reply, [entry], lang=lang)
                if not ok_p:
                    reply = reply_p

            edchat.append_turn(state, "bot", reply, meta=tool_result, project_id=p.get("id"))
            save_state(primary, state)
            out = {
                "reply": reply,
                "kind": route.get("kind"),
                "tool": tool_result,
                "conversation": isolated_conversation(state, p.get("id")),
                "chat_pending": state.get("chat_pending"),
                "gaps": state.get("gaps") or [],
                "gap_summary": ds.gaps_summary(state.get("gaps") or []),
            }
            if tool_result and tool_result.get("actions"):
                out["actions"] = tool_result["actions"]
            if chat_extras and chat_extras.get("open_ended"):
                out["offer_document"] = True
                out["estimate_eur"] = chat_extras["estimate_eur"]
                out["corpus_brief"] = chat_extras["corpus_brief"]
            return self._send(200, out)

        if path == "/api/conversation/mark-hint":
            pid = _pid(body)
            try:
                p, folders, primary = resolve_project(pid)
            except ValueError as e:
                return self._send(400, {"error": str(e)})
            except LookupError as e:
                return self._send(404, {"error": str(e)})
            state = load_state(primary)
            state["assist_hint_shown"] = True
            save_state(primary, state)
            return self._send(200, {"ok": True})

        if path == "/api/doc/gap-assist":
            pid = _pid(body)
            try:
                p, folders, primary, index = load_project_index(pid, "no")
                log_chat_isolation("doc/gap-assist", pid, p, folders, primary, index)
            except ValueError as e:
                return self._send(400, {"error": str(e)})
            except LookupError as e:
                return self._send(404, {"error": str(e)})
            except IsolationError as e:
                return self._send(500, {"error": str(e)})
            if not KEY_SET:
                return self._send(503, {"error": "ANTHROPIC_API_KEY er ikke satt"})
            key = body.get("key")
            state = load_state(primary)
            # reload with user_facts from THIS project's state
            index = load_index(folders, "no", state.get("user_facts"))
            cands = fc.search_fact_candidates(index, key)
            if cands:
                return self._send(200, {"step": 1, "candidates": cands, "project_id": p["id"]})
            chat_block = build_project_chat_context(
                p, folders, primary, state, index, lang=body.get("lang", "no"))["text"]
            proposal = fc.gap_assist_search(index, key, body.get("lang", "no"),
                                           project_context=chat_block)
            return self._send(200, {"step": 2, "proposal": proposal, "project_id": p["id"],
                                    "cost_eur": round(fc.LEDGER[-1]["eur"], 4) if fc.LEDGER else 0})

        if path == "/api/files/upload":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            if not KEY_SET:
                return self._send(503, {"error": "ANTHROPIC_API_KEY er ikke satt"})
            import base64
            import agent_truth as atruth
            name = (body.get("name") or "upload.txt").replace("\\", "_").replace("/", "_")
            raw = body.get("content_b64")
            if not raw:
                return self._send(400, {"error": "content_b64 er påkrevd"})
            folder = Path(p["folders"][0])
            dest = atruth.additive_dest(folder / name)  # WO 0.22 C1 — never overwrite
            assert_source_immutable_write(dest, engine_root=ROOT)
            dest.write_bytes(base64.b64decode(raw))
            name = dest.name
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            open_keys = {g["key"] for g in state.get("gaps", []) if g.get("key")}
            cache_dir = fpaths.cache_dir(folder)
            entry = fc.index_file(dest, body.get("lang", "no"), cache_dir, rel_name=name)
            matches = fc.match_new_facts_to_gaps(entry, open_keys, aliases=fact_aliases())
            return self._send(200, {"indexed": name, "facts": len(entry.get("facts", [])),
                                    "covers": list(matches.keys()),
                                    "matches": {k: {"fact_id": v["id"], "value": v["value"],
                                                    "unit": v.get("unit")} for k, v in matches.items()},
                                    "cost_eur": round(sum(l["eur"] for l in fc.LEDGER[-2:]), 4)})

        if path == "/api/chat/attach":
            """LEARNING_AND_BOUNDARIES §1 — classify drop and route into existing pipelines."""
            import base64
            name = (body.get("name") or "upload.bin").replace("\\", "_").replace("/", "_")
            raw_b64 = body.get("content_b64")
            token = (body.get("token") or "").strip()
            raw = None
            if token:
                ATTACH_STAGING.mkdir(parents=True, exist_ok=True)
                cands = list(ATTACH_STAGING.glob(f"{token}_*"))
                if cands:
                    raw = cands[0].read_bytes()
                    name = cands[0].name.split("_", 1)[-1]
                    cands[0].unlink(missing_ok=True)
            if raw is None:
                if not raw_b64:
                    return self._send(400, {"error": "content_b64 er påkrevd"})
                try:
                    raw = base64.b64decode(raw_b64)
                except Exception:
                    return self._send(400, {"error": "Ugyldig content_b64"})
            choice = body.get("choice")  # form | project
            pending_gap = (body.get("pending_gap") or body.get("gap_key") or "").strip() or None
            surface = body.get("surface") or "editor"
            # WORKORDER_0.25 A3 — hub never borrows another project's id
            pid = None if surface == "hub" else _pid(body)

            clf = chattach.classify(name, raw=raw, user_choice=choice)
            kind = clf["kind"]

            if kind == "ambiguous":
                ATTACH_STAGING.mkdir(parents=True, exist_ok=True)
                token = uuid.uuid4().hex[:12]
                (ATTACH_STAGING / f"{token}_{name}").write_bytes(raw)
                return self._send(200, {
                    "kind": "ambiguous",
                    "name": name,
                    "token": token,
                    "reply": "Skjema til gjenbruk, eller prosjektfil?",
                    "buttons": [
                        {"id": "form", "label": "Skjema til gjenbruk", "token": token},
                        {"id": "project", "label": "Prosjektfil", "token": token},
                    ],
                })

            if kind == "form_template":
                # WORKORDER_0.30 — ONE offer with field summary, no questions
                import form_model as fm
                ATTACH_STAGING.mkdir(parents=True, exist_ok=True)
                token = uuid.uuid4().hex[:12]
                stage = ATTACH_STAGING / f"{token}_{name}"
                stage.write_bytes(raw)
                peek = chattach.peek_text_bytes(raw, name)
                det = fm.detect_form_shaped(peek, name)
                # Filename form signals count even without text peek (skjema.jpg)
                if not det.get("form_shaped") and chattach.FORM_NAME.search(name):
                    det = {**det, "form_shaped": True}
                # Quick offline structure for the offer summary (0 tokens)
                draft = fm.offline_extract_form_structure(peek, name)
                summary = fm.form_summary_for_offer(draft)
                if not peek and Path(name).suffix.lower() in chattach.PHOTO_EXT:
                    summary = f"bildet «{name}» (skjema)"
                offer = fm.form_propose_reply(summary, filled=det.get("filled"), lang="no")
                return self._send(200, {
                    "kind": "form_propose",
                    "name": name,
                    "token": token,
                    "filled": det.get("filled"),
                    "summary": summary,
                    "reply": offer["reply"],
                    "actions": offer["actions"],
                    "buttons": [
                        {"id": a["id"], "label": a["label"], "token": token}
                        for a in offer["actions"]
                    ] + [{"id": "as_project", "label": "Prosjektfil", "token": token}],
                })

            # project_material
            # WORKORDER_0.25 A2/A3 — hub cold-start: index without gap-match
            if surface == "hub" or not pid:
                if surface != "hub":
                    return self._send(200, {
                        "kind": "need_project",
                        "name": name,
                        "reply": ("Jeg trenger et prosjekt først — opprett eller åpne en mappe, "
                                  "så legger jeg filen der og indekserer."),
                    })
                import hub_session as hses
                ATTACH_STAGING.mkdir(parents=True, exist_ok=True)
                token = hses.new_stage_token()
                stage = ATTACH_STAGING / f"hub_{token}_{name}"
                stage.write_bytes(raw)
                lang = hub.detect_lang(body.get("message") or name)
                caption = name
                fact_keys = []
                cost = 0.0
                entry = {"caption": name, "facts": [], "kind": "doc"}
                if KEY_SET:
                    try:
                        cache_dir = ATTACH_STAGING / "hub_cache"
                        cache_dir.mkdir(exist_ok=True)
                        entry = fc.index_file(stage, "no", cache_dir, rel_name=name)
                        caption = entry.get("caption") or name
                        fact_keys = sorted({
                            (f.get("key") or "") for f in (entry.get("facts") or [])
                            if f.get("key")
                        })
                        cost = round(sum(l["eur"] for l in fc.LEDGER[-2:]), 4) if fc.LEDGER else 0
                    except Exception:
                        caption = name
                else:
                    caption = f"basert på filnavnet — ikke indeksert ennå ({name})"
                session = hses.load_session()
                session.setdefault("staged", []).append({
                    "token": token, "name": name, "path": str(stage),
                    "caption": caption, "fact_keys": fact_keys,
                    "facts": entry.get("facts") or [],
                })
                ack = hses.hub_indexed_ack(name, caption, fact_keys, lang=lang, cost_eur=cost)
                hses.append_event(
                    session, "system",
                    f"[file_added] {name} | Indeksert som: {caption} | "
                    f"facts: {', '.join(fact_keys) or '—'}",
                    meta={"kind": "file_added", "name": name, "token": token})
                hses.append_event(session, "bot", ack, meta={"kind": "hub_indexed"})
                pending = hses.set_pending(
                    session, "create_project_from_staged",
                    {"token": token, "name": Path(name).stem},
                    offer_label="Opprett prosjekt →" if lang != "en" else "Create project →")
                hses.mark_action_asked(session, pending["fingerprint"])
                hses.save_session(session)
                return self._send(200, {
                    "kind": "hub_indexed",
                    "name": name,
                    "token": token,
                    "caption": caption,
                    "fact_keys": fact_keys,
                    "facts": len(entry.get("facts") or []),
                    "cost_eur": cost,
                    "reply": ack,
                    "pending_action": pending,
                    "actions": [{"id": "create_project", "label": pending["offer_label"],
                                 "token": token}],
                    "covers": [],
                    "matches": {},
                })

            try:
                p, folders, primary = resolve_project(pid)
            except (ValueError, LookupError, IsolationError) as e:
                return self._send(400, {"error": str(e)})
            if not KEY_SET:
                return self._send(503, {"error": "ANTHROPIC_API_KEY er ikke satt"})

            import agent_truth as atruth
            sub = chattach.dest_subdir(name)
            dest_dir = Path(primary) / sub
            dest_dir.mkdir(parents=True, exist_ok=True)
            # WO 0.22 C1 — additive only; never overwrite user sources
            dest = atruth.additive_dest(dest_dir / name)
            assert_source_immutable_write(dest, engine_root=ROOT)
            dest.write_bytes(raw)
            rel = f"{sub}/{dest.name}".replace("\\", "/")

            tf = body.get("template") or load_state(primary).get("active_template")
            state = load_state(primary, tf)
            open_keys = {g["key"] for g in state.get("gaps", []) if g.get("key")}
            if pending_gap and pending_gap not in open_keys:
                pending_gap = None  # don't invent foreign gap keys
            if pending_gap:
                open_keys.add(pending_gap)
            cache_dir = fpaths.cache_dir(primary)
            entry = fc.index_file(dest, body.get("lang", "no"), cache_dir, rel_name=rel)
            aliases = fact_aliases()
            # A3 — no gap-match when this project has no open gaps
            matches = (fc.match_new_facts_to_gaps(entry, open_keys, aliases=aliases)
                       if open_keys else {})
            cost = round(sum(l["eur"] for l in fc.LEDGER[-2:]), 4) if fc.LEDGER else 0

            # Gap-context photo → targeted extraction offer first
            targeted = None
            if pending_gap and Path(name).suffix.lower() in chattach.PHOTO_EXT:
                if pending_gap not in matches:
                    targeted = {
                        "key": pending_gap,
                        "file": rel,
                        "offer": True,
                        "message": (f"Åpen mangel `{pending_gap}` — skal jeg hente verdien "
                                    f"målrettet fra dette bildet?"),
                    }

            covers = list(matches.keys())
            cap = entry.get("caption") or dest.name
            reply = (f"Lagt i {sub}/ og indeksert (~€{cost:.2f}). "
                     f"Indeksert som: {cap}.")
            if covers:
                reply += (f" Dekker {len(covers)} mangler i dette prosjektet: "
                          f"{', '.join(covers)}.")
            elif targeted:
                reply += " " + targeted["message"]

            # Remember for «bruk dette bildet på forsiden» (WO 0.20 B)
            if Path(name).suffix.lower() in chattach.PHOTO_EXT or entry.get("kind") == "photo":
                state["last_indexed_media"] = {
                    "file": rel,
                    "caption": entry.get("caption") or dest.name,
                    "kind": entry.get("kind") or "photo",
                }

            edchat.append_turn(state, "system",
                               f"[file_added] {rel} | Indeksert som: {cap}",
                               meta={"kind": "file_added", "file": rel},
                               project_id=p.get("id"))
            edchat.append_turn(state, "user", f"[fil] {name}", project_id=p.get("id"))
            edchat.append_turn(state, "bot", reply, project_id=p.get("id"))
            save_state(primary, state)

            return self._send(200, {
                "kind": "project_material",
                "name": dest.name,
                "rel": rel,
                "subdir": sub,
                "facts": len(entry.get("facts", [])),
                "covers": covers,
                "matches": {k: {"fact_id": v["id"], "value": v["value"],
                                "unit": v.get("unit")} for k, v in matches.items()},
                "targeted": targeted,
                "cost_eur": cost,
                "reply": reply,
                "conversation": isolated_conversation(state, p.get("id")),
            })

        if path == "/api/chat/attach/import-template":
            """WORKORDER_0.30 — extract_form_structure → review payload."""
            import base64
            import form_model as fm
            token = (body.get("token") or "").strip()
            name = (body.get("name") or "skjema.pdf").replace("\\", "_").replace("/", "_")
            stage = None
            if token:
                ATTACH_STAGING.mkdir(parents=True, exist_ok=True)
                cands = list(ATTACH_STAGING.glob(f"{token}_*"))
                stage = cands[0] if cands else None
            raw = None
            if stage and stage.exists():
                raw = stage.read_bytes()
                name = stage.name.split("_", 1)[-1]
            elif body.get("content_b64"):
                raw = base64.b64decode(body["content_b64"])
            if not raw:
                return self._send(400, {"error": "Mangler staged fil eller content_b64"})

            ATTACH_STAGING.mkdir(parents=True, exist_ok=True)
            tmp = ATTACH_STAGING / f"import_{uuid.uuid4().hex[:8]}_{name}"
            tmp.write_bytes(raw)
            ledger_before = len(fc.LEDGER)
            try:
                # Form Engine v2 — ingest page rasters + layout extract (overlay when possible)
                import form_engine as fe
                import hashlib as _hl
                sha = _hl.sha256(raw).hexdigest()
                # Project-local cache only — never engine .foldok_ref_cache
                pid = _pid(body) if body else ""
                proj = get_project(pid) if pid else None
                primary = primary_folder(proj) if proj else None
                if primary is not None and primary.is_dir():
                    cache_dir = fpaths.cache_dir(primary)
                else:
                    cache_dir = ATTACH_STAGING / "form_cache"
                cache_dir.mkdir(parents=True, exist_ok=True)
                cache = cache_dir / f"formimport-{sha}.json"
                if cache.exists():
                    drafted = json.loads(cache.read_text(encoding="utf-8"))
                else:
                    ask_fn = None
                    if KEY_SET:
                        def ask_fn(purpose, _model, messages, max_tokens=4000):
                            return fc.ask_json(purpose, fc.SONNET, messages, max_tokens=max_tokens)
                    pkg = fe.package_from_upload(
                        raw, name, ask_fn=ask_fn, cache_dir=cache_dir)
                    drafted = fe.template_from_package(pkg)
                    # If overlay yielded too few fields, merge structure extract from text
                    n_fields = sum(len(s.get("fields") or []) for s in drafted.get("sections") or [])
                    if n_fields < 3:
                        text = (pkg.get("meta") or {}).get("text_peek") or ""
                        if not text:
                            try:
                                from markitdown import MarkItDown
                                text = (MarkItDown().convert(str(tmp)).text_content or "")[:12000]
                            except Exception:
                                text = raw[:4000].decode("utf-8", errors="ignore")
                        struct = fm.extract_form_structure(
                            text, name=name, lang="no", ask_fn=ask_fn)
                        # Keep form_package backgrounds from overlay path
                        fp = drafted.get("form_package")
                        drafted = struct
                        if fp:
                            drafted["form_package"] = fp
                            drafted["layout_mode"] = fp.get("layout_mode") or "overlay"
                    cache.write_text(json.dumps(drafted, ensure_ascii=False), encoding="utf-8")
                drafted["origin"] = "imported"
                drafted["badge"] = "Egen mal"
                drafted["import_status"] = "review"
                drafted["source_file"] = name
            finally:
                tmp.unlink(missing_ok=True)
                if stage:
                    stage.unlink(missing_ok=True)

            cost = round(sum(l["eur"] for l in fc.LEDGER[ledger_before:]), 4) if fc.LEDGER else 0
            ATTACH_STAGING.mkdir(parents=True, exist_ok=True)
            review_token = uuid.uuid4().hex[:12]
            (ATTACH_STAGING / f"{review_token}_draft.json").write_text(
                json.dumps(drafted, ensure_ascii=False), encoding="utf-8")
            review = fm.review_payload(drafted)
            summary = fm.form_summary_for_offer(drafted)
            return self._send(200, {
                "kind": "template_review",
                "token": review_token,
                "template_key": drafted.get("template_key"),
                "name_no": drafted.get("name_no") or drafted.get("name"),
                "document_species": drafted.get("document_species"),
                "origin": "imported",
                "badge": drafted.get("badge"),
                "review": review,
                "pills": [
                    {"section_key": s["section_key"], "title": s["title"],
                     "facts": [f.get("label_no") or f.get("key") for f in (s.get("fields") or [])[:8]]}
                    for s in review.get("sections") or []
                ],
                "cost_eur": cost,
                "model_calls": 1 if cost > 0 else 0,
                "reply": (
                    f"Foreslått mal **{drafted.get('name_no') or drafted.get('name')}** "
                    f"({summary}) — se gjennom feltene, så lagrer jeg den som egen mal."
                ),
            })

        if path == "/api/chat/attach/confirm-template":
            import form_model as fm
            token = (body.get("token") or "").strip()
            draft_path = ATTACH_STAGING / f"{token}_draft.json"
            if not draft_path.exists():
                return self._send(404, {"error": "Review-utkastet er utløpt — last opp på nytt"})
            drafted = json.loads(draft_path.read_text(encoding="utf-8"))
            # Optional edited structure from review UI
            if body.get("template"):
                drafted = fm.validate_form_template(body["template"])
            drafted["import_status"] = "confirmed"
            drafted["origin"] = "imported"
            drafted["badge"] = "Egen mal"
            drafted = fm.validate_form_template(drafted)
            fname = save_company_template(drafted)
            draft_path.unlink(missing_ok=True)
            caps = regen_capabilities()
            return self._send(200, {
                "ok": True,
                "file": fname,
                "template_key": drafted.get("template_key"),
                "name_no": drafted.get("name_no"),
                "document_species": drafted.get("document_species"),
                "badge": drafted.get("badge"),
                "origin": "imported",
                "templates": templates_list(),
                "reply": (f"Lagret **{drafted.get('name_no')}** som egen mal — "
                          f"den ligger i katalogen nå."),
                "capabilities": caps if isinstance(caps, dict) else None,
            })

        if path == "/api/doc/form-field":
            """WORKORDER_0.29 — set a form field value; measure/text → user_facts."""
            import form_model as fm
            try:
                p, folders, primary = resolve_project(_pid(body))
            except ValueError as e:
                return self._send(400, {"error": str(e)})
            except LookupError as e:
                return self._send(404, {"error": str(e)})
            except IsolationError as e:
                return self._send(500, {"error": str(e)})
            tf = body.get("template") or load_state(primary).get("active_template")
            state = load_state(primary, tf, project_id=p.get("id"))
            template = load_template(tf, primary) if tf else None
            if not template or not fm.is_form_fill(template):
                return self._send(400, {"error": "Aktivt dokument er ikke et skjema"})
            section = body.get("section") or ""
            key = body.get("key") or ""
            if not section or not key:
                return self._send(400, {"error": "section og key er påkrevd"})
            # Find field def for type
            fdef = None
            for s in template.get("sections") or []:
                if s.get("section_key") != section:
                    continue
                for f in s.get("fields") or []:
                    if f.get("key") == key:
                        fdef = f
                        break
            value = body.get("value")
            # Cycle rating3 if value omitted and type is rating3
            if value is None and fdef and fdef.get("type") == "rating3":
                cur = ((state.get("doc") or {}).get("sections") or {}).get(section, {}).get("fields", {}).get(key, {})
                opts = list(fm.RATING3_OPTIONS)
                cur_v = cur.get("value")
                idx = opts.index(cur_v) if cur_v in opts else -1
                value = opts[(idx + 1) % len(opts)]
            slot = fm.set_field(state, section, key, value,
                                note=body.get("note"), unit=body.get("unit"))
            if fdef:
                slot["type"] = fdef.get("type") or slot.get("type")
                slot["label_no"] = fdef.get("label_no") or slot.get("label_no")
                slot["unit"] = body.get("unit") or slot.get("unit") or fdef.get("unit")
            fact = fm.field_becomes_fact(state, section, key, slot, template_field=fdef)
            index = load_index(folders, "no", state.get("user_facts"), project_name=p.get("name"))
            state["gaps"] = ds.gaps_for_document(
                state, template, index, state.get("artifact") or {}, fc)
            content = ds.assemble_draft(state, template, state.get("artifact"))
            sync_draft_files(primary, state, template, tf, content)
            persist_doc(primary, state, tf)
            return self._send(200, {
                "ok": True,
                "section": section,
                "key": key,
                "slot": slot,
                "fact": fact,
                "gaps": state.get("gaps") or [],
                "gap_summary": ds.gaps_summary(state.get("gaps") or []),
            })

        if path == "/api/learning":
            if self.command == "DELETE" or body.get("clear"):
                learning.clear()
                return self._send(200, {"ok": True, "cleared": True, "path": str(learning.path())})
            return self._send(200, {"learning": learning.load(), "path": str(learning.path())})

        if path == "/api/learning/alias":
            frm = body.get("from") or body.get("from_key") or body.get("label")
            to = body.get("to") or body.get("to_key") or body.get("key")
            kind = body.get("kind") or "alias"
            if not frm or not to:
                return self._send(400, {"error": "from og to er påkrevd"})
            data = learning.record_alias(frm, to, kind=kind)
            return self._send(200, {"ok": True, "learning": data, "path": str(learning.path())})

        if path == "/api/doc/apply-matches":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            tf = body.get("template") or load_state(folder).get("active_template")
            state = load_state(folder, tf)
            template = load_template(tf)
            index = load_index(p["folders"], "no", state.get("user_facts"))
            keys = body.get("keys") or []
            key_facts = {}
            aliases = fact_aliases()
            for gk in keys:
                for e in index:
                    for f in e.get("facts", []):
                        if f.get("key") == gk or f.get("key") in aliases.get(gk, []):
                            key_facts[gk] = {**f, "source_location": e["file"]}
            result = ds.apply_multiple_cited(state, key_facts, template, index, state.get("artifact"), fc)
            persist_doc(folder, state, tf)
            return self._send(200, {**result, "gap_summary": ds.gaps_summary(result["gaps"])})

        if path == "/api/doc/regenerate-section":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            if not KEY_SET:
                return self._send(503, {"error": "ANTHROPIC_API_KEY er ikke satt"})
            sk = body.get("section_key")
            tf = body.get("template")
            if not sk or not tf:
                return self._send(400, {"error": "section_key og template er påkrevd"})
            return self._send(200, {"job": start_job(run_regenerate_section, p["folders"], tf, sk,
                                                    body.get("lang", "no"), body.get("instruction"))})

        if path == "/api/doc/accept-regen":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            sk = body.get("section_key")
            new_md = body.get("new_md")
            tf = body.get("template")
            if not sk or new_md is None or not tf:
                return self._send(400, {"error": "section_key, new_md og template er påkrevd"})
            state = load_state(folder, tf)
            template = load_template(tf)
            sections = state.setdefault("doc", {}).setdefault("sections", {})
            prev = sections.get(sk, {}).get("md", "")
            sections.setdefault(sk, {})["md"] = new_md
            sections[sk]["updated"] = ds.iso_now()
            sections[sk]["gaps"] = ds.gaps_from_md(new_md, sk)
            ds.add_version(state, "ai", "section", f"Godtok regenerert seksjon {sk}",
                           section=sk, prev_md=prev)
            index = load_index(p["folders"], "no", state.get("user_facts"))
            state["gaps"] = ds.gaps_for_document(state, template, index, state.get("artifact"), fc)
            persist_doc(folder, state, tf)
            return self._send(200, {"section": sections[sk], "gaps": state["gaps"],
                                    "gap_summary": ds.gaps_summary(state["gaps"])})

        if path == "/api/doc/edit-section":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            sk = body.get("section_key")
            md = body.get("md")
            tf = body.get("template")
            if not sk or md is None or not tf:
                return self._send(400, {"error": "section_key, md og template er påkrevd"})
            state = load_state(folder, tf)
            ensure_doc_for_template(folder, state, tf)
            template = load_template(tf)
            sections = state.setdefault("doc", {}).setdefault("sections", {})
            prev = sections.get(sk, {}).get("md", "")
            _, _, violations = fc.postprocess(sk, md, load_index(p["folders"], "no", state.get("user_facts")),
                                              state.get("artifact"))
            sections.setdefault(sk, {})["md"] = md
            sections[sk]["updated"] = ds.iso_now()
            sections[sk]["violations"] = violations
            sections[sk]["gaps"] = ds.gaps_from_md(md, sk)
            ds.add_version(state, "user", "section", "Manuell redigering", section=sk, prev_md=prev)
            state["gaps"] = ds.gaps_for_document(state, template, load_index(p["folders"], "no", state.get("user_facts")),
                                                state.get("artifact"), fc)
            persist_doc(folder, state, tf)
            return self._send(200, {"section": sections[sk], "violations": violations,
                                    "gaps": state["gaps"], "gap_summary": ds.gaps_summary(state["gaps"])})

        if path == "/api/doc/revert":
            p = get_project(_pid(body))
            if not p:
                return self._send(404, {"error": "unknown project"})
            folder = p["folders"][0]
            sk = body.get("section_key")
            tf = body.get("template")
            if not sk or not tf:
                return self._send(400, {"error": "section_key og template er påkrevd"})
            state = load_state(folder, tf)
            template = load_template(tf)
            versions = state.get("versions", [])
            prev_md = None
            for v in versions:
                if v.get("section") == sk and v.get("prev_md") is not None:
                    prev_md = v["prev_md"]
                    break
            if prev_md is None:
                return self._send(404, {"error": "Ingen tidligere versjon for denne seksjonen"})
            sections = state.setdefault("doc", {}).setdefault("sections", {})
            cur = sections.get(sk, {}).get("md", "")
            sections.setdefault(sk, {})["md"] = prev_md
            sections[sk]["updated"] = ds.iso_now()
            sections[sk]["gaps"] = ds.gaps_from_md(prev_md, sk)
            ds.add_version(state, "user", "section", f"Tilbakestill seksjon {sk}", section=sk, prev_md=cur)
            state["gaps"] = ds.gaps_for_document(state, template, load_index(p["folders"], "no", state.get("user_facts")),
                                                state.get("artifact"), fc)
            persist_doc(folder, state, tf)
            return self._send(200, {"section": sections[sk], "gaps": state["gaps"]})

        return self._send(404, {"error": "not found"})


class Server(ThreadingHTTPServer):
    # On Windows SO_REUSEADDR lets a second instance silently bind the same
    # port and steal requests — refuse instead so a stale server is obvious.
    allow_reuse_address = False


def main():
    port = 8766
    try:
        srv = Server(("127.0.0.1", port), Handler)
    except OSError:
        sys.exit(f"Port {port} er opptatt — en annen workbench kjører allerede? "
                 f"Stopp den (eller: Get-NetTCPConnection -LocalPort {port} | "
                 f"ForEach-Object {{ Stop-Process -Id $_.OwningProcess -Force }})")
    print(f"Foldok workbench v{(ROOT / 'VERSION').read_text().strip()}", flush=True)
    print(f"  engine:  foldok_compile.py (real API calls{'' if KEY_SET else ' — NO KEY SET, browse only'})", flush=True)
    print(f"  open:    http://127.0.0.1:{port}", flush=True)
    srv.serve_forever()


if __name__ == "__main__":
    main()
